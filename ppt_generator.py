"""
PPT 생성 로직: analyze_template → generate_plan → edit_slide 루프
"""
from __future__ import annotations

import json
import re
import shutil
import subprocess
import sys
import xml.etree.ElementTree as ET
import zipfile

# namespace prefix 등록 — ET.tostring이 ns0/ns1 대신 올바른 prefix를 사용하도록
ET.register_namespace('p', 'http://schemas.openxmlformats.org/presentationml/2006/main')
ET.register_namespace('a', 'http://schemas.openxmlformats.org/drawingml/2006/main')
ET.register_namespace('r', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships')
ET.register_namespace('p14', 'http://schemas.microsoft.com/office/powerpoint/2010/main')
ET.register_namespace('a14', 'http://schemas.microsoft.com/office/drawing/2010/main')
ET.register_namespace('a15', 'http://schemas.microsoft.com/office/drawing/2014/main')
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path


# 프로젝트 디렉토리 자체를 스킬 디렉토리로 사용
SKILL_DIR = Path(__file__).parent
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

# ── Plan 콘텐츠 스키마 ────────────────────────────────────────
# role별 content 필드 요구사항 정의
PLAN_CONTENT_SCHEMA: dict[str, dict] = {
    "cover":     {"required": ["subtitle"], "optional": ["date"]},
    "toc":       {"required": ["items"]},
    "closing":   {"required": [], "optional": []},
    # 모든 본문 슬라이드 공통 필드:
    #   section_no: "1.1" 형식 번호
    #   section_title: 사이드바 본문제목 (20자 이내)
    #   section_desc: 사이드바 본문설명 (3줄 이내)
    #   body: 슬라이드 유형별 상세 내용
    "content":   {"required": ["section_title"], "optional": ["section_desc","body"]},
    "timeline":  {"required": ["section_title"], "recommended": ["periods"]},
    "quarterly": {"required": ["section_title"], "recommended": ["quarters"]},
    "steps":     {"required": ["section_title"], "recommended": ["steps"]},
    "flow":      {"required": ["section_title"], "recommended": ["keywords","solutions"]},
    "comparison":{"required": ["section_title"], "recommended": ["before","after"]},
}


# ── 레이아웃 레지스트리 ────────────────────────────────────────
# CLAUDE.md 카탈로그 기반: 역할별 최적 슬라이드 + 편집기 함수명
# shape 구조:
#   slide24: ID=8(제목), ID=7(본문1/bullets), ID=10(본문2/body)        — 순수 텍스트 2블록
#   slide30: ID=8(제목), ID=9(부제목), ID=28-31(Step1~4)               — 4단계 프로세스
#   slide32: ID=8(제목), ID=9(부제목), ID=25(본문제목), ID=26(본문설명글)  — 상단텍스트+하단3열(텍스트/이미지 삽입형)
#   slide38: ID=8(제목), ID=13-15(keyword×3), ID=7/10/11(solution×3)  — 3행 흐름도
#   slide8:  ID=2(섹션제목/40pt), ID=4(서브항목 목록)                    — 섹션 구분
LAYOUT_REGISTRY: dict[str, dict] = {
    "slide24.xml": {"label": "2블록 텍스트", "best_for": ["content", "body"],
                    "editor": "_edit_slide24"},
    "slide30.xml": {"label": "4단계 스텝",   "best_for": ["steps", "process"],
                    "editor": "_edit_slide30"},
    "slide32.xml": {"label": "상단텍스트+하단3열콘텐츠", "best_for": ["content", "three_points"],
                    "editor": "_edit_slide32"},
    "slide38.xml": {"label": "3행 흐름도",   "best_for": ["flow", "architecture"],
                    "editor": "_edit_slide38"},
    "slide8.xml":  {"label": "섹션 구분",    "best_for": ["section"],
                    "editor": "_edit_slide8"},
    "slide6.xml":  {"label": "표지",         "best_for": ["cover"],
                    "editor": "edit_cover_slide"},
    "slide7.xml":  {"label": "목차",         "best_for": ["toc"],
                    "editor": "edit_toc_slide"},
    "slide46.xml": {"label": "감사합니다",   "best_for": ["closing"],
                    "editor": "noop"},
    "slide9.xml":  {"label": "3열 이미지+제목+설명", "best_for": ["three_col", "comparison", "services"],    "editor": "_edit_slide9"},
    "slide10.xml": {"label": "3열+인사이트",          "best_for": ["insights", "metrics", "comparison"],      "editor": "_edit_slide10"},
    "slide11.xml": {"label": "3열 이미지카드",         "best_for": ["features", "three_col", "highlights"],   "editor": "_edit_slide11"},
    "slide12.xml": {"label": "3열+서브헤딩",           "best_for": ["three_col", "content"],                  "editor": "_edit_slide12"},
    "slide14.xml": {"label": "4열 아이콘+인사이트",    "best_for": ["four_features", "pillars", "benefits"],  "editor": "_edit_slide14"},
    "slide16.xml": {"label": "4열 아이콘+서브헤딩",    "best_for": ["four_col", "features", "services"],      "editor": "_edit_slide16"},
    "slide17.xml": {"label": "4열 이미지+서브헤딩",    "best_for": ["four_col", "roadmap", "phases"],         "editor": "_edit_slide17"},
    "slide25.xml": {"label": "이미지+우측3열",         "best_for": ["image_text", "overview", "content"],     "editor": "_edit_slide25"},
    "slide26.xml": {"label": "이미지+3항목",           "best_for": ["image_text", "detail", "content"],       "editor": "_edit_slide26"},
    "slide27.xml": {"label": "이미지+우측3행(2줄)",    "best_for": ["image_text", "panel", "content"],        "editor": "_edit_slide27"},
    "slide28.xml": {"label": "이미지+우측3행(3줄)",    "best_for": ["image_text", "panel", "content"],        "editor": "_edit_slide28"},
    "slide34.xml": {"label": "2이미지+키워드",         "best_for": ["keywords", "highlight", "comparison"],   "editor": "_edit_slide34"},
    "slide37.xml": {"label": "3구역 텍스트",           "best_for": ["text_heavy", "explanation", "detail"],   "editor": "_edit_slide37"},
    "slide42.xml": {"label": "대형 본문",              "best_for": ["body", "deep_dive", "content"],          "editor": "_edit_slide42"},
}


# ── 1. 템플릿 분석 ────────────────────────────────────────────

def analyze_template(work_dir: Path) -> list[dict]:
    """언팩된 슬라이드 XML을 읽어 슬라이드 목록과 텍스트를 반환한다."""
    unpacked = work_dir / "unpacked"
    slides_dir = unpacked / "ppt" / "slides"
    if not slides_dir.exists():
        raise FileNotFoundError(f"unpacked slides not found: {slides_dir}")

    slides = []
    for xml_path in sorted(slides_dir.glob("slide*.xml")):
        try:
            tree = ET.parse(xml_path)
            root = tree.getroot()
        except ET.ParseError as e:
            print(f"  ⚠ {xml_path.name} parse error: {e}", file=sys.stderr)
            continue

        texts: list[str] = []
        for t_elem in root.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}t"):
            if t_elem.text and t_elem.text.strip():
                texts.append(t_elem.text.strip())

        slides.append({
            "file": xml_path.name,
            "path": str(xml_path),
            "texts": texts,
        })

    return slides


# ── 2. 계획 수립 ──────────────────────────────────────────────

def _extract_layout_order(pptx_path: Path) -> list[str]:
    """PPTX에서 슬라이드 template_file 순서를 추출한다."""
    import zipfile, re as _re
    with zipfile.ZipFile(pptx_path, "r") as z:
        xml = z.read("ppt/presentation.xml").decode("utf-8")
    # sldIdLst 안의 r:id 순서대로 슬라이드 파일명 매핑
    rels_xml = None
    with zipfile.ZipFile(pptx_path, "r") as z:
        rels_xml = z.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
    rid_to_file = dict(_re.findall(r'Id="(rId\d+)"[^>]*Target="slides/([^"]+)"', rels_xml))
    rids = _re.findall(r'<p:sldId[^>]+r:id="(rId\d+)"', xml)
    return [rid_to_file[r] for r in rids if r in rid_to_file]


def generate_plan_with_claude(
    topic: str,
    audience: str,
    n_slides: int,
    slide_info: list[dict],
    memory: dict,
    layout_order: list[str] | None = None,
) -> dict:
    """
    Claude API를 사용해 주제에 맞는 실제 콘텐츠 계획을 생성한다.
    layout_order가 주어지면 template_file 순서를 고정하고 내용만 생성한다.
    ANTHROPIC_API_KEY 없으면 None을 반환해 폴백으로 넘긴다.
    """
    import os
    import json as _json

    api_key = os.environ.get("ANTHROPIC_API_KEY")
    vertex_project = os.environ.get("ANTHROPIC_VERTEX_PROJECT_ID")
    vertex_region = os.environ.get("CLOUD_ML_REGION")
    aws_region = os.environ.get("AWS_REGION")
    explicit = os.environ.get("PPT_SKILL_BACKEND", "auto")  # auto|vertex|bedrock|anthropic
    # switch-provider.sh가 쓰는 CLAUDE_CODE_USE_* 변수 반영
    cc_bedrock = os.environ.get("CLAUDE_CODE_USE_BEDROCK", "0")
    cc_vertex  = os.environ.get("CLAUDE_CODE_USE_VERTEX", "0")

    # 백엔드 결정 — PPT_SKILL_BACKEND > CLAUDE_CODE_USE_* > 자동 감지
    if explicit == "vertex":
        use_vertex, use_bedrock = True, False
    elif explicit == "bedrock":
        use_vertex, use_bedrock = False, True
    elif explicit == "anthropic":
        use_vertex, use_bedrock = False, False
    elif cc_bedrock == "1":
        # switch-provider.sh bedrock
        use_vertex, use_bedrock = False, True
    elif cc_vertex == "1":
        # switch-provider.sh vertex
        use_vertex, use_bedrock = True, False
    elif cc_bedrock == "0" and cc_vertex == "0":
        # switch-provider.sh direct — 두 플래그가 명시적으로 0이면 직접 API
        use_vertex, use_bedrock = False, False
    else:
        # 환경 변수 미설정 시 자동 감지
        use_vertex = bool(vertex_project or (vertex_region and not aws_region))
        use_bedrock = bool(aws_region) and not use_vertex

    if not api_key and not use_vertex and not use_bedrock:
        return None

    available_layouts = {s["file"]: s["texts"][:3] for s in slide_info}
    layout_hints = memory.get("slide_layout_hints", {})
    known_issues = []
    for tmpl_data in memory.get("template_patterns", {}).values():
        known_issues = [i["issue"] for i in tmpl_data.get("known_issues", [])]

    system = (
        "당신은 프로 컨설턴트 수준의 PowerPoint 슬라이드 기획자입니다.\n"
        "콘텐츠를 먼저 설계하고, 그 콘텐츠에 가장 적합한 template_file을 선택합니다.\n"
        "반드시 JSON만 출력하고 다른 텍스트는 포함하지 않습니다.\n\n"
        "=== 사용 가능한 template_file 카탈로그 ===\n"
        "slide6.xml   → cover    : 표지 (전용, 변경 불가)\n"
        "slide7.xml   → toc      : 목차 (전용, 변경 불가)\n"
        "slide8.xml   → section  : 챕터 구분 (사용자 명시 요청 시에만 — 자동 선택 금지)\n"
        "slide29.xml  → timeline : 연도별 타임라인(2023→2026) 또는 월별(2026.1→2026.6)\n"
        "slide31.xml  → quarterly: 분기별 Q1→Q2→Q3→Q4 레이아웃\n"
        "slide33.xml  → quarterly: slide31 변형 — 상단 설명 + Q1-Q4열\n"
        "slide13.xml  → content  : 3가지 기능/특징/장점 카드 (제목+설명 3열)\n"
        "slide15.xml  → content  : 3가지 핵심 가치 (제목+설명 3열)\n"
        "slide14.xml  → content  : 4가지 기능/구성요소 (제목+설명+Insight 4열)\n"
        "slide16.xml  → content  : 4가지 구성요소 컴팩트 (제목+설명 4열)\n"
        "slide9.xml   → content  : 3가지 사례/제품 (이미지+제목+설명 3열)\n"
        "slide10.xml  → content  : 3가지 사례 + 하단 Insight 배너\n"
        "slide12.xml  → content  : 3가지 도구/기술 (이미지+우측 텍스트)\n"
        "slide30.xml  → steps    : 4단계 프로세스 (Step1→Step2→Step3→Step4)\n"
        "slide11.xml  → content  : 3가지 이미지 카드 (이미지+제목+설명 3열, 기능 강조)\n"
        "slide17.xml  → content  : 4가지 이미지+서브헤딩 (이미지+제목+설명 4열, 로드맵/단계)\n"
        "slide25.xml  → content  : 이미지+우측3열 (좌측 이미지, 우측 overview+3항목)\n"
        "slide26.xml  → content  : 이미지+3항목 상세 (좌측 이미지, 우측 overview+3항목+설명)\n"
        "slide27.xml  → content  : 이미지+우측3행 패널 (좌측 이미지, 우측 overview+3행 — body_title 길이로 slide28 자동 전환)\n"
        "slide34.xml  → content  : 2이미지+키워드 (키워드 4개+설명+이미지 2개)\n"
        "slide37.xml  → content  : 3구역 텍스트 (대형 설명+보조 설명 2개, 텍스트 중심)\n"
        "slide42.xml  → content  : 대형 본문 (단일 대형 텍스트 박스, 심층 설명)\n"
        "slide40.xml  → content  : 도넛 차트 3개 (비율·KPI 달성도 — chart_data 필수)\n"
        "slide41.xml  → content  : 막대 차트 (항목별 다시리즈 수치 비교 — chart_data 필수)\n"
        "slide43.xml  → content  : 막대 차트 (시계열 증감 추이, 음수 가능 — chart_data 필수)\n"
        "slide32.xml  → content  : 상단 텍스트+하단 3열 콘텐츠 (본문설명글+3가지 핵심 포인트/이미지)\n"
        "slide36.xml  → content  : As-is/To-be 벤다이어그램 (현황→목표 비교)\n"
        "slide38.xml  → flow     : 3행 흐름도 keyword→solution→service (아키텍처/파이프라인)\n"
        "slide46.xml  → closing  : 감사합니다 (전용, 변경 불가)\n\n"
        "=== template 선택 기준 (반드시 콘텐츠 형태와 일치시킬 것) ===\n"
        "- 서술형 설명+3가지 핵심 포인트 → slide32 (상단: 본문설명글, 하단 3열: bullets 3개 또는 이미지)\n"
        "- 3가지 기능/장점 나열 → slide13 (items 3개 + descriptions 3개 필수)\n"
        "- 4가지 기능/구성요소 → slide14 또는 slide16 (items 4개 + descriptions 4개)\n"
        "- 3가지 사례·제품·도구(시각 강조) → slide9/slide10/slide12 "
        "(items + descriptions + image_descriptions 3개)\n"
        "- 4단계 프로세스 → slide30 (steps 4개 필수)\n"
        "- 아키텍처/파이프라인 흐름 → slide38 (keywords + solutions + services/details 필수)\n"
        "- 현황↔목표·이전↔이후·비교 → slide36 (as_is/to_be 각 키워드 필수)\n"
        "- 짧은 2가지 요점 비교 → slide24 (body 2개)\n"
        "- 연도별 로드맵/연혁 → slide29 (periods=[{label,content}] 필수, 시계열일 때만)\n"
        "- 분기별 계획 → slide31/slide33 (quarters 필수, 시계열일 때만)\n"
        "- 챕터 시작 구분 → slide8 (사용자가 '섹션 구분' 명시 요청 시에만)\n"
        "- 비율·KPI 달성도(원형 게이지 3개) → slide40 (chart_data 도넛 3개, 각 값 3개)\n"
        "- 항목별 수치 막대 비교 → slide41 (chart_data 1개, 범주 4·시리즈 최대 3)\n"
        "- 시계열 증감 막대 추이 → slide43 (chart_data 1개, 범주 6·시리즈 2)\n"
        "★ 중요: 실사진/아이콘 이미지는 제공되지 않는다. 이미지가 필수인 레이아웃은 쓰지 말 것.\n"
        "★ 중요: slide29/31/33(타임라인)은 연도·분기 등 '시간 흐름' 데이터일 때만. "
        "그 외 서술형+3포인트 구조는 slide32를 기본으로 사용.\n"
        "★ 레이아웃별 필수 필드가 없으면 자동으로 slide32(텍스트)로 교체되니, "
        "선택한 레이아웃에 맞는 content 필드를 반드시 채울 것.\n\n"
        "=== 본문 슬라이드 헤더 3존 규칙 (cover/toc/closing 제외 모든 슬라이드 필수) ===\n"
        "① title(대제목)   : 이 슬라이드가 속한 목차 챕터 제목 — TOC items 중 해당 번호의 텍스트 그대로\n"
        "                    (예: TOC 2번 항목이 '핵심 아키텍처 및 처리 모델'이면 title='핵심 아키텍처 및 처리 모델')\n"
        "② subtitle(중제목): 대제목보다 구체적인 핵심 키워드나 기술명을 표현하는 곳.\n"
        "   역할: 대제목의 추상적인 챕터명을 좀 더 구체적인 기술/개념 키워드로 표현한다.\n"
        "         예) 대제목='핵심 기술'이면 중제목='Constitutional AI & RLHF' 처럼\n"
        "             그 챕터에서 실제로 다루는 핵심 기술명·개념명·프레임워크명을 직접 기재\n"
        "   규칙:\n"
        "   - 반드시 채워야 한다. 슬라이드가 1장이어도 구체적인 기술/개념명을 작성\n"
        "   - 숫자(01, 02) 단독 금지 — 반드시 의미 있는 키워드 포함\n"
        "   - 대제목이 포괄적일수록 중제목은 더 구체적인 기술명/개념명으로 작성\n"
        "   - 같은 챕터 내 동일한 주제 그룹 슬라이드 → 동일한 중제목\n"
        "   - 같은 챕터 내 다른 주제 그룹 슬라이드 → 다른 중제목\n"
        "   도출 방법: 그 슬라이드에서 실제 다루는 핵심 기술명/개념명을 중제목으로 설정\n"
        "             예) 대='핵심 기술', 본문='Constitutional AI 개요' → 중='Constitutional AI & RLHF'\n"
        "             예) 대='시장 현황', 본문='시장 규모·성장률' → 중='글로벌 AI 시장 동향'\n"
        "   예) 챕터 2 '핵심 아키텍처'에 슬라이드 3장:\n"
        "     slide A: subtitle='JobManager & TaskManager', section_title='2.1 JobManager 역할'\n"
        "     slide B: subtitle='JobManager & TaskManager', section_title='2.2 TaskManager 구성'\n"
        "     slide C: subtitle='State Backend 설계',       section_title='2.3 State Backend'\n"
        "   예) 챕터 1에 슬라이드 1장:\n"
        "     slide A: subtitle='실시간 스트리밍 시장 규모', section_title='1. 실시간 스트리밍 시장 동향'\n"
        "③ section_title   : 이 슬라이드가 실제로 표현하는 제목 — content 내부에 위치\n"
        "   - 중제목의 세부 분류 역할. 항상 title(대제목)과 다른 텍스트여야 한다\n"
        "   - 챕터 내 단독 슬라이드: 'N. 제목' (예: '1. 실시간 스트리밍 시장 동향')\n"
        "   - 챕터 내 복수 슬라이드: 'N.M 제목' (예: '2.1 JobManager 역할')\n"
        "   ★ 4단계 중복 금지 규칙 (반드시 준수):\n"
        "     대제목(title) ≠ 중제목(subtitle) ≠ section_title ≠ 슬라이드 내 소제목\n"
        "     내용이 비슷해도 되지만 동일한 텍스트는 절대 금지\n"
        "     예시(올바름): 대='ML 기술 기반' / 중='학습 패러다임 진화' / 본='1. 3대 핵심 트렌드' / 소='1.1 Transformer'\n"
        "     예시(잘못됨): 대='ML 패러다임' / 중='ML 패러다임 전환' / 본='ML 패러다임 트렌드' ← 모두 동일 주제 반복\n\n"
        "=== role별 content 필드 ===\n"
        "cover  : subtitle(30자 이내 1줄), date\n"
        "toc    : items(문자열 배열)\n"
        "steps  : section_title(소제목), section_desc(본문제목 설명 2줄 이내), "
        "steps(4개, 각 12자 이내 짧은 한 줄 — 예 '1단계: 평가', '4단계: 배포'. 15자 초과 금지, 잘림), "
        "step_descs(정확히 4개 문자열 배열 — 각 문자열은 \\n으로 구분된 bullet 항목 3개 이상. 예: ['항목1\\n항목2\\n항목3', '항목A\\n항목B\\n항목C', ...]. ⚠️ \\n 없이 문장 1개만 쓰면 상단 구역에 bullet 1줄만 표시됨. 각 bullet 40자 이내), "
        "step_meta(정확히 4개 dict, 각 dict는 정확히 4개 키-값 쌍 — 하단 4열×4행 grid에 표시. 예: [{'출시': '2023.03', '특징': '텍스트 전용', '성능': '...', '비고': '...'}, ...]. ⚠️ 4개 미만이면 하단 grid 공백).\n"
        "timeline: section_title, section_desc, "
        "periods(4개 — 각 항목: {label: '2026', period: 'Jan~Mar', items: ['□ 항목1', '□ 항목2'], "
        "kpi: '핵심지표', risk: '낮음/중간/높음', team: '담당팀'} 형식). "
        "⚠️ items는 문자열 배열 필수 — 누락 시 연도별 콘텐츠 박스 공백.\n"
        "flow   : section_title(소제목), section_desc(3~5줄 풍부하게), keywords(3개, 왼쪽 첫번째 컬럼 키워드), solutions(3개, 파란색 화살표 본문), details(3개, 두번째 컬럼 보조 설명), services(3개, 오른쪽 네번째 컬럼 세부 설명 — 누락 시 우측 박스 공백)\n"
        "content(slide13/15): section_title, section_desc(3~5줄 풍부하게), "
        "items(3개, 각 16자 이내), descriptions(3개, 각 45자 이내), "
        "image_descriptions(이미지 영역에 넣어야 할 이미지 상세 설명 3개 — 실제 이미지 아님)\n"
        "content(slide14/16): section_title, section_desc(3~5줄 풍부하게), "
        "items(4개, 16자 이내), descriptions(4개, 40자 이내), "
        "image_descriptions(이미지 영역 상세 설명 4개)\n"
        "content(slide9/11/12): section_title, section_desc, items(3개), descriptions(3개), "
        "image_descriptions(각 칸 이미지 상세 설명 3개)\n"
        "content(slide10): section_title, section_desc, items/descriptions/image_descriptions(3개) + insights(배너 1개)\n"
        "content(slide11): section_title, section_desc, items(3개, 14자 이내), descriptions(3개, 45자 이내)\n"
        "content(slide17): section_title, section_desc(3~5줄 풍부하게), items(4개, 각 16자 이내), descriptions(4개, 각 40자 이내 — 초과 시 잘림), image_descriptions(이미지 영역 상세 설명 4개)\n"
        "content(slide25/26): section_title, section_desc(3~5줄 풍부하게), items(3개, 각 16자 이내), descriptions(4개 — 첫째=overview 40자 이내, 나머지 3개=항목별 설명 각 40자 이내 — 초과 시 잘림)\n"
        "content(slide27): section_title, section_desc(3~5줄 풍부하게), items(3개, 각 16자 이내), descriptions(4개 — 첫째=overview 40자 이내, 나머지 3개=항목별 설명 각 40자 이내), image_descriptions(이미지 상세 설명 1개)\n"
        "content(slide28): section_title(45자 이내 — 초과 시 파란 배너에서 잘림), section_desc(3~4줄 풍부하게, 80자 이내), items(3개, 각 20자 이내), descriptions(3개, 각 55자 이내 4줄 이내 — 초과 시 잘림), image_descriptions(이미지 상세 설명 1개)\n"
        "content(slide31/slide33): section_title, section_desc, "
        "quarters(4개 필수 — 각 항목: {period: 'Q1', kpi: '핵심지표', risk: '낮음/중간/높음', effort: '규모', "
        "items: ['□ 항목1', '□ 항목2', '□ 항목3']} 형식). "
        "⚠️ items는 반드시 문자열 배열(3개 이상). items 누락 시 하단 콘텐츠 박스 전체 공백. "
        "kpi/risk/effort/items 내용은 주제에 맞게 자유롭게 구성(핵심목표·리스크·규모에 국한되지 않음).\n"
        "content(slide34): section_title, section_desc, keywords(4개, 10자 이내), descriptions(1개, 본문 설명)\n"
        "content(slide37): section_title, section_desc, bullets(3개 — 첫째=대형 설명 70자 이내, 나머지 2개=보조 설명 각 60자 이내 — 초과 시 잘림)\n"
        "content(slide42): section_title, section_desc(본문 설명 3줄 이내 — 좌측 사이드바 ID=28 박스에 들어감), descriptions(3개 — 각 원형 아래 텍스트박스에 표시되는 핵심 포인트, 각 40자 이내)\n"
        "content(slide32): section_title, section_desc, "
        "body(상단 본문설명글 3줄 이내), bullets(하단 3열 핵심 포인트 3개 필수), "
        "image_descriptions(하단 3열 이미지 상세 설명 3개)\n"
        "content(slide35): section_title, section_desc, before(3개 키워드), after(4개 키워드)\n"
        "content(slide36): section_title, section_desc, as_is(원형 키워드 3개, 각 8자 이내 — 길면 잘림), to_be(원형 키워드 3개, 각 8자 이내), explains(우측 설명 2개 — [좌측진영 설명, 우측진영 설명], 각 2~3줄), compare_labels(좌/우 라벨 2개, 예: ['Anthropic','OpenAI'])\n"
        "content(slide40): section_title, section_desc, chart_data(도넛 3개 = 리스트 3항목, 각 {title:'지표명(20자 이내)', values:[정확히 숫자 3개 — 세그먼트, 합 100 권장]}). ⚠️ 값은 정확히 3개(템플릿 고정·범주 라벨 미표시). 적게 주면 나머지 세그먼트는 0\n"
        "content(slide41): section_title, section_desc, chart_data(리스트 1항목 = [{categories:[범주 정확히 4개, 각 10자 이내], series:[{name:'시리즈명', values:[숫자 4개]} — 최대 3시리즈]}]). 시리즈 적게 주면 나머지 막대는 0\n"
        "content(slide43): section_title, section_desc, chart_data(리스트 1항목 = [{categories:[범주 정확히 6개, 각 10자 이내], series:[{name:'시리즈명', values:[숫자 6개, 음수 가능]} — 정확히 2시리즈]}])\n"
        "content(slide21): section_title, section_desc, bullets(3개 — 첫째=하단 핵심 배너, 나머지 2개=좌측 텍스트블록), image_descriptions(이미지 영역 상세 설명 1개)\n"
        "content(slide22/slide24): section_title, section_desc, bullets(2~3개), body, image_descriptions(이미지/영상 영역 상세 설명 1개)\n"
        "closing: 없음\n\n"
        "★★ 콘텐츠 풍부도(모든 레이아웃 공통, 최우선): 모든 텍스트 필드는 빈약한 1줄·단순 명사구로 끝내지 말고, "
        "해당 영역의 허용 분량을 충분히 활용해 '풍부해 보이게' 작성한다. 특히 section_desc(본문제목 설명글)는 3~5줄로 맥락·배경·요점을 담아 상세히, "
        "overview/main_desc(상단 콘텐츠)는 2~3줄로 충분히, item descriptions는 핵심을 구체적으로 서술한다. "
        "단 ⚠️ 과다 금지 — 각 필드의 허용 자수/줄수를 넘기면 '…'로 잘리거나 텍스트박스가 겹쳐 지저분해진다. "
        "'한도 내에서 꽉 채워 풍부하게 보이되, 넘치지 않게' 가 핵심 (잘림·겹침 발생 시 분량을 줄여 재작성).\n"
        "★ 이미지 영역 규칙: 이미지 슬롯이 있는 모든 레이아웃은 image_descriptions를 반드시 제공한다. "
        "⚠️ 단순 명사구(예: '아키텍처 다이어그램') 금지 — 나중에 자동 이미지 생성의 프롬프트로 쓸 수 있도록 "
        "'무엇을(핵심 객체·구성요소)+어떻게 표현(구도/관계)'를 담은 한 문장으로 구체 서술한다. "
        "단 한 문장(약 40~70자)으로 간결하게 — 두세 문장·과도한 색상/좌표 나열은 박스에서 지저분하게 잘리니 금지 "
        "(예: 'JobManager가 TaskManager 3개에 작업을 분배하는 분산 처리 구조도'). "
        "누락 시 엔진이 자동 파생하나, 품질을 위해 plan 단계에서 한 문장으로 구체적으로 채울 것.\n"
        "★ flow details 규칙: slide39 details의 각 항목은 문자열 또는 '여러 줄 문자열/리스트'로 줄 수 있다 — 한 칸에 2개 이상 보조 항목이 필요하면 줄바꿈(\\n)으로 구분하거나 리스트로 제공.\n"
        "★ subtitle(중제목) 필수 규칙: cover/toc/closing 제외 모든 슬라이드 top-level에 반드시 포함.\n"
        "  - 위 role별 content 목록에 없어도 항상 별도 top-level 필드로 작성\n"
        "  - 숫자(01, 02) 금지, 의미 있는 주제어만 (예: '학습 패러다임 진화', '아키텍처 구조')\n"
        "★ sub_heading(소제목 바) 규칙: 소제목 바의 '| 번호 키워드1·키워드2·키워드3 —'까지는 엔진이 items로 자동 생성한다.\n"
        "  - plan은 대시(—) 뒤에 붙을 '명사형 요약구'를 sub_heading_tail 필드(content 내)로 제공한다.\n"
        "  - ⭐ sub_heading_tail은 반드시 명사·명사형 어미로 종결('…차이점','…개요','…전략'). 서술형 종결('~다','~합니다') 금지\n"
        "  - 예: items=['설립 배경','미션·철학','핵심 투자자'] + sub_heading_tail='두 회사의 출발점과 지향의 차이점'\n"
        "★ section_desc(본문제목 설명글) 필수 규칙: cover/toc/closing 제외 모든 슬라이드 content 내에 반드시 포함.\n"
        "  - ⭐ 문체: 완결 문장(1~3문장)을 '~합니다/~입니다' 정중체로 서술. 개조식·명사형 종결('…개요','A — B' 대시 나열) 금지. 비워두면 안 됨.\n"
        "  - 최대 60자 이내 — 사이드바 박스(cx≈173pt)에 맞게 작성. 초과 시 텍스트박스에서 잘림\n"
        "  - ⚠️ 슬라이드별 제한: slide30(steps) 40자 이내, slide36/37/42 60자 이내, 그 외 60자 이내\n"
        "★ 텍스트 길이 엄수: 위에 명시된 글자 수 제한은 절대 준수. 넘으면 화면에서 '…'로 잘려 보임.\n\n"
        "출력 형식:\n"
        '{"title":"...","topic":"...","audience":"...","n_slides":N,'
        '"slides":[{"index":1,"template_file":"slideN.xml",'
        '"role":"cover|toc|section|content|steps|flow|closing",'
        '"title":"목차챕터제목(대제목)","subtitle":"세부주제명(중제목 — 숫자 금지, 의미있는 주제어 필수)","content":{"section_title":"N. 소제목",...}}]}'
    )

    layout_summary = "\n".join(f"  {f}: {t}" for f, t in list(available_layouts.items())[:15])
    hints_summary = "\n".join(f"  {k}: {v}" for k, v in layout_hints.items())

    # content_limits: harness에서 읽어 LLM에게 전달 — 초과 시 코드가 자르지 않고 LLM이 처음부터 맞춰 작성
    _cl = _load_ltm().get("content_limits", {})
    def _fmt_limits(d, indent=0):
        lines = []
        prefix = "  " * indent
        for k, v in d.items():
            if k.startswith("_"): continue
            if isinstance(v, dict):
                lines.append(f"{prefix}{k}:")
                lines.extend(_fmt_limits(v, indent+1))
            else:
                lines.append(f"{prefix}{k}: {v}")
        return lines
    content_limits_summary = "\n".join(_fmt_limits(_cl)) if _cl else ""

    layout_lock_section = ""
    if layout_order:
        order_str = "\n".join(f"  {i+1}. {f}" for i, f in enumerate(layout_order))
        layout_lock_section = (
            f"\n\n⚠️ 레이아웃 고정 모드: 아래 template_file 순서를 반드시 그대로 사용하세요. "
            f"template_file 변경 금지 — 내용(content)만 주제에 맞게 새로 작성하세요.\n"
            f"{order_str}\n"
            f"슬라이드 수는 위 목록 기준({len(layout_order)}장)으로 맞추세요."
        )

    user_prompt = (
        f"주제: {topic}\n"
        f"대상 청중: {audience}\n"
        f"슬라이드 수: {len(layout_order) if layout_order else n_slides}\n\n"
        f"사용 가능한 슬라이드 파일 (파일명: 샘플 텍스트):\n{layout_summary}\n\n"
        f"레이아웃 힌트:\n{hints_summary}\n\n"
        + (f"=== 필드별 글자 수 제한 (반드시 준수 — 초과 시 코드가 자르지 않고 직접 잘려 보임) ===\n{content_limits_summary}\n\n" if content_limits_summary else "")
        + f"주의할 이슈: {', '.join(known_issues) if known_issues else '없음'}"
        f"{layout_lock_section}\n\n"
        f"위 조건에 맞는 슬라이드 계획을 JSON으로 작성하세요. "
        f"각 슬라이드의 콘텐츠는 {topic} 전문 지식을 반영해 구체적으로 작성하세요."
    )

    messages = [{"role": "user", "content": user_prompt}]
    model = os.environ.get("ANTHROPIC_DEFAULT_SONNET_MODEL", "claude-sonnet-4-6")
    # Vertex AI는 us.anthropic. 접두사 미지원 → 제거
    if use_vertex and model.startswith("us.anthropic."):
        model = model[len("us.anthropic."):]

    if use_bedrock:
        # project-steer 패턴: boto3 직접 사용
        import boto3 as _boto3
        from botocore.config import Config as _BotoConfig
        bedrock = _boto3.client("bedrock-runtime", region_name=aws_region or "us-east-1",
                                config=_BotoConfig(read_timeout=300, connect_timeout=10))
        bedrock_model = model
        if not (bedrock_model.startswith("us.anthropic.") or bedrock_model.startswith("anthropic.")):
            bedrock_model = "us.anthropic.claude-sonnet-4-6"
        body = _json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 8192,
            "system": system,
            "messages": messages,
        })
        resp = bedrock.invoke_model(modelId=bedrock_model, body=body)
        raw = _json.loads(resp["body"].read())["content"][0]["text"].strip()
    else:
        try:
            import anthropic
        except ImportError:
            return None

        if use_vertex:
            client = anthropic.AnthropicVertex(
                project_id=vertex_project or "",
                region=vertex_region or "us-east5",
            )
        else:
            client = anthropic.Anthropic(api_key=api_key)

        response = client.messages.create(
            model=model,
            max_tokens=8192,
            thinking={"type": "adaptive"},
            system=system,
            messages=messages,
        )
        raw = response.content[0].text.strip()

    # JSON 블록 추출 (```json ... ``` 형식 대응)
    if raw.startswith("```"):
        raw = re.sub(r"^```[a-z]*\n?", "", raw)
        raw = re.sub(r"\n?```$", "", raw)

    plan = _json.loads(raw)
    plan["_generated_by"] = "vertex" if use_vertex else ("bedrock" if use_bedrock else "anthropic")
    return plan


def generate_plan(topic: str, audience: str, n_slides: int, slide_info: list[dict]) -> dict:
    """
    슬라이드 계획(plan.json)을 생성한다.
    Claude Code 스킬 모드에서는 Claude 자신이 이 함수를 대체하지만,
    폴백 모드에서는 단순 규칙 기반 계획을 만든다.
    """
    available = [s["file"] for s in slide_info]

    # 사용 가능한 슬라이드 파일을 역할에 매핑 (GS Neotek 템플릿 기준)
    role_map: dict[str, str] = {}
    for f in available:
        name = Path(f).stem  # e.g. "slide6"
        idx = int(re.search(r"\d+", name).group())
        if idx == 6:
            role_map["cover"] = f
        elif idx == 7:
            role_map["toc"] = f
        elif idx == 8:
            role_map["section"] = f
        elif idx == 14:
            role_map["three_col"] = f
        elif idx == 30:
            role_map["steps"] = f
        elif idx == 46:
            role_map["closing"] = f

    def pick(preferred: str, fallback_idx: int) -> str:
        if preferred in role_map:
            return role_map[preferred]
        if available:
            return available[min(fallback_idx, len(available) - 1)]
        return "slide1.xml"

    slides_plan: list[dict] = []

    # 표지
    slides_plan.append({
        "index": 1,
        "template_file": pick("cover", 0),
        "role": "cover",
        "title": topic,
        "content": {"subtitle": f"대상: {audience}"},
    })

    # 목차
    if n_slides >= 3:
        slides_plan.append({
            "index": 2,
            "template_file": pick("toc", 1),
            "role": "toc",
            "title": "목차",
            "content": {"items": [f"Chapter {i}" for i in range(1, min(n_slides - 1, 7) + 1)]},
        })

    # 본문 슬라이드
    for i in range(3, n_slides):
        slides_plan.append({
            "index": i,
            "template_file": pick("three_col", 2),
            "role": "content",
            "title": f"{topic} — 주요 내용 {i - 2}",
            "content": {"body": f"{topic}에 대한 심층 내용 {i - 2}"},
        })

    # 마지막: 감사합니다
    slides_plan.append({
        "index": n_slides,
        "template_file": pick("closing", -1),
        "role": "closing",
        "title": "감사합니다",
        "content": {},
    })

    return {
        "title": topic,
        "topic": topic,
        "audience": audience,
        "n_slides": n_slides,
        "slides": slides_plan,
    }


# ── 3. 슬라이드 XML 편집 ─────────────────────────────────────

_PLACEHOLDER_RE = re.compile(
    r"lorem\s+ipsum|작성해주세요|TODO|\[insert|placeholder",
    re.IGNORECASE,
)

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

# ── 슬라이드 존 맵 (analyze_zones.py로 템플릿에서 자동 생성) ──────────
# 각 본문 슬라이드의 5개 존 + 본문구역 하위 존(shape ID 리스트)을 정의.
# 단일 진실 공급원(single source of truth) — 편집기는 하드코딩 ID 대신 이걸 참조.
_ZONE_MAP_CACHE: dict | None = None
_SLIDE_CATALOG_CACHE: dict | None = None
_ZONE_FILL_CACHE: dict | None = None
_TOC_CONFIG_CACHE: dict | None = None
_COLLISION_RULES_CACHE: dict | None = None
_COMMON_FORMATTING_CACHE: dict | None = None
_SLIDE_SHAPE_IDS_CACHE: dict | None = None
_LTM_CACHE: dict | None = None


def _load_ltm() -> dict:
    """harness/long_term_memory.json 로드 (캐시). 없으면 빈 dict."""
    global _LTM_CACHE
    if _LTM_CACHE is not None:
        return _LTM_CACHE
    for p in (SKILL_DIR / "harness" / "long_term_memory.json",
              Path(__file__).parent / "harness" / "long_term_memory.json"):
        try:
            if p.exists():
                _LTM_CACHE = json.loads(p.read_text())
                return _LTM_CACHE
        except Exception:
            pass
    _LTM_CACHE = {}
    return _LTM_CACHE


def _load_zone_map() -> dict:
    """harness/layout_zone_map.json 로드 (캐시). 없으면 빈 dict."""
    global _ZONE_MAP_CACHE
    if _ZONE_MAP_CACHE is not None:
        return _ZONE_MAP_CACHE
    for p in (SKILL_DIR / "harness" / "layout_zone_map.json",
              Path(__file__).parent / "harness" / "layout_zone_map.json"):
        try:
            if p.exists():
                _ZONE_MAP_CACHE = json.loads(p.read_text())
                return _ZONE_MAP_CACHE
        except Exception:
            pass
    _ZONE_MAP_CACHE = {}
    return _ZONE_MAP_CACHE


def _load_slide_catalog() -> dict:
    """harness/slide_catalog.json 로드 (캐시). 없으면 빈 dict."""
    global _SLIDE_CATALOG_CACHE
    if _SLIDE_CATALOG_CACHE is not None:
        return _SLIDE_CATALOG_CACHE
    for p in (SKILL_DIR / "harness" / "slide_catalog.json",
              Path(__file__).parent / "harness" / "slide_catalog.json"):
        try:
            if p.exists():
                _SLIDE_CATALOG_CACHE = json.loads(p.read_text())
                return _SLIDE_CATALOG_CACHE
        except Exception:
            pass
    _SLIDE_CATALOG_CACHE = {}
    return _SLIDE_CATALOG_CACHE


def _load_zone_fill_rules() -> dict:
    """harness/zone_fill_rules.json 로드 (캐시). 없으면 하드코딩 폴백."""
    global _ZONE_FILL_CACHE
    if _ZONE_FILL_CACHE is not None:
        return _ZONE_FILL_CACHE
    for p in (SKILL_DIR / "harness" / "zone_fill_rules.json",
              Path(__file__).parent / "harness" / "zone_fill_rules.json"):
        try:
            if p.exists():
                _ZONE_FILL_CACHE = json.loads(p.read_text())
                return _ZONE_FILL_CACHE
        except Exception:
            pass
    _ZONE_FILL_CACHE = None  # 로드 실패 시 _ZONE_FILL_RULES 하드코딩 폴백
    return {}


def _load_placeholder_patterns() -> "re.Pattern | None":
    """harness/placeholder_patterns.json에서 placeholder 패턴 로드.
    JSON 존재 시 컴파일된 re.Pattern 반환, 없으면 None → 하드코딩 폴백 사용."""
    for p in (SKILL_DIR / "harness" / "placeholder_patterns.json",
              Path(__file__).parent / "harness" / "placeholder_patterns.json"):
        try:
            if p.exists():
                data = json.loads(p.read_text())
                patterns = data.get("patterns", [])
                if patterns:
                    return re.compile(
                        "(" + "|".join(patterns) + ")",
                        re.IGNORECASE,
                    )
        except Exception:
            pass
    return None


def _load_toc_config() -> dict:
    """harness/toc_layout_config.json 로드 (캐시). 없으면 빈 dict."""
    global _TOC_CONFIG_CACHE
    if _TOC_CONFIG_CACHE is not None:
        return _TOC_CONFIG_CACHE
    for p in (SKILL_DIR / "harness" / "toc_layout_config.json",
              Path(__file__).parent / "harness" / "toc_layout_config.json"):
        try:
            if p.exists():
                _TOC_CONFIG_CACHE = json.loads(p.read_text())
                return _TOC_CONFIG_CACHE
        except Exception:
            pass
    _TOC_CONFIG_CACHE = {}
    return _TOC_CONFIG_CACHE


def _load_collision_rules() -> dict:
    """harness/collision_resolution.json 로드 (캐시). 없으면 빈 dict."""
    global _COLLISION_RULES_CACHE
    if _COLLISION_RULES_CACHE is not None:
        return _COLLISION_RULES_CACHE
    for p in (SKILL_DIR / "harness" / "collision_resolution.json",
              Path(__file__).parent / "harness" / "collision_resolution.json"):
        try:
            if p.exists():
                _COLLISION_RULES_CACHE = json.loads(p.read_text())
                return _COLLISION_RULES_CACHE
        except Exception:
            pass
    _COLLISION_RULES_CACHE = {}
    return _COLLISION_RULES_CACHE


def _load_common_formatting() -> dict:
    """harness/common_formatting.json 로드 (캐시). 없으면 빈 dict."""
    global _COMMON_FORMATTING_CACHE
    if _COMMON_FORMATTING_CACHE is not None:
        return _COMMON_FORMATTING_CACHE
    for p in (SKILL_DIR / "harness" / "common_formatting.json",
              Path(__file__).parent / "harness" / "common_formatting.json"):
        try:
            if p.exists():
                _COMMON_FORMATTING_CACHE = json.loads(p.read_text())
                return _COMMON_FORMATTING_CACHE
        except Exception:
            pass
    _COMMON_FORMATTING_CACHE = {}
    return _COMMON_FORMATTING_CACHE


def _ppt_lang() -> str:
    """harness/common_formatting.json의 locale.lang 반환. 기본값 ko-KR."""
    return _load_common_formatting().get("locale", {}).get("lang", "ko-KR")


def _load_slide_shape_ids() -> dict:
    """harness/slide_shape_ids.json 로드 (캐시). 없으면 빈 dict."""
    global _SLIDE_SHAPE_IDS_CACHE
    if _SLIDE_SHAPE_IDS_CACHE is not None:
        return _SLIDE_SHAPE_IDS_CACHE
    for p in (SKILL_DIR / "harness" / "slide_shape_ids.json",
              Path(__file__).parent / "harness" / "slide_shape_ids.json"):
        try:
            if p.exists():
                _SLIDE_SHAPE_IDS_CACHE = json.loads(p.read_text())
                return _SLIDE_SHAPE_IDS_CACHE
        except Exception:
            pass
    _SLIDE_SHAPE_IDS_CACHE = {}
    return _SLIDE_SHAPE_IDS_CACHE


def _zone(template_file: str) -> dict:
    """주어진 slideN.xml의 존 정의 반환 (base 이름으로도 조회)."""
    zm = _load_zone_map()
    if template_file in zm:
        return zm[template_file]
    base = re.sub(r"_c\d+\.xml$", ".xml", template_file)  # 중복 복사본(slide38_c1.xml) 대응
    return zm.get(base, {})

# 표지 레이아웃 상수 (GS Neotek 2026 템플릿 기준, EMU 단위)
_COVER_TITLE_ID = "12"       # 대제목 shape id
_COVER_SUBTITLE_ID = "6"     # 소제목 shape id
_COVER_DATE_ID = "15"        # 날짜 shape id
_COVER_TITLE_FONT_PT = 48
_COVER_SUBTITLE_FONT_PT = 32
_EMU_PER_PT = 12700
# 가시성 최적 자간 권장값 (OOXML spcPct):
#   대제목: 100% (lnSpc) — 빽빽하지 않고 깔끔한 밀도
#   소제목: 100% (lnSpc) — 단일 줄이므로 줄간격 자체 의미 없으나 일관성 유지
_COVER_TITLE_LNSPC = 100000   # 100 % × 1000
# 동적 여백 — 타이포그래피 권장: 대제목 폰트의 0.5×, 0.4× 배율
_GAP_DATE_TITLE_RATIO = 0.5   # date → title 간격 = 0.5 × title_font_height
_GAP_TITLE_SUBTITLE_RATIO = 0.4  # title → subtitle 간격


def _make_cover_run_xml(text: str, font_pt: int, font_name: str) -> str:
    """한국어 혼합 텍스트용 단일 run XML 조각 반환 (latin + ea + cs 모두 지정)."""
    sz = font_pt * 100  # hundredths of pt
    # lang=ko-KR 으로 설정해야 LibreOffice가 ea 폰트로 한국어를 올바르게 렌더링
    return (
        f'<a:r xmlns:a="{_NS_A}">'
        f'<a:rPr lang="{_ppt_lang()}" altLang="en-US" sz="{sz}" dirty="0">'
        f'<a:latin typeface="{font_name}" pitchFamily="2" charset="-127"/>'
        f'<a:ea typeface="{font_name}" pitchFamily="2" charset="-127"/>'
        f'<a:cs typeface="{font_name}" pitchFamily="2" charset="-127"/>'
        f'</a:rPr>'
        f'<a:t>{text}</a:t>'
        f'</a:r>'
    )


# ── 폰트 메트릭 (fonttools 기반 실제 렌더링) ──────────────────────────────
# PowerPoint는 Latin/CJK 스크립트 전환 시 약 5% 추가 간격을 삽입한다.
# fonttools advance width × 1.05 를 실효 폭으로 사용.
_RENDER_OVERHEAD = 1.05

_FONT_PATHS = {
    "Pretendard SemiBold": "/Users/toule/Library/Fonts/Pretendard-SemiBold.otf",
    "Pretendard":          "/Users/toule/Library/Fonts/Pretendard-Regular.otf",
    "Pretendard Medium":   "/Users/toule/Library/Fonts/Pretendard-Medium.otf",
    "Pretendard Bold":     "/Users/toule/Library/Fonts/Pretendard-Bold.otf",
}
_FONT_METRIC_CACHE: dict = {}  # font_name → (tt, upm, cmap, hmtx)


def _load_font_metrics(font_name: str):
    """fonttools TTFont 메트릭 로드 (캐시). 없으면 None."""
    if font_name in _FONT_METRIC_CACHE:
        return _FONT_METRIC_CACHE[font_name]
    path = _FONT_PATHS.get(font_name)
    if not path:
        # 이름에서 Pretendard 변형 추론
        for key in _FONT_PATHS:
            if key.lower() in font_name.lower() or font_name.lower() in key.lower():
                path = _FONT_PATHS[key]
                break
    if not path:
        _FONT_METRIC_CACHE[font_name] = None
        return None
    try:
        from fontTools.ttLib import TTFont as _TTFont
        tt = _TTFont(path)
        metrics = (tt, tt['head'].unitsPerEm, tt.getBestCmap(), tt['hmtx'].metrics)
        _FONT_METRIC_CACHE[font_name] = metrics
        return metrics
    except Exception:
        _FONT_METRIC_CACHE[font_name] = None
        return None


def _char_width_emu(ch: str, font_pt: int, font_name: str = "Pretendard SemiBold") -> int:
    """
    한 글자의 advance width를 EMU로 반환한다.
    fonttools 사용 가능하면 실제 폰트 메트릭, 아니면 보정된 계수 fallback.

    보정 계수 (fonttools Pretendard SemiBold 실측):
      한글: 0.864  (구: 0.85)
      영문 평균: 0.63  (구: 0.52 — 'W'=0.97, 'A'=0.69, 'I'=0.26 등 분산 큼)
      공백: 0.237  (구: 0.28)
    """
    metrics = _load_font_metrics(font_name)
    if metrics is not None:
        _, upm, cmap, hmtx = metrics
        cp = ord(ch)
        gid = cmap.get(cp) if cmap else None
        if gid is not None:
            adv, _ = hmtx.get(gid, (int(upm * 0.5), 0))
            return int(adv / upm * font_pt * _EMU_PER_PT)

    # ── fallback: 보정된 계수 ──
    cp = ord(ch)
    if (0xAC00 <= cp <= 0xD7A3 or 0x3131 <= cp <= 0x318E or 0x4E00 <= cp <= 0x9FFF):
        return int(font_pt * 0.864 * _EMU_PER_PT)
    if ch == ' ':
        return int(font_pt * 0.237 * _EMU_PER_PT)
    return int(font_pt * 0.630 * _EMU_PER_PT)


def _text_width_emu(text: str, font_pt: int, font_name: str = "Pretendard SemiBold") -> int:
    """텍스트 전체 advance width를 EMU로 반환 (렌더 오버헤드 포함)."""
    raw = sum(_char_width_emu(ch, font_pt, font_name) for ch in text)
    return int(raw * _RENDER_OVERHEAD)


def _estimate_lines(text: str, cx_emu: int, font_pt: int,
                    font_name: str = "Pretendard SemiBold") -> int:
    """실제 폰트 메트릭 기반으로 텍스트가 몇 줄로 렌더될지 추정."""
    if not text:
        return 1
    line_width = 0
    lines = 1
    for ch in text:
        cw = int(_char_width_emu(ch, font_pt, font_name) * _RENDER_OVERHEAD)
        if line_width + cw > cx_emu:
            lines += 1
            line_width = cw
        else:
            line_width += cw
    return lines


def _truncate_to_lines(text: str, cx_emu: int, font_pt: int, max_lines: int,
                        font_name: str = "Pretendard SemiBold", ellipsis: bool = True) -> str:
    """실제 폰트 메트릭 기반으로 max_lines 초과 시 자른다.
    ellipsis=True면 '…'를 붙이고, False면 '…' 없이 깔끔히 끊는다(이미지 캡션 등 — 지저분한 '…' 방지)."""
    if not text:
        return text
    _ell = "…" if ellipsis else ""
    # 개행 문자도 줄 바꿈으로 계산
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    ellipsis_w = int(_char_width_emu("…", font_pt, font_name) * _RENDER_OVERHEAD) if ellipsis else 0
    line_width = 0
    lines = 1
    for i, ch in enumerate(text):
        if ch == '\n':
            lines += 1
            if lines > max_lines:
                cut = i
                return text[:cut] + _ell
            line_width = 0
            continue
        cw = int(_char_width_emu(ch, font_pt, font_name) * _RENDER_OVERHEAD)
        if line_width + cw > cx_emu:
            lines += 1
            if lines > max_lines:
                cut = i
                while cut > 0 and line_width > cx_emu - ellipsis_w:
                    cut -= 1
                    line_width -= int(_char_width_emu(text[cut], font_pt, font_name) * _RENDER_OVERHEAD)
                return text[:cut] + _ell
            line_width = cw
        else:
            line_width += cw
    return text


def _write_xml(root: ET.Element, xml_path: Path) -> None:
    """네임스페이스 처리를 안전하게 유지하며 XML을 저장한다."""
    xml_str = ET.tostring(root, encoding="unicode")
    with open(xml_path, "w", encoding="utf-8") as f:
        f.write('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n')
        f.write(xml_str)


def _find_shape_by_id(root: ET.Element, shape_id: str):
    for sp in root.iter(f"{{{_NS_P}}}sp"):
        cpr = sp.find(f"{{{_NS_P}}}nvSpPr/{{{_NS_P}}}cNvPr")
        if cpr is not None and cpr.get("id") == shape_id:
            return sp
    return None


def _set_shape_text_single_run(sp: ET.Element, text: str, font_pt: int, font_name: str,
                                lnspc_pct: int = None) -> None:
    """
    shape의 기존 rPr(색상·효과 포함)을 보존하면서 텍스트만 교체한다.
    - 첫 번째 run의 rPr을 복사해 lang만 ko-KR로 변경
    - 나머지 run 제거, lnSpc 추가
    """
    ns_a = _NS_A
    ns_p = _NS_P
    txBody = sp.find(f"{{{ns_p}}}txBody")
    if txBody is None:
        return

    # 첫 번째 p의 첫 번째 r에서 rPr 복사 (색상·폰트 스타일 보존용)
    original_rPr = None
    for p in txBody.findall(f"{{{ns_a}}}p"):
        r = p.find(f"{{{ns_a}}}r")
        if r is not None:
            rPr_elem = r.find(f"{{{ns_a}}}rPr")
            if rPr_elem is not None:
                import copy
                original_rPr = copy.deepcopy(rPr_elem)
            break

    # 기존 <a:p> 모두 제거
    for p in txBody.findall(f"{{{ns_a}}}p"):
        txBody.remove(p)

    # 새 <a:p> 구성
    p_elem = ET.SubElement(txBody, f"{{{ns_a}}}p")

    # pPr (줄간격)
    pPr = ET.SubElement(p_elem, f"{{{ns_a}}}pPr")
    if lnspc_pct is not None:
        lnSpc = ET.SubElement(pPr, f"{{{ns_a}}}lnSpc")
        ET.SubElement(lnSpc, f"{{{ns_a}}}spcPct", val=str(lnspc_pct))

    # run — 원본 rPr 보존 + lang만 ko-KR로 변경
    r_elem = ET.SubElement(p_elem, f"{{{ns_a}}}r")
    if original_rPr is not None:
        original_rPr.set("lang", _ppt_lang())
        original_rPr.set("altLang", "en-US")
        original_rPr.set("dirty", "0")
        r_elem.append(original_rPr)
    else:
        # 원본 없을 때만 새로 생성 (폴백)
        rPr = ET.SubElement(r_elem, f"{{{ns_a}}}rPr",
                            lang=_ppt_lang(), altLang="en-US",
                            sz=str(font_pt * 100), dirty="0")
        ET.SubElement(rPr, f"{{{ns_a}}}latin",
                      typeface=font_name, pitchFamily="2", charset="-127")
        ET.SubElement(rPr, f"{{{ns_a}}}ea",
                      typeface=font_name, pitchFamily="2", charset="-127")
        ET.SubElement(rPr, f"{{{ns_a}}}cs",
                      typeface=font_name, pitchFamily="2", charset="-127")

    t_elem = ET.SubElement(r_elem, f"{{{ns_a}}}t")
    t_elem.text = text


def _set_shape_cy(sp: ET.Element, cy: int) -> None:
    """shape의 xfrm ext cy 값을 교체한다."""
    xfrm = sp.find(f".//{{{_NS_A}}}xfrm")
    if xfrm is None:
        return
    ext = xfrm.find(f"{{{_NS_A}}}ext")
    if ext is not None:
        ext.set("cy", str(cy))


def _set_shape_cx(sp: ET.Element, cx: int) -> None:
    """shape의 xfrm ext cx 값을 교체한다."""
    xfrm = sp.find(f".//{{{_NS_A}}}xfrm")
    if xfrm is None:
        return
    ext = xfrm.find(f"{{{_NS_A}}}ext")
    if ext is not None:
        ext.set("cx", str(cx))


def _set_shape_y(sp: ET.Element, y: int) -> None:
    """shape의 xfrm off y 값을 교체한다."""
    xfrm = sp.find(f".//{{{_NS_A}}}xfrm")
    if xfrm is None:
        return
    off = xfrm.find(f"{{{_NS_A}}}off")
    if off is not None:
        off.set("y", str(y))


def _get_shape_geometry(sp: ET.Element) -> tuple[int, int, int, int]:
    """(x, y, cx, cy) 반환. 없으면 (0,0,0,0)."""
    xfrm = sp.find(f".//{{{_NS_A}}}xfrm")
    if xfrm is None:
        return (0, 0, 0, 0)
    off = xfrm.find(f"{{{_NS_A}}}off")
    ext = xfrm.find(f"{{{_NS_A}}}ext")
    x = int(off.get("x", 0)) if off is not None else 0
    y = int(off.get("y", 0)) if off is not None else 0
    cx = int(ext.get("cx", 0)) if ext is not None else 0
    cy = int(ext.get("cy", 0)) if ext is not None else 0
    return (x, y, cx, cy)


def _enable_autofit(sp: ET.Element) -> None:
    """bodyPr의 noAutofit을 normAutofit으로 전환한다."""
    ns_a = _NS_A
    bodyPr = sp.find(f".//{{{ns_a}}}bodyPr")
    if bodyPr is None:
        return
    for child in list(bodyPr):
        tag = child.tag.split("}")[-1]
        if tag in ("noAutofit", "spAutoFit", "normAutofit"):
            bodyPr.remove(child)
    ET.SubElement(bodyPr, f"{{{ns_a}}}normAutofit")


def edit_cover_slide(xml_path: Path, title: str, subtitle: str, date_str: str) -> None:
    """
    표지 슬라이드(slide6) 전용 편집:
    - 대제목 (shape 12): Pretendard SemiBold 48pt, lnSpc=100%, 최대 2줄
    - 소제목 (shape 6 ): Pretendard 32pt, 1줄, 대제목과 같은 cx 폭
    - 날짜   (shape 15): 그대로 삽입
    - 날짜→대제목→소제목 Y 위치를 동적으로 재배치

    타이포그래피 권장값 (발표 자료 가시성 기준):
      - 줄간격 1.0× (100%): 대형 제목은 1.0~1.15가 최적 — 지나치게 넓으면 연결감 소실
      - Date→Title 간격: title_font_height × 0.5 (여백이 확보되면서 묶임감 유지)
      - Title→Subtitle 간격: title_font_height × 0.4 (소제목은 제목에 가깝게)
    """
    tree = ET.parse(xml_path)
    root = tree.getroot()

    sp_title = _find_shape_by_id(root, _COVER_TITLE_ID)
    sp_sub = _find_shape_by_id(root, _COVER_SUBTITLE_ID)
    sp_date = _find_shape_by_id(root, _COVER_DATE_ID)

    if sp_title is None or sp_sub is None or sp_date is None:
        # shape ID가 맞지 않는 템플릿 → 폴백: 첫 번째 텍스트만 교체
        _replace_first_text_legacy(xml_path, title)
        return

    # ── 대제목 cx 가져오기 (소제목 동기화 기준) ──────────────────
    _, date_y, _, date_cy = _get_shape_geometry(sp_date)
    _, _, title_cx, _ = _get_shape_geometry(sp_title)

    # ── 1. 날짜 ────────────────────────────────────────────────
    _set_shape_text_single_run(sp_date, date_str,
                               font_pt=14, font_name="Pretendard SemiBold")

    # ── 2. 대제목 ───────────────────────────────────────────────
    title_text = _truncate_to_lines(title, title_cx, _COVER_TITLE_FONT_PT, max_lines=2)
    title_lines = _estimate_lines(title_text, title_cx, _COVER_TITLE_FONT_PT)
    title_cy = title_lines * _COVER_TITLE_FONT_PT * _EMU_PER_PT
    _enable_autofit(sp_title)
    _set_shape_cy(sp_title, title_cy)
    _set_shape_text_single_run(sp_title, title_text,
                               font_pt=_COVER_TITLE_FONT_PT,
                               font_name="Pretendard SemiBold",
                               lnspc_pct=_COVER_TITLE_LNSPC)

    # ── 3. 소제목 ───────────────────────────────────────────────
    # 최대 2줄 허용, wrap="square"로 PowerPoint 자연 줄바꿈
    sub_text = _truncate_to_lines(subtitle, title_cx, _COVER_SUBTITLE_FONT_PT, max_lines=2)
    sub_lines = min(_estimate_lines(subtitle, title_cx, _COVER_SUBTITLE_FONT_PT), 2)
    sub_cy = sub_lines * _COVER_SUBTITLE_FONT_PT * _EMU_PER_PT
    _set_shape_cx(sp_sub, title_cx)
    _set_shape_cy(sp_sub, sub_cy)
    ns_a = _NS_A
    bodyPr_sub = sp_sub.find(f".//{{{ns_a}}}bodyPr")
    if bodyPr_sub is not None:
        bodyPr_sub.attrib.pop("wrap", None)  # wrap="none" 제거 → 기본 square
        for child in list(bodyPr_sub):
            if child.tag.split("}")[-1] in ("spAutoFit", "normAutofit", "noAutofit"):
                bodyPr_sub.remove(child)
        ET.SubElement(bodyPr_sub, f"{{{ns_a}}}normAutofit")
    _set_shape_text_single_run(sp_sub, sub_text,
                               font_pt=_COVER_SUBTITLE_FONT_PT,
                               font_name="Pretendard",
                               lnspc_pct=_COVER_TITLE_LNSPC)

    # ── 4. 동적 Y 위치 재배치 ────────────────────────────────────
    # 기준: 날짜 위치는 고정, 아래로 title→subtitle 배치
    title_font_height = _COVER_TITLE_FONT_PT * _EMU_PER_PT
    gap_date_title = int(title_font_height * _GAP_DATE_TITLE_RATIO)
    gap_title_sub = int(title_font_height * _GAP_TITLE_SUBTITLE_RATIO)

    title_y = date_y + date_cy + gap_date_title
    sub_y = title_y + title_cy + gap_title_sub

    _set_shape_y(sp_title, title_y)
    _set_shape_y(sp_sub, sub_y)

    _write_xml(root, xml_path)


def _replace_first_text_legacy(xml_path: Path, new_text: str) -> None:
    """폴백: 첫 번째 <a:t>만 교체 (cover 전용 shape ID가 없는 구형 템플릿용)."""
    tree = ET.parse(xml_path)
    root = tree.getroot()
    ns_a = _NS_A
    t_elems = list(root.iter(f"{{{ns_a}}}t"))
    if not t_elems:
        return
    t_elems[0].text = new_text
    for t in t_elems[1:]:
        t.text = ""
    _write_xml(root, xml_path)


def _replace_first_text(xml_path: Path, new_text: str) -> None:
    """슬라이드 XML의 첫 번째 텍스트 run을 new_text로 교체하고 나머지 run은 비운다."""
    _replace_first_text_legacy(xml_path, new_text)


def _toc_row_xml(row_idx: int, row_y: int, number: str, item_text: str, page_num: str,
                 row_h: int = 674674, line_offset: int = 431324) -> str:
    """
    목차 슬라이드의 한 행(Row)을 구성하는 4개 shape XML 반환.
    번호·텍스트·선·페이지가 동일 Y를 공유하여 항상 정렬됨.
    wrap=none + noAutofit으로 자동줄바꿈 완전 차단.
    """
    y = row_y + row_idx * row_h
    line_y = y + line_offset
    base_id = 200 + row_idx * 4

    # 색상/폰트는 원본 템플릿 그대로 사용
    # 번호: bg2 lumMod=75000 / 텍스트: tx1 lumMod=95000,lumOff=5000 / 페이지: bg2 lumMod=75000
    NUM_X, NUM_CX = 3851538, 536801
    TXT_X, TXT_CX = 4541178, 4383747   # 선 시작점(8924925)까지 확장
    LINE_X, LINE_CX = 8924925, 3267075
    PAGE_X, PAGE_CX = 11431475, 760525

    # 비어있는 행이면 빈 텍스트로 (선도 투명하게)
    is_empty = not number and not item_text and not page_num
    line_alpha = "0" if is_empty else "100000"

    return f"""
  <p:sp>
    <p:nvSpPr><p:cNvPr id="{base_id}" name="toc_num_{row_idx}"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="{NUM_X}" y="{y}"/><a:ext cx="{NUM_CX}" cy="{row_h}"/></a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/>
    </p:spPr>
    <p:txBody>
      <a:bodyPr wrap="none" lIns="0" tIns="0" rIns="0" bIns="0" rtlCol="0" anchor="ctr"><a:noAutofit/></a:bodyPr>
      <a:lstStyle/>
      <a:p>
        <a:r>
          <a:rPr lang="{_ppt_lang()}" altLang="en-US" sz="3000" b="1" dirty="0">
            <a:solidFill><a:schemeClr val="bg2"><a:lumMod val="75000"/></a:schemeClr></a:solidFill>
            <a:latin typeface="Pretendard SemiBold" panose="02000703000000020004" pitchFamily="2" charset="-127"/>
            <a:ea typeface="Pretendard SemiBold" panose="02000703000000020004" pitchFamily="2" charset="-127"/>
            <a:cs typeface="Pretendard SemiBold" panose="02000703000000020004" pitchFamily="2" charset="-127"/>
          </a:rPr>
          <a:t>{number}</a:t>
        </a:r>
      </a:p>
    </p:txBody>
  </p:sp>
  <p:sp>
    <p:nvSpPr><p:cNvPr id="{base_id+1}" name="toc_text_{row_idx}"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="{TXT_X}" y="{y}"/><a:ext cx="{TXT_CX}" cy="{row_h}"/></a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/>
    </p:spPr>
    <p:txBody>
      <a:bodyPr wrap="none" lIns="0" tIns="0" rIns="0" bIns="0" rtlCol="0" anchor="ctr"><a:noAutofit/></a:bodyPr>
      <a:lstStyle/>
      <a:p>
        <a:r>
          <a:rPr lang="{_ppt_lang()}" altLang="en-US" sz="3000" dirty="0">
            <a:solidFill><a:schemeClr val="tx1"><a:lumMod val="95000"/><a:lumOff val="5000"/><a:alpha val="99000"/></a:schemeClr></a:solidFill>
            <a:latin typeface="Pretendard SemiBold" panose="02000703000000020004" pitchFamily="2" charset="-127"/>
            <a:ea typeface="Pretendard SemiBold" panose="02000703000000020004" pitchFamily="2" charset="-127"/>
            <a:cs typeface="Pretendard SemiBold" panose="02000703000000020004" pitchFamily="2" charset="-127"/>
          </a:rPr>
          <a:t>{item_text}</a:t>
        </a:r>
      </a:p>
    </p:txBody>
  </p:sp>
  <p:cxnSp>
    <p:nvCxnSpPr><p:cNvPr id="{base_id+2}" name="toc_line_{row_idx}"/><p:cNvCxnSpPr/><p:nvPr/></p:nvCxnSpPr>
    <p:spPr><a:xfrm><a:off x="{LINE_X}" y="{line_y}"/><a:ext cx="{LINE_CX}" cy="0"/></a:xfrm>
      <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
      <a:ln><a:solidFill><a:schemeClr val="bg2"><a:lumMod val="90000"/><a:alpha val="{line_alpha}"/></a:schemeClr></a:solidFill></a:ln>
    </p:spPr>
  </p:cxnSp>
  <p:sp>
    <p:nvSpPr><p:cNvPr id="{base_id+3}" name="toc_page_{row_idx}"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="{PAGE_X}" y="{y}"/><a:ext cx="{PAGE_CX}" cy="{row_h}"/></a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/>
    </p:spPr>
    <p:txBody>
      <a:bodyPr wrap="none" lIns="0" tIns="0" rIns="0" bIns="0" rtlCol="0" anchor="ctr"><a:noAutofit/></a:bodyPr>
      <a:lstStyle/>
      <a:p>
        <a:pPr algn="ctr"/>
        <a:r>
          <a:rPr lang="{_ppt_lang()}" altLang="en-US" sz="1400" dirty="0">
            <a:solidFill><a:schemeClr val="bg2"><a:lumMod val="75000"/></a:schemeClr></a:solidFill>
            <a:latin typeface="Pretendard" panose="02000503000000020004" pitchFamily="2" charset="-127"/>
            <a:ea typeface="Pretendard" panose="02000503000000020004" pitchFamily="2" charset="-127"/>
            <a:cs typeface="Pretendard" panose="02000503000000020004" pitchFamily="2" charset="-127"/>
          </a:rPr>
          <a:t>{page_num}</a:t>
        </a:r>
      </a:p>
    </p:txBody>
  </p:sp>"""


def _shorten_toc_item(text: str) -> str:
    """TOC 항목 텍스트가 1줄을 초과할 때 핵심 키워드만 남겨 축약한다.
    — 조사/접속어/설명구(em dash 이후, 괄호 내용) 제거 후 반환."""
    import re as _re
    # em dash(—) 또는 ': ' 이후 설명구 제거
    shortened = _re.split(r'\s*[—:]\s*', text)[0].strip()
    # 괄호 내용 제거
    shortened = _re.sub(r'\([^)]*\)', '', shortened).strip()
    # 마침표/쉼표 뒤 부연 제거
    shortened = _re.split(r'[,，]', shortened)[0].strip()
    return shortened if shortened else text


def edit_toc_slide(xml_path: Path, prs_title: str, items: list[str],
                   page_nums: list[str]) -> None:
    """
    목차 슬라이드(slide7) 편집.
    기존 shape 구조를 그대로 유지하면서 내용만 교체한다:
    - ID=10 (항목 텍스트): wrap=none + noAutofit, 각 paragraph에 항목 삽입
    - ID=11 (번호):        wrap=none + noAutofit, N개 초과분 비움
    - ID=14 (페이지번호):  wrap=none + noAutofit, N개 초과분 비움
    - ID=12~41 (선 7개):  N개 초과 행의 선은 투명(alpha=0) 처리
    - ID=42 (제목 ph):    제목 텍스트 교체
    원본 shape를 유지해야 슬라이드 레이아웃 placeholder가 노출되지 않는다.
    """
    ns_a = _NS_A
    ns_p = _NS_P
    import copy

    tree = ET.parse(xml_path)
    root = tree.getroot()

    # 하네스 설정 로드 (모든 상수를 여기서 결정)
    toc_cfg = _load_toc_config()
    TOC_TEXT_X  = toc_cfg.get("text_box", {}).get("x_emu", 4541178)
    TOC_TEXT_CX = toc_cfg.get("text_box", {}).get("cx_emu", 6790297)
    TOC_PAGE_X  = toc_cfg.get("page_number", {}).get("x_emu", 11431475)
    SLIDE_RIGHT = toc_cfg.get("slide_right_emu", 12192000)
    TOC_FONT_PT = toc_cfg.get("font_pt", 30)
    LINE_GAP    = toc_cfg.get("line_gap_emu", 350000)
    LINE_X_MAX  = toc_cfg.get("line_x_max_emu", SLIDE_RIGHT - 2_700_000)
    MAX_ITEM_CHARS = toc_cfg.get("max_item_chars", 999)
    width_factors = toc_cfg.get("text_width_factors", {})
    KO_FACTOR   = width_factors.get("korean", 1.1)
    EN_FACTOR   = width_factors.get("english", 0.55)
    LINE_IDS    = set(toc_cfg.get("line_shape_ids", ["12","13","37","38","39","40","41"]))

    MAX_ROWS = toc_cfg.get("max_rows", 7)
    # 하네스 max_item_chars 적용: 초과 항목은 말줄임표 없이 자름
    items = [it[:MAX_ITEM_CHARS] if it and len(it) > MAX_ITEM_CHARS else it for it in items]
    items_padded   = list(items)    + [""] * (MAX_ROWS - len(items))
    pages_padded   = list(page_nums) + [""] * (MAX_ROWS - len(page_nums))

    # ── 헬퍼: 텍스트박스 너비(cx) 추출 ──────────────────────────────
    def get_cx(sp: ET.Element) -> int:
        xfrm = sp.find(f".//{{{ns_a}}}xfrm")
        if xfrm is None:
            return 0
        ext = xfrm.find(f"{{{ns_a}}}ext")
        return int(ext.get("cx", 0)) if ext is not None else 0

    # ── 헬퍼: paragraph 내 모든 <a:t>를 합쳐 하나의 run으로 정리 ──
    def set_para_text(para: ET.Element, text: str) -> None:
        """기존 첫 번째 run의 rPr를 보존하고 텍스트만 교체, 나머지 run 제거.
        새 run은 endParaRPr 앞에 삽입 — 뒤에 넣으면 PowerPoint가 무시함."""
        # 첫 run rPr 복사
        first_r = para.find(f"{{{ns_a}}}r")
        orig_rPr = None
        if first_r is not None:
            rPr_e = first_r.find(f"{{{ns_a}}}rPr")
            if rPr_e is not None:
                orig_rPr = copy.deepcopy(rPr_e)
                orig_rPr.set("lang", _ppt_lang())
                orig_rPr.set("dirty", "0")

        # 기존 run/fld 전체 제거
        for r in para.findall(f"{{{ns_a}}}r"):
            para.remove(r)
        for fld in para.findall(f"{{{ns_a}}}fld"):
            para.remove(fld)

        # endParaRPr 위치 파악 — 새 run은 반드시 그 앞에 삽입
        end_rpr = para.find(f"{{{ns_a}}}endParaRPr")
        insert_idx = list(para).index(end_rpr) if end_rpr is not None else len(para)

        r_new = ET.Element(f"{{{ns_a}}}r")
        if orig_rPr is not None:
            r_new.append(orig_rPr)
        else:
            ET.SubElement(r_new, f"{{{ns_a}}}rPr", lang=_ppt_lang(), dirty="0")
        t_new = ET.SubElement(r_new, f"{{{ns_a}}}t")
        t_new.text = text
        para.insert(insert_idx, r_new)

    # ── 헬퍼: connector 선 투명 처리 ───────────────────────────────
    def set_line_alpha(cxn: ET.Element, alpha: str) -> None:
        spPr = cxn.find(f"{{{ns_p}}}spPr")
        if spPr is None:
            return
        ln = spPr.find(f"{{{ns_a}}}ln")
        if ln is None:
            return
        for schemeClr in ln.findall(f".//{{{ns_a}}}schemeClr"):
            # 기존 alpha 제거 후 새 값 설정
            for alp in schemeClr.findall(f"{{{ns_a}}}alpha"):
                schemeClr.remove(alp)
            ET.SubElement(schemeClr, f"{{{ns_a}}}alpha", val=alpha)

    # ── 1. 제목 placeholder (ID=42) ────────────────────────────────
    for sp in root.findall(f".//{{{ns_p}}}sp"):
        cpr = sp.find(f"{{{ns_p}}}nvSpPr/{{{ns_p}}}cNvPr")
        if cpr is None or cpr.get("id") != "42":
            continue
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None:
            break
        paras = txBody.findall(f"{{{ns_a}}}p")
        if paras:
            set_para_text(paras[0], prs_title)
            for p in paras[1:]:
                set_para_text(p, "")
        break

    # ── 2. 항목 텍스트 (ID=10) ─────────────────────────────────────
    for sp in root.findall(f".//{{{ns_p}}}sp"):
        cpr = sp.find(f"{{{ns_p}}}nvSpPr/{{{ns_p}}}cNvPr")
        if cpr is None or cpr.get("id") != "10":
            continue
        # cx 확장 (4.08" → 4.79")
        _set_shape_cx(sp, TOC_TEXT_CX)
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None:
            break
        paras = txBody.findall(f"{{{ns_a}}}p")
        for i, para in enumerate(paras):
            raw_text = items_padded[i] if i < MAX_ROWS else ""
            if raw_text:
                # 1줄 수렴 보장: 넓어진 cx 기준으로 자름 (font_pt=30)
                safe = _truncate_to_lines(raw_text, TOC_TEXT_CX, TOC_FONT_PT, max_lines=1)
                # 잘린 경우 → 핵심 키워드만 남겨 재시도 (조사·접속어·설명구 제거)
                if safe.endswith("…"):
                    shortened = _shorten_toc_item(raw_text)
                    safe2 = _truncate_to_lines(shortened, TOC_TEXT_CX, TOC_FONT_PT, max_lines=1)
                    if not safe2.endswith("…"):
                        safe = safe2
                    # 그래도 잘리면 더 공격적으로 축약
                    else:
                        words = shortened.split()
                        for cut in range(len(words) - 1, 0, -1):
                            candidate = " ".join(words[:cut])
                            t = _truncate_to_lines(candidate, TOC_TEXT_CX, TOC_FONT_PT, max_lines=1)
                            if not t.endswith("…"):
                                safe = t
                                break
            else:
                safe = ""
            set_para_text(para, safe)
        break

    # ── 3. 번호 열 (ID=11) — bodyPr 유지, 텍스트만 교체 ────────────
    for sp in root.findall(f".//{{{ns_p}}}sp"):
        cpr = sp.find(f"{{{ns_p}}}nvSpPr/{{{ns_p}}}cNvPr")
        if cpr is None or cpr.get("id") != "11":
            continue
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None:
            break
        paras = txBody.findall(f"{{{ns_a}}}p")
        for i, para in enumerate(paras):
            num = f"{i+1}." if (i < len(items) and items[i]) else ""
            set_para_text(para, num)
        break

    # ── 4. 페이지 번호 열 (ID=14) — bodyPr 유지, 텍스트만 교체 ──────
    for sp in root.findall(f".//{{{ns_p}}}sp"):
        cpr = sp.find(f"{{{ns_p}}}nvSpPr/{{{ns_p}}}cNvPr")
        if cpr is None or cpr.get("id") != "14":
            continue
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None:
            break
        paras = txBody.findall(f"{{{ns_a}}}p")
        for i, para in enumerate(paras):
            set_para_text(para, pages_padded[i] if i < MAX_ROWS else "")
        break

    # ── 5. 수평선 7개 ────────────────────────────────────────────────
    # 가장 긴 항목 기준으로 모든 선의 x 시작점을 통일 (텍스트 겹침 방지)

    def _est_width_emu(text, font_pt: int) -> int:
        """한글/영문 혼합 텍스트의 렌더 폭 추정 (EMU)."""
        if not isinstance(text, str):
            text = str(text) if text else ""
        ko = sum(1 for c in text if ord(c) > 0x1000)
        en = len(text) - ko
        return int((ko * font_pt * KO_FACTOR + en * font_pt * EN_FACTOR) * _EMU_PER_PT)

    max_text_emu = max(
        (_est_width_emu(it, TOC_FONT_PT) for it in items if it),
        default=0,
    )
    # 밑줄: 가장 긴 항목 텍스트 끝 + 여백 → 슬라이드 우측 끝까지
    # 페이지번호는 밑줄 위에 우측 정렬 (템플릿 동일 구조)
    line_x_start = TOC_TEXT_X + max_text_emu + LINE_GAP
    line_x_start = min(line_x_start, LINE_X_MAX)  # 하네스 line_x_max_emu 캡 적용
    line_cx_len  = SLIDE_RIGHT - line_x_start  # 슬라이드 끝까지

    line_cxns = []
    for cxn in root.findall(f".//{{{ns_p}}}cxnSp"):
        cpr = cxn.find(f"{{{ns_p}}}nvCxnSpPr/{{{ns_p}}}cNvPr")
        if cpr is not None and cpr.get("id") in LINE_IDS:
            line_cxns.append(cxn)

    line_cxns.sort(key=lambda c: int(
        (c.find(f".//{{{ns_a}}}off") or ET.Element("x")).get("y", "0")))

    for i, cxn in enumerate(line_cxns):
        # x, cx 업데이트
        off = cxn.find(f".//{{{ns_a}}}off")
        ext = cxn.find(f".//{{{ns_a}}}ext")
        if off is not None:
            off.set("x", str(line_x_start))
        if ext is not None:
            ext.set("cx", str(max(line_cx_len, 0)))
        # 빈 행 투명 처리
        alpha = "100000" if i < len(items) else "0"
        set_line_alpha(cxn, alpha)

    _write_xml(root, xml_path)
    ET.parse(xml_path)  # 유효성 검증


def _apply_formatting_safe(sp: ET.Element, zone_fmt: dict, common_fmt: dict) -> None:
    """
    하네스 기반 속성을 XML에 안전하게 적용.
    기존 구조 보존, 속성 값만 변경.

    zone_fmt 구조 (두 가지):
    1. 단일 zone: {"bodyPr": {...}, "pPr": {...}, "rPr": {...}}
    2. 라인별: {"bodyPr": {...}, "line1": {"pPr":{...}, "rPr":{...}}, "line2": {...}}
    """
    import copy as _copy
    ns_a = _NS_A
    ns_p = _NS_P

    txBody = sp.find(f"{{{ns_p}}}txBody")
    if txBody is None:
        return

    # bodyPr 적용
    if "bodyPr" in zone_fmt:
        bodyPr = txBody.find(f"{{{ns_a}}}bodyPr")
        if bodyPr is None:
            bodyPr = ET.Element(f"{{{ns_a}}}bodyPr")
            txBody.insert(0, bodyPr)
        for k, v in zone_fmt["bodyPr"].items():
            bodyPr.set(k, str(v))

    # 라인별 or 단일 pPr/rPr 적용
    has_line_fmt = "line1" in zone_fmt or "line2" in zone_fmt
    paras = txBody.findall(f"{{{ns_a}}}p")

    for i, para in enumerate(paras):
        line_fmt = zone_fmt.get(f"line{i+1}", {}) if has_line_fmt else zone_fmt

        # pPr 적용 (기존 요소 수정, 없으면 추가)
        if "pPr" in line_fmt:
            pPr_cfg = line_fmt["pPr"]
            pPr = para.find(f"{{{ns_a}}}pPr")
            if pPr is None:
                pPr = ET.Element(f"{{{ns_a}}}pPr")
                para.insert(0, pPr)
            if "algn" in pPr_cfg:
                pPr.set("algn", pPr_cfg["algn"])
            if "lnSpc_pct" in pPr_cfg:
                lnSpc = pPr.find(f"{{{ns_a}}}lnSpc")
                if lnSpc is None:
                    lnSpc = ET.SubElement(pPr, f"{{{ns_a}}}lnSpc")
                for ch in list(lnSpc): lnSpc.remove(ch)
                ET.SubElement(lnSpc, f"{{{ns_a}}}spcPct", val=str(pPr_cfg["lnSpc_pct"]))
            if "lnSpc_pts" in pPr_cfg:
                lnSpc = pPr.find(f"{{{ns_a}}}lnSpc")
                if lnSpc is None:
                    lnSpc = ET.SubElement(pPr, f"{{{ns_a}}}lnSpc")
                for ch in list(lnSpc): lnSpc.remove(ch)
                ET.SubElement(lnSpc, f"{{{ns_a}}}spcPts", val=str(pPr_cfg["lnSpc_pts"]))
            if "spcBef_pts" in pPr_cfg:
                spcBef = pPr.find(f"{{{ns_a}}}spcBef")
                if spcBef is None:
                    spcBef = ET.SubElement(pPr, f"{{{ns_a}}}spcBef")
                for ch in list(spcBef): spcBef.remove(ch)
                ET.SubElement(spcBef, f"{{{ns_a}}}spcPts", val=str(pPr_cfg["spcBef_pts"]))

        # rPr 적용 (기존 rPr 보존 + 속성만 수정, 새로 생성 금지)
        if "rPr" in line_fmt:
            rPr_cfg = line_fmt["rPr"]
            for run in para.findall(f"{{{ns_a}}}r"):
                rPr = run.find(f"{{{ns_a}}}rPr")
                if rPr is None:
                    continue  # rPr 없으면 건드리지 않음 (새로 생성 금지)
                if "sz" in rPr_cfg:
                    rPr.set("sz", str(rPr_cfg["sz"]))
                if "b" in rPr_cfg:
                    if rPr_cfg["b"]:
                        rPr.set("b", "1")
                    elif "b" in rPr.attrib:
                        del rPr.attrib["b"]
                if "typeface" in rPr_cfg:
                    tf = rPr_cfg["typeface"]
                    for tag in ("latin", "ea", "cs"):
                        elem = rPr.find(f"{{{ns_a}}}{tag}")
                        if elem is not None:
                            elem.set("typeface", tf)
                if "color" in rPr_cfg:
                    sf = rPr.find(f"{{{ns_a}}}solidFill")
                    if sf is None:
                        sf = ET.SubElement(rPr, f"{{{ns_a}}}solidFill")
                    for ch in list(sf): sf.remove(ch)
                    ET.SubElement(sf, f"{{{ns_a}}}srgbClr", val=rPr_cfg["color"])
                elif "color_scheme" in rPr_cfg:
                    sf = rPr.find(f"{{{ns_a}}}solidFill")
                    if sf is None:
                        sf = ET.SubElement(rPr, f"{{{ns_a}}}solidFill")
                    for ch in list(sf): sf.remove(ch)
                    sc = ET.SubElement(sf, f"{{{ns_a}}}schemeClr",
                                       val=rPr_cfg["color_scheme"])
                    if "lumMod" in rPr_cfg:
                        ET.SubElement(sc, f"{{{ns_a}}}lumMod", val=str(rPr_cfg["lumMod"]))
                    if "lumOff" in rPr_cfg:
                        ET.SubElement(sc, f"{{{ns_a}}}lumOff", val=str(rPr_cfg["lumOff"]))


def edit_slide(work_dir: Path, slide_plan: dict) -> bool:
    """
    단일 슬라이드를 계획에 따라 편집한다.
    성공하면 True, 실패하면 False를 반환한다.
    """
    slides_dir = work_dir / "unpacked" / "ppt" / "slides"
    xml_path = slides_dir / slide_plan["template_file"]

    if not xml_path.exists():
        print(f"  ⚠ 슬라이드 파일 없음: {xml_path}", file=sys.stderr)
        return False

    # 원본 백업
    backup = xml_path.with_suffix(".xml.bak")
    shutil.copy2(xml_path, backup)

    try:
        role = slide_plan.get("role", "")
        content = slide_plan.get("content", {})

        if role == "cover":
            edit_cover_slide(
                xml_path,
                title=slide_plan["title"],
                subtitle=content.get("subtitle", ""),
                date_str=content.get("date", "2026.00.00"),
            )
        elif role == "toc":
            # items 필드가 없으면 toc_items, bullets, body 순으로 폴백
            items = (
                content.get("items")
                or [i["text"] if isinstance(i, dict) else i
                    for i in content.get("toc_items", [])]
                or content.get("bullets", [])
                or []
            )
            page_nums = content.get("page_nums", [])
            # prs_title: ID=42 "발표 제목" placeholder → plan 전체 제목
            prs_title_for_toc = slide_plan.get("prs_title") or slide_plan["title"]
            edit_toc_slide(xml_path, prs_title_for_toc, items, page_nums)
        elif role == "closing":
            # 마무리 슬라이드(slide46): 일체 편집하지 않음
            # slide46은 레이아웃에서 "감사합니다." 텍스트를 제공
            # — 편집하면 레이아웃 placeholder가 덮어씌워져 텍스트가 사라짐
            # placeholder 잔여 텍스트만 클리어
            try:
                tree_c = ET.parse(xml_path)
                _clear_residual_placeholders(tree_c.getroot())
                _write_xml(tree_c.getroot(), xml_path)
            except ET.ParseError:
                pass

        else:
            # template_file 기반 전용 편집기 디스패치 (모듈 레벨 _SLIDE_EDITORS 자동 등록 참조)
            tmpl_name = slide_plan.get("template_file", "")
            editor = _SLIDE_EDITORS.get(tmpl_name)
            if editor is None:
                # slide32_c2.xml → slide32.xml 처럼 _cN suffix 벗겨서 기본 편집기 재조회
                import re as _re
                base_name = _re.sub(r"_c\d+(?=\.xml$)", "", tmpl_name)
                editor = _SLIDE_EDITORS.get(base_name)
            if editor:
                editor(xml_path, slide_plan)
            else:
                # 전용 편집기 없으면 존 맵 기반 제너릭 편집기 (이미지 그리드 등)
                _edit_zonemap_slide(xml_path, slide_plan)

        # 편집 후 남은 placeholder 자동 제거
        try:
            tree = ET.parse(xml_path)
            if _clear_residual_placeholders(tree.getroot()):
                _write_xml(tree.getroot(), xml_path)
        except ET.ParseError:
            pass

        # 유효성 검증
        ET.parse(xml_path)
        print(f"  ✓ slide {slide_plan['index']} ({xml_path.name}) 편집 완료")
        backup.unlink(missing_ok=True)
        return True

    except ET.ParseError as e:
        print(f"  ✗ XML 오류, 원본 복원: {e}", file=sys.stderr)
        shutil.copy2(backup, xml_path)
        backup.unlink(missing_ok=True)
        return False


# ── 4. 슬라이드 트리밍 + 클린 & 패킹 ────────────────────────

def generate_excel_for_charts(work_dir: Path, plan: dict, output_dir: Path) -> Path | None:
    """
    차트 슬라이드(slide40~43)가 있을 때 Excel 데이터 파일을 생성한다.
    생성된 xlsx를 PPT와 같은 폴더에 저장 — 사용자가 PPT에서 '데이터 편집'으로 열 수 있다.
    """
    chart_slides = [s for s in plan.get("slides", [])
                    if s.get("template_file","") in
                    {"slide40.xml","slide41.xml","slide43.xml"}]
    if not chart_slides:
        return None
    try:
        import openpyxl as _xl
    except ImportError:
        print("  ⚠ openpyxl 없음 — Excel 생성 건너뜀 (pip install openpyxl)")
        return None

    wb = _xl.Workbook()
    for slide in chart_slides:
        tmpl = slide.get("template_file","")
        chart_data = slide.get("content", {}).get("chart_data", {})
        sheet_name = f"{slide['index']}p_{tmpl.replace('.xml','')}"[:31]
        ws = wb.create_sheet(sheet_name)

        # 기본 헤더
        ws.append(["항목", "값", "비고"])
        if isinstance(chart_data, dict) and chart_data:
            for k, v in chart_data.items():
                ws.append([k, v, ""])
        elif isinstance(chart_data, list) and chart_data:
            # 새 스키마(list): 각 차트의 범주/시리즈를 표로 (임베디드 xlsx는 _update_chart_data가 갱신)
            for ce in chart_data:
                if not isinstance(ce, dict):
                    continue
                cats = ce.get("categories", [])
                series = ce.get("series") or ([{"name": ce.get("title", ""), "values": ce.get("values", [])}]
                                              if ce.get("values") is not None else [])
                ws.append([ce.get("title", "")] + [sr.get("name", "") for sr in series])
                _nr = max([len(cats)] + [len(sr.get("values", [])) for sr in series] + [0])
                for ri in range(_nr):
                    ws.append([(cats[ri] if ri < len(cats) else "")]
                              + [(sr.get("values", [])[ri] if ri < len(sr.get("values", [])) else "") for sr in series])
        else:
            # 템플릿별 기본 데이터 예시
            if "40" in tmpl:
                for item in [("KPI 1", 20), ("KPI 2", 35), ("KPI 3", 45)]:
                    ws.append(list(item) + [""])
            elif "41" in tmpl:
                for item in [("항목A", 83), ("항목B", 76), ("항목C", 65),
                              ("항목D", 58), ("항목E", 57)]:
                    ws.append(list(item) + [""])
            elif "42" in tmpl:
                for item in [("메인 KPI", 65), ("보조1", 75),
                              ("보조2", 75), ("보조3", 75)]:
                    ws.append(list(item) + [""])
            elif "43" in tmpl:
                for month, val in enumerate([2, 4, -2, 11, 4, 5], 1):
                    ws.append([f"{month}월", val, ""])

    if wb.sheetnames:
        wb.remove(wb["Sheet"])  # 기본 시트 제거
    xlsx_path = output_dir / f"{plan.get('topic','chart').replace(' ','_')}_data.xlsx"
    wb.save(str(xlsx_path))
    print(f"  ✓ Excel 데이터 파일 생성: {xlsx_path.name}")
    return xlsx_path


# ── 차트 데이터 주입 (slide40/41/43 Excel 연동 완성) ──────────────────
# 차트 캐시(chartN.xml의 numCache/strCache) + 임베디드 xlsx를 plan content.chart_data로
# 갱신해, 템플릿 더미값(도넛 5/20/75 등) 대신 실제 수치를 표시한다.
# ⚠ 데이터포인트 '개수'는 절대 바꾸지 않는다 — c:pt 추가/삭제 시 PowerPoint가 PDF
#    export에서 조용히 실패(차트 렌더 불가). 기존 c:pt의 c:v 텍스트만 제자리 덮어쓴다.
_CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"

def _register_chart_ns(chart_path: Path) -> None:
    """차트 XML 루트의 모든 xmlns 접두사를 등록 — 저장 시 prefix 보존(ns0 방지)."""
    try:
        head = chart_path.read_text(encoding="utf-8")[:2000]
    except Exception:
        return
    for m in re.finditer(r'xmlns:(\w+)="([^"]+)"', head):
        ET.register_namespace(m.group(1), m.group(2))

def _chart_overwrite_cache(parent: ET.Element, values: list) -> None:
    """c:cat/c:val/c:tx 아래 numCache|strCache의 기존 c:pt c:v 텍스트만 제자리 덮어쓴다.
    개수·ptCount·dPt·공식 등 구조는 불변(개수 변경 시 PowerPoint PDF export 실패).
    values가 적으면 남는 포인트는 0(숫자)/''(문자), 많으면 초과분은 버린다."""
    C = _CHART_NS
    def Q(t): return f"{{{C}}}{t}"
    ref = parent.find(Q("numRef"))
    if ref is None: ref = parent.find(Q("strRef"))
    if ref is None: return
    cache = ref.find(Q("numCache"))
    if cache is None: cache = ref.find(Q("strCache"))
    if cache is None: return
    pts = cache.findall(Q("pt"))
    if not pts: return  # 기존 포인트 없음(예: 도넛 cat) → 개수 변경 금지 위해 건너뜀
    is_num = cache.tag == Q("numCache")
    def _isnum(x):
        try: float(str(x)); return True
        except Exception: return False
    # 숫자 캐시에 비숫자 문자열을 쓰면 타입 불일치로 PowerPoint가 손상 처리 →
    # numRef/numCache를 strRef/strCache로 변환(개수·구조는 유지)
    if is_num and values and not all(_isnum(x) for x in values):
        ref.tag = Q("strRef"); cache.tag = Q("strCache")
        fc = cache.find(Q("formatCode"))
        if fc is not None: cache.remove(fc)  # strCache에는 formatCode 없음
        is_num = False
    for i, pt in enumerate(pts):
        v = pt.find(Q("v"))
        if v is None: continue
        v.text = str(values[i]) if i < len(values) else ("0" if is_num else "")

def _chart_embed_path(chart_path: Path, embed_dir: Path):
    """차트 rels에서 연동된 embeddings/*.xlsx 경로 반환(없으면 None — 예: OLE 객체)."""
    rels = chart_path.parent / "_rels" / (chart_path.name + ".rels")
    if not rels.exists(): return None
    m = re.search(r'Target="[^"]*embeddings/([^"]+\.xlsx)"', rels.read_text(encoding="utf-8"))
    return (embed_dir / m.group(1)) if m else None

def _update_chart_xlsx(embed_path: Path, cats: list, series: list) -> None:
    """임베디드 워크북 갱신: A열=범주, B/C/…=시리즈 값, 1행=시리즈명(데이터 편집 일관성)."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(embed_path))
    except Exception:
        return
    ws = wb.worksheets[0]
    for r in range(1, 40):
        for c in range(1, 10):
            ws.cell(row=r, column=c, value=None)
    for j, s in enumerate(series):
        ws.cell(row=1, column=2 + j, value=s.get("name", ""))
    nrows = max([len(cats)] + [len(s.get("values", [])) for s in series] + [0])
    for r in range(nrows):
        if r < len(cats): ws.cell(row=2 + r, column=1, value=cats[r])
        for j, s in enumerate(series):
            vals = s.get("values", [])
            if r < len(vals): ws.cell(row=2 + r, column=2 + j, value=vals[r])
    try: wb.save(str(embed_path))
    except Exception: pass

def _update_one_chart(chart_path: Path, embed_path, entry: dict) -> None:
    """chartN.xml 캐시 + 임베디드 xlsx를 entry로 갱신(개수 불변 제자리 덮어쓰기).
    entry: {"categories":[...], "series":[{"name","values":[...]}], "title"?}
           또는 단일시리즈 단축 {"categories":[...], "values":[...], "title"?}"""
    C = _CHART_NS
    def Q(t): return f"{{{C}}}{t}"
    _register_chart_ns(chart_path)
    try:
        root = ET.parse(chart_path).getroot()
    except Exception:
        return
    cats = [str(c) for c in (entry.get("categories") or [])]
    series = entry.get("series")
    if not series and entry.get("values") is not None:
        series = [{"name": entry.get("title", ""), "values": entry.get("values")}]
    series = series or []
    for i, ser in enumerate(root.findall(f".//{Q('ser')}")):
        sdata = series[i] if i < len(series) else None
        tx = ser.find(Q("tx"))
        if tx is not None:
            if sdata and sdata.get("name"):
                _chart_overwrite_cache(tx, [sdata["name"]])
            elif sdata is None:
                _chart_overwrite_cache(tx, [""])      # 미제공 시리즈 라벨 비움
        cat = ser.find(Q("cat"))
        if cat is not None and cats:
            _chart_overwrite_cache(cat, cats)
        val = ser.find(Q("val"))
        if val is not None:
            _chart_overwrite_cache(val, sdata.get("values", []) if sdata else [])
    _write_xml(root, chart_path)
    if embed_path is not None and embed_path.exists():
        _update_chart_xlsx(embed_path, cats, series)

def _update_chart_data(work_dir: Path, plan: dict) -> None:
    """차트 레이아웃(slide40/41/43)의 차트 데이터를 plan content.chart_data로 주입.
    harness/chart_map.json: 슬라이드→차트파일(시각 순서). chart_data[i] → charts[i].
    각 차트 데이터포인트 개수는 템플릿 고정 — chart_map data_hint 개수에 맞춰 제공."""
    cmap_path = Path(__file__).parent / "harness" / "chart_map.json"
    if not cmap_path.exists(): return
    try:
        cmap = json.loads(cmap_path.read_text(encoding="utf-8"))
    except Exception:
        return
    charts_dir = work_dir / "unpacked" / "ppt" / "charts"
    embed_dir  = work_dir / "unpacked" / "ppt" / "embeddings"
    updated = 0
    for slide in plan.get("slides", []):
        spec = cmap.get(slide.get("template_file", ""))
        if not spec: continue
        chart_data = slide.get("content", {}).get("chart_data") or []
        for ci, cf in enumerate(spec.get("charts", [])):
            entry = chart_data[ci] if ci < len(chart_data) else None
            chart_path = charts_dir / cf
            if not entry or not chart_path.exists(): continue
            _update_one_chart(chart_path, _chart_embed_path(chart_path, embed_dir), entry)
            updated += 1
    if updated:
        print(f"  ✓ 차트 데이터 주입: {updated}개 차트 (캐시 제자리 갱신+임베디드 xlsx)")


def _add_slide_to_presentation(work_dir: Path, slide_filename: str,
                                before_filename: str | None = None) -> None:
    """복사된 슬라이드 파일을 presentation.xml과 rels에 등록한다.
    before_filename: 이 파일의 rId 앞에 삽입 (None이면 맨 뒤).
    올바른 위치에 삽입해야 감사합니다 슬라이드가 마지막에 유지된다.
    """
    import re as _re
    prs_path  = work_dir / "unpacked" / "ppt" / "presentation.xml"
    rels_path = work_dir / "unpacked" / "ppt" / "_rels" / "presentation.xml.rels"
    if not prs_path.exists() or not rels_path.exists():
        return

    rels_raw = rels_path.read_text(encoding="utf-8")
    prs_raw  = prs_path.read_text(encoding="utf-8")

    # 새 rId 생성
    existing_ids = [int(m) for m in _re.findall(r'Id="rId(\d+)"', rels_raw)]
    new_rid_num  = max(existing_ids, default=0) + 1
    new_rid      = f"rId{new_rid_num}"

    # rels에 추가 (순서 무관)
    slide_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
    new_rel    = f'<Relationship Id="{new_rid}" Type="{slide_type}" Target="slides/{slide_filename}"/>'
    rels_new   = rels_raw.replace("</Relationships>", f"  {new_rel}\n</Relationships>")
    rels_path.write_text(rels_new, encoding="utf-8")

    # presentation.xml sldIdLst에 올바른 위치에 삽입
    existing_sld_ids = [int(m) for m in _re.findall(r'<p:sldId\b[^>]*\bid="(\d+)"', prs_raw)]
    new_sld_id = max(existing_sld_ids, default=255) + 1
    new_sld    = f'<p:sldId id="{new_sld_id}" r:id="{new_rid}"/>'

    if before_filename:
        # before_filename의 rId를 찾아서 그 앞에 삽입
        before_rid = None
        for m in _re.finditer(
            r'<Relationship\b([^>]+)/>', rels_raw
        ):
            attrs = m.group(1)
            if before_filename in attrs:
                rid_m = _re.search(r'\bId="(rId\d+)"', attrs)
                if rid_m:
                    before_rid = rid_m.group(1)
                    break
        if before_rid:
            # before_rid를 가진 sldId 앞에 삽입
            prs_new = _re.sub(
                rf'(\s*<p:sldId\b[^>]*\br:id="{_re.escape(before_rid)}"[^/]*/>\s*)',
                f'\n  {new_sld}\\1',
                prs_raw, count=1
            )
            prs_path.write_text(prs_new, encoding="utf-8")
            return

    # before_filename 없거나 못 찾으면 맨 뒤
    prs_new = _re.sub(r'(</p:sldIdLst>)', f'  {new_sld}\n  \\1', prs_raw, count=1)
    prs_path.write_text(prs_new, encoding="utf-8")


def _register_slide_content_type(work_dir: Path, slide_filename: str) -> None:
    """복사된 슬라이드를 [Content_Types].xml에 등록한다."""
    import re as _re
    ct_path = work_dir / "unpacked" / "[Content_Types].xml"
    if not ct_path.exists():
        return
    ct_raw = ct_path.read_text(encoding="utf-8")
    slide_ct = (
        'application/vnd.openxmlformats-officedocument.presentationml.slide+xml'
    )
    new_entry = (
        f'<Override PartName="/ppt/slides/{slide_filename}" '
        f'ContentType="{slide_ct}"/>'
    )
    if slide_filename not in ct_raw:
        ct_new = ct_raw.replace("</Types>", f"  {new_entry}\n</Types>")
        ct_path.write_text(ct_new, encoding="utf-8")


def _trim_to_plan_slides(work_dir: Path, plan: dict) -> None:
    """
    plan.json에 명시된 슬라이드만 남기고, plan의 index 순서대로 재정렬한다.
    순서 보장이 핵심 — 기존에는 템플릿 원본 순서를 따라 plan 순서와 불일치했음.
    """
    prs_path = work_dir / "unpacked" / "ppt" / "presentation.xml"
    rels_path = work_dir / "unpacked" / "ppt" / "_rels" / "presentation.xml.rels"
    if not prs_path.exists() or not rels_path.exists():
        return

    # plan의 index 순서대로 파일명 목록 (순서 보장용)
    ordered_files = [s.get("template_file", "") for s in
                     sorted(plan.get("slides", []), key=lambda x: x.get("index", 0))]
    plan_files = set(ordered_files)
    if not plan_files:
        return

    prs_raw  = prs_path.read_text(encoding="utf-8")
    rels_raw = rels_path.read_text(encoding="utf-8")

    # rId → 파일명 매핑
    rid_to_file: dict[str, str] = {}
    for m in re.finditer(r'<Relationship\b([^>]+)/>', rels_raw):
        attrs = m.group(1)
        rid = re.search(r'\bId="(rId\d+)"', attrs)
        tgt = re.search(r'\bTarget="[^"]*/(slide[\w]+\.xml)"', attrs)
        typ = re.search(r'\bType="([^"]+)"', attrs)
        if rid and tgt and typ and 'slide' in typ.group(1) and 'Layout' not in typ.group(1):
            rid_to_file[rid.group(1)] = tgt.group(1)

    file_to_rid = {v: k for k, v in rid_to_file.items()}

    # plan에 없는 sldId 수집
    keep_rids = {file_to_rid[f] for f in plan_files if f in file_to_rid}
    remove_rids: set[str] = set()
    for rid in rid_to_file:
        if rid not in keep_rids:
            remove_rids.add(rid)

    # 불필요한 슬라이드 제거
    new_prs = re.sub(
        r'\s*<p:sldId\b[^/]*/>',
        lambda m: '' if (r := re.search(r'\br:id="(rId\d+)"', m.group(0)))
                        and r.group(1) in remove_rids else m.group(0),
        prs_raw,
    )
    # ── plan의 index 순서대로 sldIdLst 재정렬 ──────────────────
    # 기존 순서는 템플릿 원본 순서를 따름 → plan 순서와 불일치
    existing_slds: list[str] = re.findall(r'\s*<p:sldId\b[^/]*/>', new_prs)
    rid_to_sld = {}
    for sld in existing_slds:
        r_m = re.search(r'\br:id="(rId\d+)"', sld)
        if r_m:
            rid_to_sld[r_m.group(1)] = sld

    # plan 순서대로 sldId 태그를 재배열
    ordered_slds = []
    for fname in ordered_files:
        rid = file_to_rid.get(fname)
        if rid and rid in rid_to_sld:
            ordered_slds.append(rid_to_sld[rid].strip())

    # sldIdLst 내부를 ordered_slds로 교체
    new_sld_block = "\n    " + "\n    ".join(ordered_slds) + "\n  "
    new_prs2 = re.sub(
        r'(<p:sldIdLst>).*?(</p:sldIdLst>)',
        rf'\1{new_sld_block}\2',
        new_prs, flags=re.DOTALL, count=1
    )
    prs_path.write_text(new_prs2, encoding="utf-8")
    print(f"  ✓ 슬라이드 트리밍 + 순서 재정렬: {len(plan_files)}장 (plan index 순서 보장)")


# ── 5. 클린 & 패킹 ───────────────────────────────────────────

def pack_output(work_dir: Path, output_path: Path,
                skip_validation: bool = False) -> bool:
    """clean.py → pack.py 순서로 실행해 pptx를 생성한다."""
    scripts = SKILL_DIR / "scripts"
    template = work_dir / "template.pptx"
    unpacked = work_dir / "unpacked"

    result = subprocess.run(
        [sys.executable, str(scripts / "clean.py"), str(unpacked)],
        capture_output=True, text=True,
    )
    if result.returncode != 0:
        print(f"  ⚠ clean.py stderr: {result.stderr}", file=sys.stderr)

    cmd = [
        sys.executable, str(scripts / "office" / "pack.py"),
        str(unpacked), str(output_path),
        "--original", str(template),
    ]
    if skip_validation:
        cmd += ["--validate", "false"]

    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"  ✗ pack.py 실패:\n{result.stderr}", file=sys.stderr)
        return False

    print(f"  ✓ 패킹 완료: {output_path}")
    restructure_sections(output_path)
    return True


# ── 5. 시각 QA ───────────────────────────────────────────────

def _pdf_via_powerpoint(pptx_path: Path, pdf_path: Path) -> bool:
    """
    macOS AppleScript으로 PowerPoint에 PDF 내보내기를 요청한다.
    실패 시 2초 대기 후 1회 재시도.
    PowerPoint 없으면 False 반환.

    실패 원인 분석 (2026-06-02):
    - PowerPoint 자체는 연속 호출에도 정상 동작 (5초/회)
    - 간헐적 실패는 Vision Fix 루프에서 pack.py 직후 즉시 호출 시
      이전 PowerPoint 프로세스가 완전히 정리되기 전에 새 요청이 들어오는 경우
    - 해결책: 실패 시 2초 대기 후 재시도
    """
    import shutil as _shutil, time as _time
    if not _shutil.which("osascript"):
        return False

    # 문제: PowerPoint에 블로킹 다이얼로그가 있으면 AppleScript 전체가 차단됨
    # 해결: pkill로 완전 종료 → 새 인스턴스로 열기 → PDF 저장
    # open 명령은 로딩 완료를 기다리지 않으므로 동적 대기 필요 (최대 30초)

    # macOS 샌드박스 제한 해결:
    # PowerPoint는 ~/.ppt-skill/runs/... 같은 숨김 경로에 접근 시
    # "Grant File Access" 다이얼로그를 띄움 → 무인 파이프라인에서 타임아웃
    # 해결: result 폴더 안 tmp/ 에 복사 후 변환 (사용자가 PowerPoint 접근 허용한 위치)
    import shutil as _shutil2

    result_dir = Path.home() / "Documents" / "claude" / "ppt-skill" / "result"
    tmp_dir  = result_dir / "tmp"
    tmp_dir.mkdir(parents=True, exist_ok=True)
    tmp_pptx = tmp_dir / "_ppt_qa_tmp.pptx"
    tmp_pdf  = tmp_dir / "_ppt_qa_tmp.pdf"

    try:
        _shutil2.copy2(str(pptx_path), str(tmp_pptx))
    except Exception as e:
        print(f"  ⚠ 임시 파일 복사 실패: {e}")
        return False

    script = f"""
on run
    tell application "Microsoft PowerPoint"
        repeat while (count of presentations) > 0
            close presentation 1 saving no
        end repeat
        open (POSIX file "{tmp_pptx}")
        set maxWait to 30
        set waited to 0
        repeat while (count of presentations) = 0 and waited < maxWait
            delay 1
            set waited to waited + 1
        end repeat
        if (count of presentations) = 0 then
            error "프레젠테이션 열기 실패"
        end if
        save active presentation in (POSIX file "{tmp_pdf}") as save as PDF
        close active presentation saving no
    end tell
end run
"""
    try:
        result = subprocess.run(
            ["osascript", "-e", script],
            capture_output=True, text=True, timeout=180,
        )
        if result.returncode == 0 and tmp_pdf.exists():
            # 결과를 원래 경로로 이동
            _shutil2.move(str(tmp_pdf), str(pdf_path))
            return True
        if result.stderr:
            print(f"  ⚠ PowerPoint 오류: {result.stderr[:100]}")
    except subprocess.TimeoutExpired:
        print("  ⚠ PowerPoint 타임아웃 (180s)")
    finally:
        # 임시 파일 정리
        tmp_pptx.unlink(missing_ok=True)
        tmp_pdf.unlink(missing_ok=True)

    return False


def visual_qa(work_dir: Path, output_path: Path) -> list[str]:
    """
    PowerPoint로 PDF 변환 후 pdftoppm으로 슬라이드 이미지를 생성한다.
    LibreOffice 폴백 없음 — PowerPoint만 사용 (정확한 렌더링 보장).
    PowerPoint 실패 시 빈 리스트 반환 (QA 건너뜀).
    """
    import shutil as _shutil

    pdf_path = work_dir / "output.pdf"

    # PowerPoint (유일한 QA 소스)
    if not _shutil.which("osascript"):
        print("  ⚠ osascript 없음 (macOS 전용) — 시각 QA 건너뜀")
        return []

    ok = _pdf_via_powerpoint(output_path, pdf_path)
    if not ok:
        print("  ⚠ PowerPoint PDF 변환 실패 — 시각 QA 건너뜀")
        print("    (원인: PowerPoint 미설치, 타임아웃, 또는 파일 접근 오류)")
        return []

    print("  ✓ PowerPoint PDF 변환 완료")
    qa_source = "powerpoint"

    # 이미지 변환
    slides_img_dir = work_dir / "qa_images"
    slides_img_dir.mkdir(exist_ok=True)
    prefix = str(slides_img_dir / "slide")
    try:
        subprocess.run(
            ["pdftoppm", "-jpeg", "-r", "120", str(pdf_path), prefix],
            capture_output=True, text=True,
        )
    except FileNotFoundError:
        print("  ⚠ pdftoppm 없음 — QA 이미지 생성 건너뜀 (brew install poppler 권장)")
        return []
    images = sorted(slides_img_dir.glob("*.jpg"))
    print(f"  ✓ QA 이미지 {len(images)}장 ({qa_source}): {slides_img_dir}")

    qa_report = {
        "images": [str(p) for p in images],
        "pdf": str(pdf_path),
        "source": qa_source,
        "note": "텍스트 오버플로우·요소 겹침·플레이스홀더 잔여를 확인하세요.",
    }
    (work_dir / "qa_report.json").write_text(
        __import__("json").dumps(qa_report, ensure_ascii=False, indent=2)
    )
    return [str(p) for p in images]


# ── 6. 섹션 정리 ─────────────────────────────────────────────

def restructure_sections(output_path: Path) -> None:
    """
    생성된 pptx의 섹션을 표지/목차·간지/본문/마무리 4개로 정리한다.
    - 가이드·변화·타임라인·흐름 섹션: 제거 (슬라이드는 이미 없음)
    - 그래프·차트 등 나머지 섹션: 슬라이드를 본문으로 이동 후 제거
    - 매번 생성되는 PPT에 자동 적용된다.
    """
    import re as _re
    import shutil as _shutil

    KEEP = {'표지', '목차/간지', '본문', '마무리'}
    MOVE_TO_BODY = set()   # 본문으로 이동
    DELETE = set()          # 그냥 삭제

    tmp = output_path.with_suffix(".tmp.pptx")
    try:
        with zipfile.ZipFile(output_path) as zin:
            if 'ppt/presentation.xml' not in zin.namelist():
                return
            prs_raw = zin.read('ppt/presentation.xml').decode('utf-8')

        # 섹션 이름 수집
        all_names = _re.findall(r'<p14:section\b[^>]*name="([^"]+)"', prs_raw)
        for name in all_names:
            if name not in KEEP:
                MOVE_TO_BODY.add(name)

        if not MOVE_TO_BODY:
            return   # 이미 정리된 상태

        def get_section(xml, name):
            pat = rf'<p14:section\b[^>]*name="{_re.escape(name)}"[^>]*>.*?</p14:section>'
            m = _re.search(pat, xml, _re.DOTALL)
            return m.group(0) if m else ''

        prs_new = prs_raw

        # 비정규 섹션의 sldId를 본문으로 이동
        move_tags = []
        for name in MOVE_TO_BODY:
            sec = get_section(prs_new, name)
            move_tags += _re.findall(r'<p14:sldId\b[^/]*/>', sec)

        if move_tags:
            body_sec = get_section(prs_new, '본문')
            if body_sec:
                body_updated = _re.sub(
                    r'(</p14:sldIdLst>)',
                    ''.join(move_tags) + r'\1',
                    body_sec, count=1,
                )
                prs_new = prs_new.replace(body_sec, body_updated, 1)

        # 비정규 섹션 제거
        for name in MOVE_TO_BODY:
            sec = get_section(prs_new, name)
            if sec:
                prs_new = prs_new.replace(sec, '')

        with zipfile.ZipFile(output_path) as zin, \
             zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename == 'ppt/presentation.xml':
                    zout.writestr(item, prs_new.encode('utf-8'))
                else:
                    zout.writestr(item, zin.read(item.filename))

        _shutil.move(str(tmp), str(output_path))
        remaining = _re.findall(r'<p14:section\b[^>]*name="([^"]+)"', prs_new)
        print(f"  ✓ 섹션 정리 완료: {remaining}")

    except Exception as e:
        print(f"  ⚠ 섹션 정리 실패 (무시): {e}")
        if tmp.exists():
            tmp.unlink()


# ── 7. 폰트 컴플라이언스 검사 ─────────────────────────────────
_TEMPLATE_FONT = "Pretendard"   # GS Neotek 템플릿 기본 폰트

def check_font_compliance(slides_dir: Path, font: str = _TEMPLATE_FONT) -> list[str]:
    """
    모든 슬라이드 XML을 검사해 지정 폰트(기본: Pretendard) 미사용 shape를 보고.
    반환: [경고 문자열 목록]
    """
    warnings = []
    if not slides_dir.exists():
        return warnings
    for xml_path in sorted(slides_dir.glob("slide*.xml")):
        if "Layout" in xml_path.name or "Master" in xml_path.name:
            continue
        try:
            tree = ET.parse(xml_path); root = tree.getroot()
        except ET.ParseError:
            continue
        for sp in root.iter(f"{{{_NS_P}}}sp"):
            cpr = sp.find(f"{{{_NS_P}}}nvSpPr/{{{_NS_P}}}cNvPr")
            sid = cpr.get("id","?") if cpr is not None else "?"
            for rPr in sp.findall(f".//{{{_NS_A}}}rPr"):
                for tag in ("latin", "ea", "cs"):
                    elem = rPr.find(f"{{{_NS_A}}}{tag}")
                    if elem is not None:
                        tf = elem.get("typeface","")
                        if tf and tf not in (font, "+mj-lt", "+mn-lt", "+mj-ea", "+mn-ea"):
                            warnings.append(
                                f"{xml_path.name} shape {sid}: {tag}={tf!r} (expected {font})"
                            )
    return warnings


# ── 8. 콘텐츠 검증 ───────────────────────────────────────────

def verify_content(output_path: Path) -> list[str]:
    """pptx 텍스트를 추출해 플레이스홀더 잔여 여부를 확인한다."""
    extract_script = """
import sys
from pptx import Presentation
prs = Presentation(sys.argv[1])
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    print(t)
"""
    result = subprocess.run(
        [sys.executable, "-c", extract_script, str(output_path)],
        capture_output=True, text=True,
    )
    issues = [
        line.strip()
        for line in result.stdout.splitlines()
        if _PLACEHOLDER_RE.search(line)
    ]
    return issues


# ── Plan 검증 ────────────────────────────────────────────────

# planning_constraints에서 파싱한 금지/허용 슬라이드 목록
# 사용 불가 슬라이드 (이미지 전용, 차트 고정, 구조적 편집 불가)
_BANNED_SLIDES: set[str] = {
    # 섹션 구분: 명시적 요청 시에만 (section role 자동 선택 금지)
    "slide8.xml",
    # 이미지가 화면 대부분을 차지해 텍스트 설명으로 대체가 어려운 레이아웃
    "slide18.xml", "slide19.xml", "slide20.xml", "slide23.xml",
    "slide21.xml", "slide22.xml",
    # layout_variants 자동 선택 전용 (LLM 직접 지정 금지)
    "slide28.xml",
    # (slide40/41/43 차트는 Excel 연동 완성으로 un-ban — _update_chart_data가 데이터 주입)
}

# ── 콘텐츠 유형별 슬라이드 카탈로그 ───────────────────────────────
# Claude가 콘텐츠를 분석한 후 아래 목록에서 template_file을 선택
_ALLOWED_CONTENT_SLIDES: list[str] = [
    # 텍스트 위주 (가장 안전 — 이미지/도형 불필요, 미배정 시 기본 폴백)
    "slide32.xml",  # ✅ 상단 텍스트 + 하단 3열 콘텐츠 (텍스트/이미지 삽입형)
    # 프로세스/흐름
    "slide38.xml",  # ✅ 3행 흐름도 (keyword→Solution→Service) — flow 표준
    "slide30.xml",  # ✅ 4단계 스텝 프로세스
    # 시간축 레이아웃 (반드시 연도/분기 시계열 콘텐츠일 때만)
    "slide29.xml",  # ✅ 연도별/월별 타임라인 (2023→2026)
    "slide31.xml",  # ✅ 분기별 Q1→Q2→Q3→Q4
    "slide33.xml",  # ✅ 분기별 (상단 설명+Q4열, slide31 변형)
    # 비교/분석
    "slide36.xml",  # 🟡 As-is/To-be 벤다이어그램
    # 카드/아이콘 (아이콘 영역은 비워둠 — placeholder 텍스트 미노출)
    "slide13.xml",  # 🟡 3열 아이콘카드 (3가지 기능/특징/장점)
    "slide15.xml",  # 🟡 3열 대형아이콘 (3가지 핵심 가치)
    "slide14.xml",  # 🟡 4열 아이콘카드 + Insight
    "slide16.xml",  # 🟡 4열 아이콘 컴팩트
    # 이미지+텍스트 (이미지 자리는 '어떤 이미지' 설명 텍스트로 채움)
    "slide9.xml",   # 3열 이미지+제목+설명
    "slide10.xml",  # 3열 이미지+제목+설명 + Insight 배너
    "slide11.xml",  # 3행 이미지(좌)+설명
    "slide12.xml",  # 3열 이미지+우측 텍스트
    "slide17.xml",  # 2x2 이미지+설명
    "slide24.xml",  # ✅ 2블록 텍스트 (bullets + body)
    "slide25.xml",  # ✅ 이미지+우측3열
    "slide26.xml",  # ✅ 이미지+3항목
    "slide27.xml",  # ✅ 이미지+우측3행 (body_title ≤2줄; slide28로 자동 전환 가능)
    "slide35.xml",  # ✅ Before→After 비교
    "slide34.xml",  # ✅ 2이미지+키워드×4 버블 개념도
    "slide37.xml",  # ✅ 3구역 텍스트 비교
    "slide39.xml",  # ✅ 4열 흐름도 (keyword/solution/detail/service)
    "slide42.xml",  # ✅ 대형 본문 텍스트 박스
    # slide28: layout_variants로 자동 선택 (LLM 직접 지정 금지)
    # 차트 레이아웃 (Excel 연동 완성 — content.chart_data로 캐시+임베디드 xlsx 갱신)
    "slide40.xml",  # ✅ 도넛 차트 3개 (KPI/비율)
    "slide41.xml",  # ✅ 막대 차트 (다범주 비교)
    "slide43.xml",  # ✅ 막대 차트 (시계열/증감)
]

# 레이아웃별 필수 콘텐츠 필드 — 하나도 없으면 텍스트 배너(slide32)로 리맵.
# 빈 타임라인 막대·빈 카드·빈 흐름도가 배포되는 것을 코드 레벨에서 차단한다.
_LAYOUT_CONTENT_REQ: dict[str, list[str]] = {
    "slide29.xml": ["periods"],
    "slide31.xml": ["quarters"],
    "slide33.xml": ["quarters"],
    "slide30.xml": ["steps"],
    "slide40.xml": ["chart_data"],
    "slide41.xml": ["chart_data"],
    "slide43.xml": ["chart_data"],
    "slide38.xml": ["keywords"],
    "slide35.xml": ["before", "after"],
    "slide36.xml": ["as_is", "to_be"],
    "slide13.xml": ["items"],
    "slide15.xml": ["items"],
    # slide9-12, slide14, slide16-17, slide25-26, slide34, slide37, slide42 는
    # 에디터 내부에서 items/bullets/descriptions 폴백 처리를 하므로 가드 불필요.
}


def _has_content_field(content: dict, fields: list[str]) -> bool:
    """content(또는 content.body) 안에 주어진 필드 중 하나라도 값이 있으면 True."""
    body = content.get("body") if isinstance(content.get("body"), dict) else {}
    for f in fields:
        if content.get(f) or (isinstance(body, dict) and body.get(f)):
            return True
    return False

_CLOSING_SLIDES: list[str] = ["slide46.xml", "slide44.xml", "slide45.xml"]
_CONTENT_TITLE_ID   = "8"
_CONTENT_BULLETS_ID = "11"  # 호환성 유지


def _resolve_layout_variants(plan: dict) -> list[str]:
    """slide_catalog.json layout_variants 규칙에 따라 plan의 template_file을 자동 교체.
    예: slide27.xml + section_title ≥3줄 → slide28.xml 자동 전환.
    반환: 변경 로그"""
    import math as _math
    cat = _load_slide_catalog()
    variants = cat.get("layout_variants", {})
    if not variants:
        return []

    changes: list[str] = []
    for slide in plan.get("slides", []):
        tmpl = slide.get("template_file", "")
        rule = variants.get(tmpl)
        if not rule:
            continue

        content = slide.get("content", {})
        body = content.get("body", {})
        field = rule.get("selection_field", "section_title")
        text = (content.get(field)
                or (body.get(field) if isinstance(body, dict) else "")
                or "")

        method = rule.get("selection_method", "line_count")
        if method == "line_count":
            chars_per_line = rule.get("chars_per_line", 15)
            n_lines = _math.ceil(len(str(text)) / chars_per_line) if chars_per_line > 0 and text else 1

            selected = tmpl
            for threshold in rule.get("thresholds", []):
                min_l = threshold.get("min_lines", 0)
                max_l = threshold.get("max_lines", 9999)
                if min_l <= n_lines <= max_l:
                    selected = threshold.get("use", tmpl)
                    break

            if selected != tmpl:
                changes.append(
                    f"slide {slide['index']}: {tmpl} → {selected} "
                    f"(section_title {n_lines}줄)"
                )
                slide["template_file"] = selected

    return changes


def enforce_plan_constraints(plan: dict, slide_info: list[dict]) -> tuple[dict, list[str]]:
    """
    생성된 plan의 template_file이 금지 목록에 있거나 중복되면 자동 교체한다.
    planning_constraints를 Claude에게만 맡기지 않고 코드 레벨에서 강제.
    반환: (수정된 plan, 변경 로그)
    """
    # slide_catalog.json 우선, 없으면 하드코딩 폴백
    _cat = _load_slide_catalog()
    _banned    = set(_cat["banned_slides"].keys())    if _cat.get("banned_slides")         else _BANNED_SLIDES
    # verified + unverified 합집합. 구버전 키(allowed_content_slides) 폴백 유지
    if _cat.get("verified_slides") or _cat.get("unverified_slides"):
        _allowed = (list(_cat.get("verified_slides", {}).keys()) +
                    list(_cat.get("unverified_slides", {}).keys()))
    else:
        _allowed = _cat.get("allowed_content_slides", _ALLOWED_CONTENT_SLIDES)
    _cont_req  = _cat.get("layout_content_req", {})   or _LAYOUT_CONTENT_REQ

    available_files = {s["file"] for s in slide_info}
    allowed = [f for f in _allowed if f in available_files]
    if not allowed:
        allowed = [f for f in available_files
                   if f not in _banned
                   and f not in {"slide6.xml", "slide7.xml", "slide9.xml"}]

    MAX_REPEAT = 3  # 동일 레이아웃 최대 사용 횟수 (cover/toc/closing 제외)
    used_files: set[str] = set()
    used_counts: dict[str, int] = {}
    changes: list[str] = []
    allowed_pool = list(allowed)
    pool_idx = 0

    for slide in plan.get("slides", []):
        role = slide.get("role", "")
        tmpl = slide.get("template_file", "")

        # cover/toc는 전용 파일 유지
        if role in ("cover", "toc"):
            used_files.add(tmpl)
            continue

        # closing은 전용 마무리 슬라이드로 강제
        if role == "closing":
            for cs in _CLOSING_SLIDES:
                if cs in available_files and cs not in used_files:
                    if cs != tmpl:
                        changes.append(f"slide {slide['index']}: {tmpl} → {cs} (closing 전용)")
                        slide["template_file"] = cs
                        tmpl = cs
                    break
            used_files.add(tmpl)
            continue

        # 금지 목록이거나 허용 목록에 없으면 교체
        # Claude가 직접 template을 선택했으면 존중, 금지된 경우만 교체
        # section role은 사용자 명시 요청 없이는 일반 content로 변환
        if role == "section" and tmpl == "slide8.xml":
            role = "content"
            slide["role"] = "content"

        needs_replace = (tmpl in _banned) or (
            tmpl not in _allowed
            and role not in ("cover", "toc", "closing"))

        # 동일 레이아웃 MAX_REPEAT 초과 시 교체 강제
        base_tmpl = __import__("re").sub(r"_c\d+(?=\.xml$)", "", tmpl)
        if not needs_replace and used_counts.get(base_tmpl, 0) >= MAX_REPEAT:
            needs_replace = True
            changes.append(f"slide {slide['index']}: {tmpl} 반복 {used_counts[base_tmpl]}회 초과 → 교체")

        if needs_replace:
            replacement = None
            for candidate in _allowed:
                if candidate in available_files and used_counts.get(candidate, 0) < MAX_REPEAT:
                    replacement = candidate
                    break
            if replacement is None:
                # 모든 허용 레이아웃이 MAX_REPEAT 초과 시 가장 적게 쓴 것으로 교체
                replacement = min(
                    (c for c in _allowed if c in available_files),
                    key=lambda c: used_counts.get(c, 0),
                    default=None,
                )
            if replacement:
                changes.append(f"slide {slide['index']}: {tmpl} → {replacement}"
                                + f" (금지/미허용/반복 → 교체)")
                slide["template_file"] = replacement
                tmpl = replacement
            else:
                changes.append(f"slide {slide['index']}: 대체 슬라이드 없음 ({tmpl} 유지)")

        # ── 레이아웃-콘텐츠 적합성 가드 ──
        # 선택된 레이아웃의 필수 콘텐츠가 없으면 텍스트 배너(slide32)로 리맵.
        # 예) 서술형 내용에 연도 타임라인(slide29) 배정 → 빈 막대 방지.
        req = _cont_req.get(tmpl)
        if req and not _has_content_field(slide.get("content", {}), req):
            fallback = "slide32.xml" if "slide32.xml" in available_files else (
                "slide24.xml" if "slide24.xml" in available_files else None)
            if fallback and fallback != tmpl:
                changes.append(
                    f"slide {slide['index']}: {tmpl} → {fallback} "
                    f"(필수 콘텐츠 {req} 없음 → 텍스트 레이아웃)")
                slide["template_file"] = fallback
                tmpl = fallback

        used_files.add(tmpl)
        base_key = __import__("re").sub(r"_c\d+(?=\.xml$)", "", tmpl)
        used_counts[base_key] = used_counts.get(base_key, 0) + 1

    return plan, changes


def validate_plan(plan: dict) -> tuple[bool, list[str]]:
    """
    plan.json의 각 슬라이드가 role별 필수 content 필드를 가졌는지 검증.
    반환: (ok, warnings)
    """
    warnings: list[str] = []
    for slide in plan.get("slides", []):
        role = slide.get("role", "content")
        content = slide.get("content", {})
        schema = PLAN_CONTENT_SCHEMA.get(role, {})

        for field in schema.get("required", []):
            if field not in content or not content[field]:
                warnings.append(
                    f"slide {slide['index']} ({role}): 필수 필드 '{field}' 없음"
                )
        for field in schema.get("recommended", []):
            if field not in content or not content[field]:
                warnings.append(
                    f"slide {slide['index']} ({role}): 권장 필드 '{field}' 없음 — 폴백 사용"
                )
        # 본문 슬라이드(cover/toc/closing 제외)는 슬라이드 레벨 subtitle 필수
        if role not in ("cover", "toc", "closing"):
            if not slide.get("subtitle", "").strip():
                warnings.append(
                    f"slide {slide['index']} ({role}): 필수 — 슬라이드 레벨 'subtitle'(중제목) 없음"
                )

    ok = not any("필수" in w for w in warnings)
    return ok, warnings


# ── known_fixes 자동 적용 ────────────────────────────────────

def apply_known_fixes(xml_path: Path, slide_plan: dict, known_fixes: list) -> None:
    """
    long_term_memory의 known_failure_fixes 패턴을 슬라이드 편집 전에 검사·적용.
    현재 적용 가능한 픽스:
      - endParaRPr 순서 보정 (run_after_endpararpr)
      - lang=ko-KR 강제 (korean_in_en_us_run)
    """
    if not xml_path.exists() or not known_fixes:
        return

    ns_a = _NS_A
    try:
        tree = ET.parse(xml_path)
        root = tree.getroot()
        fixed = False

        for para in root.iter(f"{{{ns_a}}}p"):
            children = list(para)
            tags = [c.tag.split("}")[-1] for c in children]

            # Fix 1: <a:r>가 <a:endParaRPr> 뒤에 있으면 앞으로 이동
            if "endParaRPr" in tags and "r" in tags:
                end_idx = tags.index("endParaRPr")
                runs_after = [i for i, t in enumerate(tags) if t == "r" and i > end_idx]
                if runs_after:
                    end_elem = children[end_idx]
                    for i in sorted(runs_after, reverse=True):
                        r_elem = children[i]
                        para.remove(r_elem)
                        para.insert(end_idx, r_elem)
                    fixed = True

            # Fix 2: lang=en-US run에 한국어 텍스트 → lang=ko-KR
            for r in para.findall(f"{{{ns_a}}}r"):
                rPr = r.find(f"{{{ns_a}}}rPr")
                t   = r.find(f"{{{ns_a}}}t")
                if rPr is not None and t is not None and t.text:
                    if rPr.get("lang") == "en-US" and re.search(r"[가-힣]", t.text):
                        rPr.set("lang", _ppt_lang())
                        fixed = True

        if fixed:
            _write_xml(root, xml_path)

    except ET.ParseError:
        pass  # 편집 함수가 따로 처리


# ── Verifier 실행 파이프라인 ─────────────────────────────────

def _vfy_xml_validity(output_path: Path, work_dir: Path, **_) -> list[dict]:
    issues = []
    slides_dir = work_dir / "unpacked" / "ppt" / "slides"
    for xml_f in sorted(slides_dir.glob("slide*.xml")):
        try:
            ET.parse(xml_f)
        except ET.ParseError as e:
            issues.append({"rule": "xml_validity", "severity": "CRITICAL",
                           "file": xml_f.name, "detail": str(e)})
    return issues


def _vfy_placeholder(output_path: Path, work_dir: Path, **_) -> list[dict]:
    """plan.json에 명시된 슬라이드(실제 편집된 것)만 placeholder 검사."""
    issues = []
    if not output_path.exists():
        return issues

    # plan.json에서 편집된 슬라이드 파일명 집합
    plan_path = work_dir / "plan.json"
    edited_files: set[str] = set()
    if plan_path.exists():
        plan = json.loads(plan_path.read_text())
        edited_files = {s.get("template_file", "") for s in plan.get("slides", [])}

    extract = (
        "import sys, json; from pptx import Presentation; prs=Presentation(sys.argv[1]);"
        "[print(p.text.strip()) for s in prs.slides for sh in s.shapes"
        " if sh.has_text_frame for p in sh.text_frame.paragraphs if p.text.strip()]"
    )
    result = subprocess.run([sys.executable, "-c", extract, str(output_path)],
                            capture_output=True, text=True)

    # plan이 없으면 전체 검사, 있으면 편집된 슬라이드만 검사
    # pptx 슬라이드 텍스트를 파일 단위로 분리하기 어려우므로
    # 실제 편집 슬라이드에서만 발생할 수 있는 패턴을 검사
    PATTERNS = ["lorem", "작성해주세요", "todo", "[insert", "placeholder"]
    seen: set[str] = set()
    for line in result.stdout.splitlines():
        text = line.strip()
        if not text or text in seen:
            continue
        seen.add(text)
        if any(p in text.lower() for p in PATTERNS):
            # 편집 슬라이드에서만 실제 문제가 되는 패턴 (중복 제거)
            issues.append({"rule": "placeholder_check", "severity": "CRITICAL",
                           "detail": text})
            if len(issues) >= 5:  # 최대 5건만 보고 (템플릿 잔여 노이즈 방지)
                break
    return issues


def _vfy_run_order(output_path: Path, work_dir: Path, **_) -> list[dict]:
    issues = []
    ns_a = _NS_A
    slides_dir = work_dir / "unpacked" / "ppt" / "slides"
    for xml_f in sorted(slides_dir.glob("slide*.xml")):
        try:
            tree = ET.parse(xml_f)
            for para in tree.getroot().iter(f"{{{ns_a}}}p"):
                tags = [c.tag.split("}")[-1] for c in para]
                if "endParaRPr" in tags:
                    ei = tags.index("endParaRPr")
                    if any(t == "r" for t in tags[ei+1:]):
                        issues.append({"rule": "run_before_endpararpr",
                                       "severity": "CRITICAL", "file": xml_f.name,
                                       "detail": "<a:r> after <a:endParaRPr>"})
                        break
        except ET.ParseError:
            pass
    return issues


def _vfy_section_structure(output_path: Path, **_) -> list[dict]:
    if not output_path.exists():
        return []
    EXPECTED = {"표지", "목차/간지", "본문", "마무리"}
    try:
        with zipfile.ZipFile(output_path) as z:
            prs = z.read("ppt/presentation.xml").decode("utf-8")
        found = set(re.findall(r'<p14:section\b[^>]*name="([^"]+)"', prs))
        if found != EXPECTED:
            return [{"rule": "section_structure", "severity": "MEDIUM",
                     "detail": f"섹션 불일치: {found} ≠ {EXPECTED}"}]
    except Exception:
        pass
    return []


def _vfy_qa_done(output_path: Path, work_dir: Path, **_) -> list[dict]:
    qa = work_dir / "qa_report.json"
    if not qa.exists() or not json.loads(qa.read_text()).get("images"):
        return [{"rule": "qa_completion_check", "severity": "MEDIUM",
                 "detail": "qa_ok=false — QA 이미지 없음"}]
    return []


def _vfy_toc_paragraph_count(output_path: Path, work_dir: Path, **_) -> list[dict]:
    """slide7(목차) ID=10의 paragraph 수가 7개인지 검증.
    Vision Fix Agent가 set_paragraphs로 줄이면 번호·텍스트·페이지 높이 불일치 발생."""
    slides_dir = work_dir / "unpacked" / "ppt" / "slides"
    slide7 = slides_dir / "slide7.xml"
    if not slide7.exists():
        return []
    try:
        tree = ET.parse(slide7)
        ns_p = _NS_P
        ns_a = _NS_A
        for sp in tree.getroot().findall(f".//{{{ns_p}}}sp"):
            cpr = sp.find(f"{{{ns_p}}}nvSpPr/{{{ns_p}}}cNvPr")
            if cpr is None or cpr.get("id") != "10":
                continue
            paras = sp.findall(f".//{{{ns_a}}}p")
            if len(paras) != 7:
                return [{"rule": "toc_paragraph_count", "severity": "CRITICAL",
                         "detail": f"TOC ID=10 paragraph 수={len(paras)} (7 필요) — "
                                   "Vision Fix Agent가 구조 파괴했을 수 있음"}]
    except ET.ParseError:
        pass
    return []


def _vfy_shape_coverage(output_path: Path, work_dir: Path, **_) -> list[dict]:
    """
    템플릿 원본 shape과 생성 결과를 1:1 비교.
    템플릿에 텍스트가 있던 shape이 생성 후 비어있으면 MEDIUM 위반으로 보고.
    icon/이미지 전용 shape(텍스트="icon")은 제외.
    """
    plan_path = work_dir / "plan.json"
    if not plan_path.exists():
        return []
    plan = json.loads(plan_path.read_text())

    template_pptx = SKILL_DIR / "template" / "2026_PPT Template.pptx"
    if not template_pptx.exists():
        return []

    ns_p = _NS_P
    ns_a = _NS_A
    issues: list[dict] = []

    # 출력 pptx에서 슬라이드 순서 읽기
    try:
        with zipfile.ZipFile(output_path) as oz:
            prs_xml  = etree.fromstring(oz.read("ppt/presentation.xml"))
            rels_xml = etree.fromstring(oz.read("ppt/_rels/presentation.xml.rels"))
    except Exception:
        return []

    NS_R  = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    rid_map = {r.get("Id"): r.get("Target","") for r in rels_xml
               if r.get("Target","").startswith("slides/slide")}
    out_slide_order = [rid_map[s.get(f"{{{NS_R}}}id")]
                       for s in prs_xml.iter(f"{{{NS_P}}}sldId")
                       if s.get(f"{{{NS_R}}}id") in rid_map]

    def _shapes_text(z, slide_path: str) -> dict[str, str]:
        try:
            root = etree.fromstring(z.read(slide_path))
        except Exception:
            return {}
        result = {}
        for sp in root.iter(f"{{{ns_p}}}sp"):
            for el in sp.iter():
                if el.tag.endswith("}cNvPr"):
                    sid = el.get("id", "?"); break
            txBody = sp.find(f".//{{{ns_p}}}txBody")
            if txBody is None: continue
            text = "".join(t.text or "" for t in txBody.iter(f"{{{ns_a}}}t")).strip()
            result[sid] = text
        return result

    slides_in_plan = plan.get("slides", [])
    with zipfile.ZipFile(template_pptx) as tz, zipfile.ZipFile(output_path) as oz:
        for slide_plan, out_slide_rel in zip(slides_in_plan, out_slide_order):
            role = slide_plan.get("role", "")
            if role in ("cover", "toc", "closing"):
                continue
            tmpl_file = slide_plan.get("template_file", "")
            # _cN suffix 제거해 기본 템플릿 파일명 찾기
            import re as _re2
            base_file = _re2.sub(r"_c\d+(?=\.xml$)", "", tmpl_file)
            tmpl_path = f"ppt/slides/{base_file}"
            if tmpl_path not in tz.namelist():
                continue

            tmpl_shapes = _shapes_text(tz, tmpl_path)
            out_shapes  = _shapes_text(oz, f"ppt/{out_slide_rel}")

            for sid, orig_text in tmpl_shapes.items():
                if not orig_text or orig_text.lower() == "icon":
                    continue  # 원래 비어있거나 icon 전용 shape은 검사 제외
                out_text = out_shapes.get(sid, "").strip()
                if not out_text:
                    issues.append({
                        "rule": "shape_coverage",
                        "severity": "MEDIUM",
                        "detail": (f"슬라이드 {slide_plan.get('index','?')} "
                                   f"({tmpl_file}) ID={sid}: "
                                   f"템플릿에 있던 텍스트 사라짐 "
                                   f"(원본: {orig_text[:30]!r})")
                    })
    return issues


def _vfy_zone1_subtitle(output_path: Path, work_dir: Path, **_) -> list[dict]:
    """
    본문 슬라이드(cover/toc/closing 제외)의 plan subtitle 필드 누락 검사.
    subtitle 없으면 Zone1 ID=9(중제목)이 빈 채로 출력됨.
    """
    plan_path = work_dir / "plan.json"
    if not plan_path.exists():
        return []
    plan = json.loads(plan_path.read_text())
    issues: list[dict] = []
    for slide in plan.get("slides", []):
        role = slide.get("role", "content")
        if role in ("cover", "toc", "closing"):
            continue
        if not slide.get("subtitle", "").strip():
            issues.append({
                "rule": "zone1_subtitle_empty",
                "severity": "HIGH",
                "detail": (f"slide {slide.get('index','?')} ({slide.get('template_file','')}) "
                           f"subtitle(중제목) 없음 — Zone1 ID=9 비어있음. "
                           f"plan에 콘텐츠에 맞는 subtitle 추가 필요."),
            })
    return issues


_VERIFIER_REGISTRY: dict[str, object] = {
    "xml_validity":          _vfy_xml_validity,
    "placeholder_check":     _vfy_placeholder,
    "run_before_endpararpr": _vfy_run_order,
    "section_structure":     _vfy_section_structure,
    "qa_completion_check":   _vfy_qa_done,
    "toc_paragraph_count":   _vfy_toc_paragraph_count,
    "shape_coverage":        _vfy_shape_coverage,
    "zone1_subtitle_empty":  _vfy_zone1_subtitle,
}


def execute_verifier_rules(output_path: Path, work_dir: Path) -> list[dict]:
    """
    ~/.ppt-skill/harness/verifier_rules.json 에 정의된 규칙을 실행하고
    위반 목록을 반환한다. CRITICAL 위반이 있으면 호출자가 재시도 트리거.
    """
    rules_path = SKILL_DIR / "harness" / "verifier_rules.json"
    if not rules_path.exists():
        return []

    rules_cfg = json.loads(rules_path.read_text()).get("rules", {})
    all_issues: list[dict] = []

    for rule_name, cfg in rules_cfg.items():
        fn = _VERIFIER_REGISTRY.get(rule_name)
        if fn is None:
            continue
        try:
            issues = fn(output_path=output_path, work_dir=work_dir)
            # 규칙 설정의 severity 우선 (없으면 함수 기본값)
            for iss in issues:
                iss.setdefault("severity", cfg.get("severity", "MEDIUM"))
            all_issues.extend(issues)
        except Exception as e:
            print(f"  ⚠ verifier [{rule_name}] 실행 오류: {e}")

    critical = [i for i in all_issues if i["severity"] == "CRITICAL"]
    if critical:
        print(f"  ✗ CRITICAL 위반 {len(critical)}건:")
        for i in critical[:3]:
            print(f"    [{i['rule']}] {i['detail']}")
    elif all_issues:
        print(f"  ⚠ verifier 경고 {len(all_issues)}건 (non-critical)")
    else:
        print("  ✓ 모든 verifier 규칙 통과")

    return all_issues


# ── 슬라이드별 전용 편집기 ────────────────────────────────────────

def _edit_slide8(xml_path: Path, slide_plan: dict) -> None:
    """slide8 (섹션 구분): ID=2(섹션 제목 40pt), ID=4(서브항목 목록), ID=5(번호)."""
    import copy as _copy
    ns_p, ns_a = _NS_P, _NS_A
    title   = slide_plan.get("title", "")
    content = slide_plan.get("content", {})
    items   = content.get("items") or content.get("bullets") or []

    try:
        tree = ET.parse(xml_path)
        root = tree.getroot()
    except ET.ParseError:
        return

    def _set(sp, text):
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None: return
        first_r = sp.find(f".//{{{ns_a}}}r")
        orig_rPr = None
        if first_r is not None:
            rPr_e = first_r.find(f"{{{ns_a}}}rPr")
            if rPr_e is not None:
                orig_rPr = _copy.deepcopy(rPr_e)
                orig_rPr.set("lang", _ppt_lang()); orig_rPr.set("dirty", "0")
        for p in txBody.findall(f"{{{ns_a}}}p"):
            for r in p.findall(f"{{{ns_a}}}r"): p.remove(r)
            end = p.find(f"{{{ns_a}}}endParaRPr")
            idx = list(p).index(end) if end is not None else len(p)
            r_new = ET.Element(f"{{{ns_a}}}r")
            if orig_rPr is not None: r_new.append(_copy.deepcopy(orig_rPr))
            ET.SubElement(r_new, f"{{{ns_a}}}t").text = text
            p.insert(idx, r_new)
            break

    _s8 = _load_slide_shape_ids().get("slide8", {})
    sp2 = _find_shape_by_id(root, _s8.get("section_title_id", "2"))
    if sp2: _set(sp2, title)

    # ID=4: 서브항목 (여러 paragraph)
    sp4 = _find_shape_by_id(root, _s8.get("subitems_id", "4"))
    if sp4 and items:
        txBody = sp4.find(f"{{{ns_p}}}txBody")
        if txBody:
            orig_rPr = None
            for r in sp4.findall(f".//{{{ns_a}}}r"):
                rPr_e = r.find(f"{{{ns_a}}}rPr")
                if rPr_e is not None:
                    orig_rPr = _copy.deepcopy(rPr_e)
                    orig_rPr.set("lang", _ppt_lang()); orig_rPr.set("dirty", "0")
                    break
            paras = txBody.findall(f"{{{ns_a}}}p")
            for i, para in enumerate(paras):
                for r in para.findall(f"{{{ns_a}}}r"): para.remove(r)
                text = items[i] if i < len(items) else ""
                end = para.find(f"{{{ns_a}}}endParaRPr")
                idx = list(para).index(end) if end is not None else len(para)
                r_new = ET.Element(f"{{{ns_a}}}r")
                if orig_rPr: r_new.append(_copy.deepcopy(orig_rPr))
                ET.SubElement(r_new, f"{{{ns_a}}}t").text = text
                para.insert(idx, r_new)

    _clear_residual_placeholders(root)
    _write_xml(root, xml_path)


# ── 템플릿 기본 폰트 (GS Neotek: Pretendard) ─────────────────────
_TEMPLATE_FONT = "Pretendard"

# ── 사이드바 동적 리사이징 상수 (3p slide13 기준) ────────────────
_SIDEBAR_LINE_HEIGHT_EMU = 314_603   # 629206 cy / 2줄
# _TEMPLATE_FONT 는 line 1632 에서 선언됨 (check_font_compliance 앞)
_SIDEBAR_REF_GAP_EMU     = 217_772   # body_desc.y - (body_title.y + body_title.cy)
_SIDEBAR_CHARS_PER_LINE  = 10        # 한국어 ~20pt, cx=2,200,409 EMU 기준


def _count_sidebar_lines(text: str, cx_emu: int = 2_200_409, font_pt=None) -> int:
    """
    body_title 텍스트 줄 수 추정 (cx 기반 실제 폭 계산).
    한글/CJK = font_pt pt, ASCII = font_pt*0.6 pt, 공백 = font_pt*0.35 pt.
    cx_emu: 텍스트박스 폭 (기본값 = 사이드바 실측 2,200,409 EMU = 173.3pt).
    font_pt: None이면 _BODY_TITLE_FONT_PT(20pt) 사용.
    """
    import math
    if font_pt is None:
        font_pt = _BODY_TITLE_FONT_PT

    cx_pt = cx_emu / 12700.0

    def _cw(c: str) -> float:
        if '가' <= c <= '힣' or '一' <= c <= '鿿':
            return font_pt
        if c in (' ', '\t', '\n'):
            return font_pt * 0.35
        return font_pt * 0.6

    w = sum(_cw(c) for c in (text or ""))
    return max(1, math.ceil(w / cx_pt))


def _get_shape_xfrm_elems(sp):
    """shape spPr > xfrm 의 off, ext 엘리먼트 반환"""
    spPr = sp.find(f"{{{_NS_P}}}spPr")
    if spPr is None: return None, None
    xfrm = spPr.find(f"{{{_NS_A}}}xfrm")
    if xfrm is None: return None, None
    return xfrm.find(f"{{{_NS_A}}}off"), xfrm.find(f"{{{_NS_A}}}ext")


def _resize_sidebar_and_reposition_desc(root, label_id: str, desc_id, label_text: str) -> None:
    """
    body_title: 텍스트 줄 수에 맞게 cy 동적 확장 (템플릿값 이상으로만).
    body_desc:
      - cy / y: 템플릿 원본값 유지 (변경하지 않음)
      - 빈 단락 = max(0, 3 - desc_text_lines) 추가 — 글자 높이 정렬
    """
    import copy as _copy
    ns_a = _NS_A
    lbl = _find_shape_by_id(root, label_id)
    if lbl is None: return
    off_lbl, ext_lbl = _get_shape_xfrm_elems(lbl)
    if off_lbl is None or ext_lbl is None: return

    title_cx    = int(ext_lbl.get('cx', 2_200_409))
    title_lines = _count_sidebar_lines(label_text, cx_emu=title_cx)
    title_cy    = title_lines * _SIDEBAR_LINE_HEIGHT_EMU
    cur_cy      = int(ext_lbl.get('cy', 0))
    if title_cy > cur_cy:                  # 넘칠 때만 확장, 템플릿값 이하로 축소 금지
        ext_lbl.set('cy', str(title_cy))

    if not desc_id: return
    desc = _find_shape_by_id(root, desc_id)
    if desc is None: return

    # desc: cy/y는 건드리지 않음 — 빈 단락만 추가
    desc_text = ''.join(el.text or '' for el in desc.iter(f'{{{ns_a}}}t')).strip()
    _, ext_desc = _get_shape_xfrm_elems(desc)
    desc_cx    = int(ext_desc.get('cx', title_cx)) if ext_desc is not None else title_cx
    desc_lines = _count_sidebar_lines(desc_text, cx_emu=desc_cx, font_pt=12.0) if desc_text else 0

    txBody = desc.find(f'{{{_NS_P}}}txBody')
    if txBody is None: return

    for p in list(txBody.findall(f'{{{ns_a}}}p')):
        if p.find(f'{{{ns_a}}}r') is None:
            txBody.remove(p)

    padding = max(0, 3 - desc_lines)
    if padding > 0:
        end_rpr = None
        for p in txBody.findall(f'{{{ns_a}}}p'):
            e = p.find(f'{{{ns_a}}}endParaRPr')
            if e is not None:
                end_rpr = _copy.deepcopy(e); break
        for _ in range(padding):
            p_new = ET.Element(f'{{{ns_a}}}p')
            if end_rpr is not None:
                p_new.append(_copy.deepcopy(end_rpr))
            txBody.append(p_new)


def _resolve_shape_bounds(sp) -> dict | None:
    """shape의 실제 좌표 반환 {x,y,cx,cy}. 레이아웃 상속(cx=0)이면 None."""
    spPr = sp.find(f"{{{_NS_P}}}spPr")
    if spPr is None: return None
    xfrm = spPr.find(f"{{{_NS_A}}}xfrm")
    if xfrm is None: return None
    off = xfrm.find(f"{{{_NS_A}}}off")
    ext = xfrm.find(f"{{{_NS_A}}}ext")
    if off is None or ext is None: return None
    x = int(off.get('x', 0)); y = int(off.get('y', 0))
    cx = int(ext.get('cx', 0)); cy = int(ext.get('cy', 0))
    if cx == 0: return None
    return {'x': x, 'y': y, 'cx': cx, 'cy': cy}


def _boxes_overlap(a: dict, b: dict, margin: int = 0) -> bool:
    """두 박스가 겹치는지 확인 (margin: 최소 유지 간격)."""
    return not (
        a['x'] + a['cx'] + margin <= b['x'] or
        b['x'] + b['cx'] + margin <= a['x'] or
        a['y'] + a['cy'] + margin <= b['y'] or
        b['y'] + b['cy'] + margin <= a['y']
    )


def _resolve_overlaps(root, changed_sid: str, gap: int = 50_000) -> None:
    """
    changed_sid shape의 크기 변동 후 같은 x 열에 있는 인접 shape들을 검사해
    겹침이 발생하면 아래 shape를 밀어내어 겹침을 해소한다.
    gap: 최소 유지 간격 (EMU, 기본 50,000 ≈ 4mm)
    """
    # 변경된 shape 좌표
    changed_sp = _find_shape_by_id(root, changed_sid)
    if changed_sp is None: return
    cb = _resolve_shape_bounds(changed_sp)
    if cb is None: return

    # 모든 shape를 y 순서로 정렬하여 변경된 shape 아래에 있는 것 검사
    ns_p = _NS_P; ns_a = _NS_A
    all_shapes = []
    for sp in root.iter(f'{{{ns_p}}}sp'):
        cpr = sp.find(f'{{{ns_p}}}nvSpPr/{{{ns_p}}}cNvPr')
        if cpr is None: continue
        sid = cpr.get('id', '?')
        if sid == changed_sid: continue
        b = _resolve_shape_bounds(sp)
        if b is None: continue
        all_shapes.append((sid, sp, b))

    # x 범위가 겹치고 y가 변경된 shape 아래에 있는 shape 검사
    for sid, sp, b in all_shapes:
        # x 범위 겹침 확인 (같은 컬럼)
        x_overlap = (b['x'] < cb['x'] + cb['cx']) and (b['x'] + b['cx'] > cb['x'])
        if not x_overlap: continue
        # y가 변경 shape 아래에 있는 경우만
        if b['y'] < cb['y']: continue
        # 겹침 확인
        if _boxes_overlap(cb, b, margin=gap):
            # 아래 shape를 밀어냄
            new_y = cb['y'] + cb['cy'] + gap
            spPr = sp.find(f'{{{ns_p}}}spPr')
            if spPr is None: continue
            xfrm = spPr.find(f'{{{ns_a}}}xfrm')
            if xfrm is None: continue
            off = xfrm.find(f'{{{ns_a}}}off')
            if off is not None:
                off.set('y', str(new_y))
                # 재귀적으로 이 shape도 아래 shape들에 영향 줄 수 있으므로 cb 갱신
                b_new = _resolve_shape_bounds(sp)
                if b_new:
                    b.update(b_new)


_CHAPTER_MAP_CACHE: dict | None = None


def _load_chapter_map() -> dict:
    global _CHAPTER_MAP_CACHE
    if _CHAPTER_MAP_CACHE is not None:
        return _CHAPTER_MAP_CACHE
    p = Path(__file__).parent / "harness" / "chapter_map.json"
    if p.exists():
        _CHAPTER_MAP_CACHE = json.loads(p.read_text()).get("chapters", {})
    else:
        _CHAPTER_MAP_CACHE = {}
    return _CHAPTER_MAP_CACHE


def _infer_chapter_title(slide_plan: dict) -> str:
    """plan.title을 대제목으로 반환. 캐시된 챕터맵이 있으면 section_title 번호 기준 보정."""
    plan_title = slide_plan.get("title", "")
    # _CHAPTER_MAP_CACHE가 실행 중 주입된 경우에만 추론 시도 (chapter_map.json 정적 파일 금지)
    if _CHAPTER_MAP_CACHE:
        content = slide_plan.get("content", {})
        sec_title = content.get("section_title", "")
        if sec_title:
            parts = sec_title.split(".")
            if parts and parts[0].strip().isdigit():
                ch = _CHAPTER_MAP_CACHE.get(parts[0].strip())
                if ch:
                    return ch
    return plan_title


def _apply_common_zones(root, slide_plan: dict, template_file: str) -> None:
    """
    모든 본문 슬라이드 공통 3-Zone 처리:
      Zone 1 헤더 바: ID=8(대제목=목차챕터제목), ID=9(중제목=챕터번호 '01','02'...)
      Zone 2 사이드바: body_title → 본문제목(section_title), body_desc → 본문설명글(section_desc)
      Zone 3 본문 구역: 각 편집기가 직접 처리
    """
    import copy as _copy
    ns_p, ns_a = _NS_P, _NS_A

    # ID=8 대제목: section_title 번호로 챕터 추론. 추론 불가 시 plan.title 사용
    chapter_title = _infer_chapter_title(slide_plan)
    subtitle  = slide_plan.get("subtitle", "")
    content   = slide_plan.get("content", {})
    sec_title = content.get("section_title", "")
    sec_desc  = content.get("section_desc", "")

    # subtitle fallback: LLM 누락 시 section_title에서 숫자 접두사 제거해 도출
    if not subtitle and sec_title:
        subtitle = re.sub(r'^\d+(\.\d+)*[\s.]*', '', sec_title).strip()

    _BODY_TITLE_SEMIBOLD = "Pretendard SemiBold"

    def _set_shape(sid: str, text: str, force_semibold: bool = False) -> None:
        sp = _find_shape_by_id(root, sid)
        if sp is None: return
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None: return
        orig_rPr = None
        for r in sp.findall(f".//{{{ns_a}}}r"):
            rPr_e = r.find(f"{{{ns_a}}}rPr")
            if rPr_e is not None:
                orig_rPr = _copy.deepcopy(rPr_e)
                orig_rPr.set("lang", _ppt_lang())
                orig_rPr.set("dirty", "0")
                # sz 보존 — 원본 명시 크기 유지 (삭제 금지)
                if force_semibold:
                    # body_title은 항상 Pretendard SemiBold 강제
                    for tag in ("latin", "ea", "cs"):
                        el = orig_rPr.find(f"{{{ns_a}}}{tag}")
                        if el is not None:
                            el.set("typeface", _BODY_TITLE_SEMIBOLD)
                        else:
                            ET.SubElement(orig_rPr, f"{{{ns_a}}}{tag}", typeface=_BODY_TITLE_SEMIBOLD)
                break
        if orig_rPr is None and force_semibold:
            # body_title만 폰트 강제 생성 — 일반 shape는 rPr 없으면 테마 상속 유지
            orig_rPr = ET.Element(f"{{{ns_a}}}rPr", dirty="0")
            for tag in ("latin", "ea", "cs"):
                ET.SubElement(orig_rPr, f"{{{ns_a}}}{tag}", typeface=_BODY_TITLE_SEMIBOLD)
        paras = txBody.findall(f"{{{ns_a}}}p")
        for pi, p in enumerate(paras):
            for r in p.findall(f"{{{ns_a}}}r"): p.remove(r)
            end = p.find(f"{{{ns_a}}}endParaRPr")
            idx = list(p).index(end) if end is not None else len(p)
            if pi == 0:
                r_new = ET.Element(f"{{{ns_a}}}r")
                if orig_rPr is not None:
                    r_new.append(_copy.deepcopy(orig_rPr))
                ET.SubElement(r_new, f"{{{ns_a}}}t").text = text
                p.insert(idx, r_new)
            # 첫 번째 이후 paragraph는 run 제거만 (빈 paragraph 유지)

    # Zone 1: 헤더 바 — 챕터 대제목(추론값) 사용
    _set_shape("8", chapter_title)
    hdr = _find_shape_by_id(root, "8")
    if hdr is not None:
        _enable_autofit(hdr)

    if subtitle:
        # 중제목 앞에 대제목 내 순번 prefix: "01 창립 비전..." / "02 안전 배포..." 형식
        # generate_pptx가 주입한 subtitle_seq 우선 사용, 없으면 section_title 번호 fallback
        seq = slide_plan.get("subtitle_seq", 0)
        if seq == 0:
            if sec_title:
                parts = sec_title.split(".")
                if len(parts) >= 2 and parts[1].strip().split()[0].isdigit():
                    seq = int(parts[1].strip().split()[0])
            if seq == 0:
                seq = 1
        subtitle_labeled = f"{seq:02d} {subtitle}"
        _set_shape("9", _truncate_to_lines(subtitle_labeled, 8_000_000, 14, 1))

    # Zone 2: 사이드바 — 존 맵(body_title/body_desc) 우선, 없으면 레거시 dict
    z = _zone(template_file)
    label_id = z.get("body_title")
    desc_id  = z.get("body_desc")
    if label_id:
        label_text = sec_title or chapter_title
        # 슬라이드별 body_title truncate 파라미터 읽기 (하네스 우선)
        _sids_label = _load_slide_shape_ids().get(template_file.replace(".xml",""), {})
        _bt_max_lines = _sids_label.get("body_title_max_lines", 3)
        _bt_font_pt   = _sids_label.get("body_title_font_pt", 16)
        _sp_label = _find_shape_by_id(root, label_id)
        _bt_cx = 2_200_000
        if _sp_label is not None:
            _xfrm = _sp_label.find(f".//{{{ns_a}}}xfrm")
            if _xfrm is not None:
                _ext = _xfrm.find(f"{{{ns_a}}}ext")
                if _ext is not None:
                    _bt_cx = int(_ext.get("cx", _bt_cx))
        label_text = _truncate_to_lines(label_text, _bt_cx, _bt_font_pt, _bt_max_lines)
        # … 발생 시 숫자 prefix 제거 후 재시도
        if label_text.endswith("…"):
            bare_label = re.sub(r'^\d+(\.\d+)*[\s.]*', '', sec_title or chapter_title).strip()
            bare_trunc = _truncate_to_lines(bare_label, _bt_cx, _bt_font_pt, _bt_max_lines)
            if not bare_trunc.endswith("…"):
                label_text = bare_trunc
        # body_title은 Pretendard SemiBold 강제 (템플릿마다 typeface 상태가 달라도 통일)
        _set_shape(label_id, label_text, force_semibold=True)
        # 1) 텍스트박스 cy 동적 확장 (줄바꿈 발생 시)
        _auto_resize_textbox(root, label_id, label_text)
        # 2) 확장된 body_title 기준으로 body_desc 위치 재조정 (겹침 방지)
        _resize_sidebar_and_reposition_desc(root, label_id, desc_id, label_text)

    # group_title: 우측 패널 네이비 박스(사이드 소제목) → subtitle 주입
    # 폰트·줄수는 하네스(slide_shape_ids) 우선, 기본 14pt·2줄 (박스 높이 ~0.5") — body_title과 동일 패턴
    group_id = z.get("group_title")
    if group_id and subtitle:
        _sids_grp = _load_slide_shape_ids().get(template_file.replace(".xml", ""), {})
        _gt_font  = _sids_grp.get("group_title_font_pt", 14)
        _gt_lines = _sids_grp.get("group_title_max_lines", 2)
        _grp_cx = 3_800_000
        _sp_grp = _find_shape_by_id(root, group_id)
        if _sp_grp is not None:
            _xfrm = _sp_grp.find(f".//{{{ns_a}}}xfrm")
            if _xfrm is not None:
                _ext = _xfrm.find(f"{{{ns_a}}}ext")
                if _ext is not None:
                    _grp_cx = int(_ext.get("cx", _grp_cx))
        grp_text = _truncate_to_lines(subtitle, _grp_cx, _gt_font, _gt_lines)
        _set_shape(group_id, grp_text, force_semibold=True)

    # section_desc fallback: LLM 누락 시 section_title 기반 1줄 desc 생성
    if not sec_desc and sec_title:
        bare = re.sub(r'^\d+(\.\d+)*[\s.]*', '', sec_title).strip()
        sec_desc = f"{bare}의 개념과 주요 내용을 다룬다." if bare else ""
    if desc_id and sec_desc:
        # 실제 shape cx 읽기
        sp_desc = _find_shape_by_id(root, desc_id)
        desc_cx = 2_200_000
        if sp_desc is not None:
            xfrm = sp_desc.find(f".//{{{ns_a}}}xfrm")
            if xfrm is not None:
                ext = xfrm.find(f"{{{ns_a}}}ext")
                if ext is not None:
                    desc_cx = int(ext.get("cx", desc_cx))
        desc_txt = _truncate_to_lines(sec_desc, desc_cx, 12, 5)
        # … 말줄임 발생 시 자동 요약 (em dash 이후 설명구 제거)
        if desc_txt.endswith("…"):
            import re as _re
            short = _re.split(r'\s*[—.]\s*', sec_desc)[0].strip()
            desc_txt2 = _truncate_to_lines(short, desc_cx, 12, 5)
            if not desc_txt2.endswith("…"):
                desc_txt = desc_txt2
        _set_shape(desc_id, desc_txt)


def _edit_slide29(xml_path: Path, slide_plan: dict) -> None:
    """slide29 (연도별/월별 타임라인):
    Zone1: 헤더(ID=8,9) | Zone2: 사이드바(ID=18,19)
    Zone3: 연도별 레이블(16/17/22/23) + rich content 삽입
    periods: [{"label":"2026","period":"Jan~Jun","items":[...],
               "kpi":"...","team":"...","risk":"낮음"}]
    """
    content = slide_plan.get("content", {})
    periods = content.get("periods", [])
    body    = content.get("body", {})
    if isinstance(body, dict):
        periods = periods or body.get("periods", [])

    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return

    _apply_common_zones(root, slide_plan, "slide29.xml")
    import copy as _copy
    ns_p, ns_a = _NS_P, _NS_A

    def _set(sid, text):
        sp = _find_shape_by_id(root, sid)
        if sp is None: return
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None: return
        orig_rPr = None
        for r in sp.findall(f".//{{{ns_a}}}r"):
            rPr_e = r.find(f"{{{ns_a}}}rPr")
            if rPr_e is not None:
                orig_rPr = _copy.deepcopy(rPr_e)
                orig_rPr.set("lang", _ppt_lang()); orig_rPr.set("dirty","0"); break
        for p in txBody.findall(f"{{{ns_a}}}p"):
            for r in p.findall(f"{{{ns_a}}}r"): p.remove(r)
            end = p.find(f"{{{ns_a}}}endParaRPr")
            idx = list(p).index(end) if end is not None else len(p)
            r_new = ET.Element(f"{{{ns_a}}}r")
            if orig_rPr: r_new.append(_copy.deepcopy(orig_rPr))
            ET.SubElement(r_new, f"{{{ns_a}}}t").text = text
            p.insert(idx, r_new); break

    # 연도 레이블 설정 — zone_map에서 읽어 하드코딩 제거
    label_ids = _zone("slide29.xml").get("body", {}).get("period_labels", ["16","17","22","23"])
    for i, period in enumerate(periods[:4]):
        if isinstance(period, dict) and i < len(label_ids):
            _set(label_ids[i], period.get("label",""))

    # rich content: 대각선 레이아웃 — 각 연도 아래에 텍스트 박스 삽입
    # slide29는 대각선 계단식 구조 (2026 최상단, 2023 최하단)
    # 각 연도 콘텐츠는 해당 연도 바 아래 빈 영역에 삽입
    # 근사 위치: 각 열의 x 오프셋과 연도별 y 오프셋 계산
    _tc = _load_common_formatting().get("timeline_colors", {})
    DARK_COLS  = _tc.get("dark_cols",  ["3C41E6","2D30B8","1E2490","1419AB"])
    RISK_COLOR = _tc.get("risk_color", {"낮음":"52C41A","중간":"F5A623","높음":"FF4B4B"})
    _s29 = _load_slide_shape_ids().get("slide29", {})
    SLIDE_W   = _s29.get("slide_w",   12_192_000)
    SIDEBAR_W = _s29.get("sidebar_w", 2_700_000)
    BAR_H     = _s29.get("bar_h",     350_000)
    CONTENT_W = SLIDE_W - SIDEBAR_W
    COL_W     = CONTENT_W // 4
    _y_base   = _s29.get("year_y_base", [750_000, 900_000, 1_050_000, 1_200_000])
    YEAR_Y    = list(_y_base)  # 실제 템플릿 bar y 위치 (slide_shape_ids.json 실측값)
    YEAR_X    = [SIDEBAR_W, SIDEBAR_W + COL_W, SIDEBAR_W + COL_W*2, SIDEBAR_W + COL_W*3]
    YEAR_CX   = [CONTENT_W, CONTENT_W - COL_W, CONTENT_W - COL_W*2, CONTENT_W - COL_W*3]

    spTree = root.find(f".//{{{ns_p}}}spTree")

    def _mk_rPr(sz, bold=False, col=None):
        at = {"lang":"ko-KR","dirty":"0","sz":str(sz)}
        if bold: at["b"] = "1"
        r = ET.Element(f"{{{ns_a}}}rPr", **at)
        if col:
            sf = ET.SubElement(r, f"{{{ns_a}}}solidFill")
            ET.SubElement(sf, f"{{{ns_a}}}srgbClr", val=col)
        for t in ("latin","ea","cs"):
            ET.SubElement(r, f"{{{ns_a}}}{t}", typeface=_TEMPLATE_FONT)
        return r

    def _add_desc_box(sid, x, y, cx, cy, lines):
        sp = ET.SubElement(spTree, f"{{{ns_p}}}sp")
        nv = ET.SubElement(sp, f"{{{ns_p}}}nvSpPr")
        ET.SubElement(nv, f"{{{ns_p}}}cNvPr", id=str(sid), name=f"PD{sid}")
        cs = ET.SubElement(nv, f"{{{ns_p}}}cNvSpPr")
        ET.SubElement(cs, f"{{{ns_a}}}spLocks", noGrp="1")
        ET.SubElement(nv, f"{{{ns_p}}}nvPr")
        s2 = ET.SubElement(sp, f"{{{ns_p}}}spPr")
        xf = ET.SubElement(s2, f"{{{ns_a}}}xfrm")
        ET.SubElement(xf, f"{{{ns_a}}}off", x=str(x), y=str(y))
        ET.SubElement(xf, f"{{{ns_a}}}ext", cx=str(cx), cy=str(cy))
        pg = ET.SubElement(s2, f"{{{ns_a}}}prstGeom", prst="rect")
        ET.SubElement(pg, f"{{{ns_a}}}avLst")
        ET.SubElement(s2, f"{{{ns_a}}}noFill")
        tb = ET.SubElement(sp, f"{{{ns_p}}}txBody")
        ET.SubElement(tb, f"{{{ns_a}}}bodyPr", wrap="square",
                      lIns="152400", tIns="114300", rIns="152400", bIns="114300", anchor="t")
        ET.SubElement(tb, f"{{{ns_a}}}lstStyle")
        for txt, sz, bold, col in lines:
            p = ET.SubElement(tb, f"{{{ns_a}}}p")
            r2 = ET.SubElement(p, f"{{{ns_a}}}r")
            r2.append(_mk_rPr(sz, bold, col))
            ET.SubElement(r2, f"{{{ns_a}}}t").text = txt

    CONTENT_CY = _s29.get("content_box_cy", 2_065_000)
    MARGIN     = _s29.get("content_box_margin", 80_000)

    sid = 400
    for i, period in enumerate(periods[:4]):
        if not isinstance(period, dict): continue
        dk       = DARK_COLS[i]
        items    = period.get("items", [])
        kpi      = period.get("kpi", "")
        team     = period.get("team", "")
        risk     = period.get("risk", "낮음")
        p_period = period.get("period", "")

        # 파란 제목박스 좌표를 직접 읽어 content box 위치·크기 계산 (ML PPT 패턴)
        lbl_id = label_ids[i] if i < len(label_ids) else None
        sp_lbl = _find_shape_by_id(root, lbl_id) if lbl_id else None
        if sp_lbl is not None:
            spPr_l = sp_lbl.find(f"{{{ns_p}}}spPr")
            xf_l   = spPr_l.find(f"{{{ns_a}}}xfrm") if spPr_l is not None else None
            off_l  = xf_l.find(f"{{{ns_a}}}off") if xf_l is not None else None
            ext_l  = xf_l.find(f"{{{ns_a}}}ext") if xf_l is not None else None
            lbl_x  = int(off_l.get("x", YEAR_X[i])) if off_l is not None else YEAR_X[i]
            lbl_y  = int(off_l.get("y", YEAR_Y[i])) if off_l is not None else YEAR_Y[i]
            lbl_cx = int(ext_l.get("cx", COL_W))     if ext_l is not None else COL_W
            lbl_cy = int(ext_l.get("cy", BAR_H))     if ext_l is not None else BAR_H
            x  = lbl_x + MARGIN
            y  = lbl_y + lbl_cy + MARGIN
            cx = lbl_cx - MARGIN * 2
        else:
            x  = YEAR_X[i] + MARGIN
            y  = YEAR_Y[i] + BAR_H + MARGIN
            cx = COL_W - MARGIN * 2

        lines = [(f"▷ {p_period}", 1100, True, dk), ("", 500, False, None)]
        lines += [(item, 1100, False, "444444") for item in items[:3]]
        if team: lines += [("", 400, False, None), (f"● {team}", 1000, False, "888888")]
        if kpi:  lines += [("", 400, False, None), (f"✓ {kpi}", 1000, False, dk)]
        if risk:
            rc = RISK_COLOR.get(risk, "888888")
            lines += [(f"리스크: {risk}", 900, False, rc)]
        if items or kpi or team or p_period:
            _add_desc_box(sid, x, y, cx, CONTENT_CY, lines)
        sid += 1

    _clear_residual_placeholders(root)
    _write_xml(root, xml_path)


def _insert_rich_quarter_content(
    root, spTree, quarters: list,
    col_x_list: list, col_cx: int, content_top_y: int,
    upper_top_pad: int = 400_000,
    sep_lower_pad: int = 1_648_141,
    lower_start_gap: int = 261_149,
    box_dims: dict | None = None,
) -> None:
    """
    분기별(Q1~Q4) 슬라이드에 rich content 삽입 공통 헬퍼.
    slide31 / slide33 에서 호출.
    box_dims: slide_shape_ids.json 슬라이드 항목에서 전달 — 치수·행동 규칙 전체 포함.

    quarters 항목 형식 (dict):
      label   : "Q1" (이미 기본 shape에 설정됨)
      period  : "Jan~Mar" (기간)
      items   : ["□ 항목1", "□ 항목2", "□ 항목3"] (체크리스트)
      kpi     : "핵심 지표" (KPI)
      team    : "담당팀"
      risk    : "낮음/중간/높음"
      effort  : "Small/Medium/Large"
    """
    import copy as _copy, math as _math
    ns_p, ns_a = _NS_P, _NS_A

    _cf = _load_common_formatting()
    _tc = _cf.get("timeline_colors", {})
    _rq = _cf.get("rich_quarter_box_settings", {})
    DARK_COLS  = _tc.get("dark_cols",  ["3C41E6","2D30B8","1E2490","1419AB"])
    LIGHT_COLS = _tc.get("light_cols", ["E8ECFC","D4DCFB","C0CCFA","AABBF9"])
    RISK_COLOR = _tc.get("risk_color", {"낮음":"52C41A","중간":"F5A623","높음":"FF4B4B"})

    # 치수: slide_shape_ids box_dims 우선, 없으면 common_formatting 기본값 폴백
    _bd = box_dims or {}
    LINE_H           = _bd.get("line_h",    190_500)
    LOGO_Y           = _bd.get("logo_y",    6_200_000)
    PAD              = _bd.get("col_pad",   60_000)
    KPI_H            = _bd.get("kpi_h",     520_000)
    RISK_H           = _bd.get("risk_h",    380_000)
    GAP              = _bd.get("lower_gap", 120_000)
    SEP_H_val        = _bd.get("sep_h",     6_000)
    SUMM_GAP         = _bd.get("summ_gap",  30_000)
    SUMM_H           = _bd.get("summ_h",    120_000)
    # 행동 규칙: common_formatting 기본값, box_dims 오버라이드 가능
    NORM_AUTOFIT     = _bd.get("normAutofit",      _rq.get("normAutofit",      True))
    RISK_WIDTH_RATIO = _bd.get("risk_width_ratio", _rq.get("risk_width_ratio", 0.65))

    def _mk_rPr(sz, bold=False, col=None):
        at = {"lang":"ko-KR","dirty":"0","sz":str(sz)}
        if bold: at["b"] = "1"
        r = ET.Element(f"{{{ns_a}}}rPr", **at)
        if col:
            sf = ET.SubElement(r, f"{{{ns_a}}}solidFill")
            ET.SubElement(sf, f"{{{ns_a}}}srgbClr", val=col)
        for t in ("latin","ea","cs"):
            ET.SubElement(r, f"{{{ns_a}}}{t}", typeface=_TEMPLATE_FONT)
        return r

    def _add_box(sid, x, y, cx, cy, fill, rows=None, border=None):
        sp = ET.SubElement(spTree, f"{{{ns_p}}}sp")
        nv = ET.SubElement(sp, f"{{{ns_p}}}nvSpPr")
        ET.SubElement(nv, f"{{{ns_p}}}cNvPr", id=str(sid), name=f"RQ{sid}")
        cs = ET.SubElement(nv, f"{{{ns_p}}}cNvSpPr")
        ET.SubElement(cs, f"{{{ns_a}}}spLocks", noGrp="1")
        ET.SubElement(nv, f"{{{ns_p}}}nvPr")
        s2 = ET.SubElement(sp, f"{{{ns_p}}}spPr")
        xf = ET.SubElement(s2, f"{{{ns_a}}}xfrm")
        ET.SubElement(xf, f"{{{ns_a}}}off", x=str(x), y=str(y))
        ET.SubElement(xf, f"{{{ns_a}}}ext", cx=str(cx), cy=str(cy))
        pg = ET.SubElement(s2, f"{{{ns_a}}}prstGeom", prst="rect")
        ET.SubElement(pg, f"{{{ns_a}}}avLst")
        if fill == "none": ET.SubElement(s2, f"{{{ns_a}}}noFill")
        else:
            sf = ET.SubElement(s2, f"{{{ns_a}}}solidFill")
            ET.SubElement(sf, f"{{{ns_a}}}srgbClr", val=fill)
        if border:
            ln = ET.SubElement(s2, f"{{{ns_a}}}ln", w="19050")
            lsf = ET.SubElement(ln, f"{{{ns_a}}}solidFill")
            ET.SubElement(lsf, f"{{{ns_a}}}srgbClr", val=border)
        if rows:
            tb = ET.SubElement(sp, f"{{{ns_p}}}txBody")
            bpr = ET.SubElement(tb, f"{{{ns_a}}}bodyPr", wrap="square",
                          lIns="91440", tIns="45720", rIns="91440", bIns="45720", anchor="ctr")
            if NORM_AUTOFIT:
                ET.SubElement(bpr, f"{{{ns_a}}}normAutofit")
            ET.SubElement(tb, f"{{{ns_a}}}lstStyle")
            for txt, sz, bold, col, algn in rows:
                p = ET.SubElement(tb, f"{{{ns_a}}}p")
                ET.SubElement(p, f"{{{ns_a}}}pPr", algn=algn)
                r2 = ET.SubElement(p, f"{{{ns_a}}}r")
                r2.append(_mk_rPr(sz, bold, col))
                ET.SubElement(r2, f"{{{ns_a}}}t").text = txt

    def _add_desc(sid, x, y, cx, cy, lines):
        sp = ET.SubElement(spTree, f"{{{ns_p}}}sp")
        nv = ET.SubElement(sp, f"{{{ns_p}}}nvSpPr")
        ET.SubElement(nv, f"{{{ns_p}}}cNvPr", id=str(sid), name=f"RD{sid}")
        cs = ET.SubElement(nv, f"{{{ns_p}}}cNvSpPr")
        ET.SubElement(cs, f"{{{ns_a}}}spLocks", noGrp="1")
        ET.SubElement(nv, f"{{{ns_p}}}nvPr")
        s2 = ET.SubElement(sp, f"{{{ns_p}}}spPr")
        xf = ET.SubElement(s2, f"{{{ns_a}}}xfrm")
        ET.SubElement(xf, f"{{{ns_a}}}off", x=str(x), y=str(y))
        ET.SubElement(xf, f"{{{ns_a}}}ext", cx=str(cx), cy=str(cy))
        pg = ET.SubElement(s2, f"{{{ns_a}}}prstGeom", prst="rect")
        ET.SubElement(pg, f"{{{ns_a}}}avLst")
        ET.SubElement(s2, f"{{{ns_a}}}noFill")
        tb = ET.SubElement(sp, f"{{{ns_p}}}txBody")
        ET.SubElement(tb, f"{{{ns_a}}}bodyPr", wrap="square",
                      lIns="152400", tIns="114300", rIns="152400", bIns="114300", anchor="t")
        ET.SubElement(tb, f"{{{ns_a}}}lstStyle")
        for txt, sz, bold, col in lines:
            p = ET.SubElement(tb, f"{{{ns_a}}}p")
            r2 = ET.SubElement(p, f"{{{ns_a}}}r")
            r2.append(_mk_rPr(sz, bold, col))
            ET.SubElement(r2, f"{{{ns_a}}}t").text = txt

    # 콘텐츠 블록 높이 계산 — 실제 items 개수와 team 유무로 동적 계산
    _max_items = min(max(
        (len(q.get("items", [])) for q in quarters if isinstance(q, dict)),
        default=3
    ), 5)
    _has_team  = any(q.get("team") for q in quarters if isinstance(q, dict))
    UPPER_LINES = 1 + 1 + _max_items + (2 if _has_team else 0)  # period+blank+items[+blank+team]
    UPPER_CONT_H = UPPER_LINES * LINE_H
    DESC_START_Y = content_top_y + upper_top_pad

    SEP_Y    = DESC_START_Y + UPPER_CONT_H + sep_lower_pad
    SUMM_Y   = SEP_Y + SEP_H_val + SUMM_GAP
    SUMM_BOT = SUMM_Y + SUMM_H + SUMM_GAP
    LOWER_START = SUMM_BOT + lower_start_gap

    # col_cx는 int(균등) 또는 list(컬럼별 개별 폭) 모두 허용
    _col_cx = col_cx if isinstance(col_cx, list) else [col_cx] * len(col_x_list)

    # ── KPI·Risk 높이 동적 계산 (4열 통일) ─────────────────────────
    # 텍스트 줄 수 추정: sz는 100분의 1 포인트 단위(1100=11pt), 1pt=12700 EMU
    _V_PAD = 91_440  # bodyPr tIns(45720) + bIns(45720)

    def _box_h(text: str, bw: int, sz: int = 1100, header_lines: int = 1) -> int:
        """배경 도형 최소 높이 = (헤더+내용 줄 수) × LINE_H + 위아래 여백"""
        char_w = sz * 127        # sz 단위 → EMU (1pt=12700, sz=100분의1pt → ×127)
        inner_w = max(1, bw - 120_000)  # lIns(91440)+rIns(91440) 여백 제외
        cpp = max(1, inner_w // char_w)
        content_lines = _math.ceil(len(text) / cpp) if text else 1
        return (header_lines + content_lines) * LINE_H + _V_PAD

    # ── 하단 메타: 데이터 주도 ───────────────────────────────────────
    # 라벨/키는 하네스(quarter_meta_slots) 정의. plan이 q["meta"] dict를 주면
    # 그 키=라벨로 우선 사용(고정 kpi/risk/effort 강제 아님). 최대 3슬롯
    # (entries[0]=상단 wide, [1]=좌, [2]=우).
    _meta_slots = _bd.get("quarter_meta_slots", [
        {"label": "핵심 목표", "key": "kpi"},
        {"label": "리스크",    "key": "risk"},
        {"label": "규모",      "key": "effort"},
    ])

    def _q_meta(q):
        qm = q.get("meta") if isinstance(q, dict) else None
        if isinstance(qm, dict) and qm:
            return [(str(k), str(v)) for k, v in qm.items()][:3]
        if isinstance(q, dict):
            return [(s.get("label", ""), str(q.get(s.get("key", ""), ""))) for s in _meta_slots][:3]
        return [(s.get("label", ""), "") for s in _meta_slots][:3]

    _meta_by_q = [_q_meta(q) for q in quarters[:4]]

    max_kpi_h  = KPI_H
    max_risk_h = RISK_H
    for i_q in range(min(len(_meta_by_q), len(_col_cx))):
        entries   = _meta_by_q[i_q]
        bw_i      = _col_cx[i_q] - PAD * 2
        risk_w_i  = int(bw_i * RISK_WIDTH_RATIO)
        eff_w_i   = bw_i - risk_w_i - 40_000
        if len(entries) >= 1:
            max_kpi_h  = max(max_kpi_h,  _box_h(entries[0][1], bw_i))
        if len(entries) >= 2:
            max_risk_h = max(max_risk_h, _box_h(entries[1][1], risk_w_i))
        if len(entries) >= 3:
            max_risk_h = max(max_risk_h, _box_h(entries[2][1], eff_w_i))
    KPI_H  = max_kpi_h
    RISK_H = max_risk_h
    # ────────────────────────────────────────────────────────────────

    all_x = col_x_list[0] + PAD
    all_w = col_x_list[-1] + _col_cx[-1] - PAD - all_x

    sid = 300
    _add_box(sid, all_x, SEP_Y, all_w, SEP_H_val, "CCCCCC"); sid += 1
    _add_box(sid, all_x, SUMM_Y, all_w, SUMM_H, "none",
             [("분기별 실행 계획", 1000, False, "999999", "ctr")]); sid += 1

    for i, q in enumerate(quarters[:4]):
        if i >= len(col_x_list): break
        col_x = col_x_list[i]
        dk = DARK_COLS[i]; lt = LIGHT_COLS[i]
        bx = col_x + PAD; bw = _col_cx[i] - PAD * 2

        if isinstance(q, dict):
            period  = q.get("period", f"Q{i+1}")
            items   = q.get("items", [])
            team    = q.get("team", "")
        else:
            period = str(q); items = []; team = ""

        desc_lines = [(f"▷ {period}", 1100, True, dk), ("", 500, False, None)]
        desc_lines += [(item, 1100, False, "444444") for item in items[:3]]
        if team: desc_lines += [("", 500, False, None), (f"● {team}", 1000, False, "888888")]
        _add_desc(sid, bx, DESC_START_Y, bw, 4_000_000, desc_lines); sid += 1

        # 하단 메타 박스 — entries[0]=상단 wide, [1]=좌, [2]=우 (라벨·값 모두 데이터 주도)
        entries = _meta_by_q[i] if i < len(_meta_by_q) else []
        y = LOWER_START
        if len(entries) >= 1 and (entries[0][0] or entries[0][1]):
            _add_box(sid, bx, y, bw, KPI_H, "F4F6FE",
                     [(entries[0][0], 900, True, dk, "ctr"), (entries[0][1], 1100, False, "111111", "ctr")],
                     border=dk); sid += 1
            y += KPI_H + GAP
        # 좌(entries[1])·우(entries[2]) 박스는 서로 독립 (F2). 메타 2개면 좌측 박스를 full-width로
        # 채워 우측 빈 공간을 없앤다 — 3개일 때만 좌(65%)/우(35%) 분할 (빈 슬롯 방지)
        if len(entries) >= 2:
            _has_third = len(entries) >= 3 and (entries[2][0] or entries[2][1])
            risk_w = int(bw * RISK_WIDTH_RATIO) if _has_third else bw
            eff_w  = bw - int(bw * RISK_WIDTH_RATIO) - 40_000
            if entries[1][0] or entries[1][1]:
                rc = RISK_COLOR.get(entries[1][1], "888888")
                _add_box(sid, bx, y, risk_w, RISK_H, lt,
                         [(entries[1][0], 900, True, dk, "ctr"), (entries[1][1], 1100, False, rc, "ctr")]); sid += 1
            if _has_third:
                _add_box(sid, bx + int(bw * RISK_WIDTH_RATIO) + 40_000, y, eff_w, RISK_H, lt,
                         [(entries[2][0], 900, True, dk, "ctr"), (entries[2][1], 1100, False, "333333", "ctr")]); sid += 1


def _edit_slide31(xml_path: Path, slide_plan: dict) -> None:
    """slide31 (분기별 Q1→Q4): Zone1+2 공통 + Zone3 Q레이블 + rich content."""
    content  = slide_plan.get("content", {})
    quarters = content.get("quarters", content.get("body", {}).get("quarters", []) if isinstance(content.get("body"), dict) else [])
    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, "slide31.xml")
    import copy as _copy; ns_p, ns_a = _NS_P, _NS_A
    def _set(sid, text):
        sp = _find_shape_by_id(root, sid)
        if sp is None: return
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None: return
        orig_rPr = next((_copy.deepcopy(r.find(f"{{{ns_a}}}rPr")) for r in sp.findall(f".//{{{ns_a}}}r") if r.find(f"{{{ns_a}}}rPr") is not None), None)
        if orig_rPr is not None: orig_rPr.set("lang", _ppt_lang()); orig_rPr.set("dirty","0")
        for p in txBody.findall(f"{{{ns_a}}}p"):
            for r in p.findall(f"{{{ns_a}}}r"): p.remove(r)
            end = p.find(f"{{{ns_a}}}endParaRPr"); idx = list(p).index(end) if end is not None else len(p)
            r_new = ET.Element(f"{{{ns_a}}}r")
            if orig_rPr: r_new.append(_copy.deepcopy(orig_rPr))
            ET.SubElement(r_new, f"{{{ns_a}}}t").text = text; p.insert(idx, r_new); break
    _s31 = _load_slide_shape_ids().get("slide31", {})
    _qlabel_ids31 = _s31.get("q_label_ids", ["23","24","27","28"])
    for i, (lid, q) in enumerate(zip(_qlabel_ids31, quarters)):
        _set(lid, q.get("label", f"Q{i+1}") if isinstance(q, dict) else str(q))
        # Q-레이블 1줄 강제: bodyPr wrap="none"
        _sp_ql = _find_shape_by_id(root, lid)
        if _sp_ql is not None:
            _tb_ql = _sp_ql.find(f"{{{ns_p}}}txBody")
            if _tb_ql is not None:
                _bpr_ql = _tb_ql.find(f"{{{ns_a}}}bodyPr")
                if _bpr_ql is not None:
                    _bpr_ql.set("wrap", "none")
    # rich content 삽입 — items가 있을 때만
    has_items = any(q.get("items") for q in quarters if isinstance(q, dict))
    if quarters and has_items:
        spTree = root.find(f".//{{{ns_p}}}spTree")
        col_x_list    = _s31.get("col_x_list",   [2_657_592, 4_897_938, 7_393_720, 9_696_217])
        col_cx        = _s31.get("col_cx_list",   [2_240_346, 2_495_782, 2_302_497, 2_495_782])
        content_top_y = _s31.get("content_top_y", 1_099_874)
        _insert_rich_quarter_content(
            root, spTree, quarters, col_x_list, col_cx, content_top_y,
            upper_top_pad=_s31.get("upper_top_pad", 400_000),
            sep_lower_pad=_s31.get("sep_lower_pad", 1_648_141),
            lower_start_gap=_s31.get("lower_start_gap", 261_149),
            box_dims=_s31,
        )
    _clear_residual_placeholders(root); _write_xml(root, xml_path)


def _edit_slide33(xml_path: Path, slide_plan: dict) -> None:
    """slide33 (분기별 변형): Zone1+2 공통 + Zone3 Q레이블 + rich content.
    dynamic_layout 하네스 설정이 있으면 body_title 높이 기반으로 bar/Q레이블을 자동 재배치."""
    content  = slide_plan.get("content", {})
    quarters = content.get("quarters", content.get("body", {}).get("quarters", []) if isinstance(content.get("body"), dict) else [])
    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, "slide33.xml")
    import copy as _copy; ns_p, ns_a = _NS_P, _NS_A
    def _set(sid, text):
        sp = _find_shape_by_id(root, sid)
        if sp is None: return
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None: return
        orig_rPr = next((_copy.deepcopy(r.find(f"{{{ns_a}}}rPr")) for r in sp.findall(f".//{{{ns_a}}}r") if r.find(f"{{{ns_a}}}rPr") is not None), None)
        if orig_rPr is not None: orig_rPr.set("lang", _ppt_lang()); orig_rPr.set("dirty","0")
        for p in txBody.findall(f"{{{ns_a}}}p"):
            for r in p.findall(f"{{{ns_a}}}r"): p.remove(r)
            end = p.find(f"{{{ns_a}}}endParaRPr"); idx = list(p).index(end) if end is not None else len(p)
            r_new = ET.Element(f"{{{ns_a}}}r")
            if orig_rPr: r_new.append(_copy.deepcopy(orig_rPr))
            ET.SubElement(r_new, f"{{{ns_a}}}t").text = text; p.insert(idx, r_new); break
    _s33 = _load_slide_shape_ids().get("slide33", {})
    _qlabel_ids33 = _s33.get("q_label_ids", ["35","36","37","38"])
    for i, (lid, q) in enumerate(zip(_qlabel_ids33, quarters)):
        _set(lid, q.get("label", f"Q{i+1}") if isinstance(q, dict) else str(q))
        # Q-레이블 1줄 강제: bodyPr wrap="none"
        _sp_ql33 = _find_shape_by_id(root, lid)
        if _sp_ql33 is not None:
            _tb_ql33 = _sp_ql33.find(f"{{{ns_p}}}txBody")
            if _tb_ql33 is not None:
                _bpr_ql33 = _tb_ql33.find(f"{{{ns_a}}}bodyPr")
                if _bpr_ql33 is not None:
                    _bpr_ql33.set("wrap", "none")

    # dynamic_layout: body_title 높이 확장 시 bar·Q레이블 등 delta_y로 이동, content_top_y 자동 조정
    content_top_y = _s33.get("content_top_y", 1_850_000)
    dyn = _s33.get("dynamic_layout", {})
    if dyn:
        bt_sp = _find_shape_by_id(root, dyn.get("body_title_id", "25"))
        bar_sp = _find_shape_by_id(root, dyn.get("bar_id", "31"))
        bt_b = _resolve_shape_bounds(bt_sp)
        bar_b = _resolve_shape_bounds(bar_sp)
        if bt_b and bar_b:
            new_bt_bottom = bt_b['y'] + bt_b['cy']
            new_bar_y = new_bt_bottom + dyn.get("gap_below_title", 69_633)
            delta_y = new_bar_y - bar_b['y']
            if delta_y != 0:
                # bar 이동
                off_bar, _ = _get_shape_xfrm_elems(bar_sp)
                if off_bar: off_bar.set('y', str(new_bar_y))
                # shift_shape_ids 이동 (Q레이블·삼각형·배경rect 모두 동일 delta_y)
                for sid in dyn.get("shift_shape_ids", []):
                    sp = _find_shape_by_id(root, sid)
                    if sp is None: continue
                    off, _ = _get_shape_xfrm_elems(sp)
                    if off: off.set('y', str(int(off.get('y', 0)) + delta_y))
                content_top_y += delta_y

    # rich content 삽입 — items가 있을 때만
    has_items33 = any(q.get("items") for q in quarters if isinstance(q, dict))
    if quarters and has_items33:
        spTree = root.find(f".//{{{ns_p}}}spTree")
        col_x_list = _s33.get("col_x_list",  [0, 3_077_029, 6_084_461, 9_184_568])
        col_cx     = _s33.get("col_cx_list", [3_077_029, 3_007_432, 3_100_107, 3_007_432])
        _insert_rich_quarter_content(
            root, spTree, quarters, col_x_list, col_cx, content_top_y,
            upper_top_pad=_s33.get("upper_top_pad", 400_000),
            sep_lower_pad=_s33.get("sep_lower_pad", 1_648_141),
            lower_start_gap=_s33.get("lower_start_gap", 261_149),
            box_dims=_s33,
        )
    _clear_residual_placeholders(root); _write_xml(root, xml_path)


def _apply_spAutoFit(root, sid: str) -> None:
    """shape의 bodyPr에 spAutoFit을 설정해 텍스트에 맞게 높이가 자동 확장되도록 한다."""
    ns_p, ns_a = _NS_P, _NS_A
    sp = _find_shape_by_id(root, sid)
    if sp is None: return
    txBody = sp.find(f"{{{ns_p}}}txBody")
    if txBody is None: return
    bodyPr = txBody.find(f"{{{ns_a}}}bodyPr")
    if bodyPr is None: return
    for child in list(bodyPr):
        if child.tag.split("}")[-1] in ("noAutofit", "spAutoFit", "normAutofit"):
            bodyPr.remove(child)
    ET.SubElement(bodyPr, f"{{{ns_a}}}spAutoFit")


def _equalize_row_heights(root, id_groups: list[list[str]]) -> None:
    """같은 행(컬럼)에 속하는 shape 그룹들의 cy를 그룹 내 최댓값으로 통일한다.
    id_groups: [[col0_shape_a, col1_shape_a, ...], [col0_shape_b, col1_shape_b, ...]]
    → 같은 인덱스(열)끼리 묶어 최대 cy를 구한 뒤 모두 같은 값으로 맞춘다.
    실제로는 슬롯 타입별로 같은 열 인덱스를 묶어야 하므로,
    id_groups는 [[row_type_a_ids...], [row_type_b_ids...]] 형태로 전달된다.
    같은 index 위치의 shape들을 한 열로 취급한다."""
    ns_p, ns_a = _NS_P, _NS_A
    if not id_groups: return
    n_cols = max(len(g) for g in id_groups)
    for col_i in range(n_cols):
        # 이 열에 속하는 모든 shape 수집
        col_ids = [g[col_i] for g in id_groups if col_i < len(g)]
        max_cy = 0
        shape_exts = []
        for sid in col_ids:
            sp = _find_shape_by_id(root, sid)
            if sp is None: continue
            spPr = sp.find(f"{{{ns_p}}}spPr")
            if spPr is None: continue
            xfrm = spPr.find(f"{{{ns_a}}}xfrm")
            if xfrm is None: continue
            ext = xfrm.find(f"{{{ns_a}}}ext")
            if ext is None: continue
            try:
                cy = int(ext.get("cy", "0"))
                max_cy = max(max_cy, cy)
                shape_exts.append(ext)
            except ValueError:
                pass
        if max_cy > 0:
            for ext in shape_exts:
                ext.set("cy", str(max_cy))


def _apply_slots(root, content: dict, slots: list, slide_plan: dict | None = None) -> None:
    """JSON slots 선언 기반 범용 슬롯 채우기 엔진.

    slot types:
      item_title / item_body : content_key 리스트 → 각 id에 순서대로 기록
      img_slot               : content_key 리스트 → "[이미지: {description}]" 포맷으로 기록
      sub_heading            : | {major}.{minor}.{sub} prefix 자동 부착
      clear                  : ids를 빈 문자열로 초기화

    item_body / img_slot 슬롯에는 spAutoFit 적용 후 같은 열 기준 최대 cy로 높이 통일.
    """
    import copy as _copy
    ns_p, ns_a = _NS_P, _NS_A
    img_fmt = _load_common_formatting().get("image_slot_format", "[이미지: {description}]")

    # sub_heading prefix 계산 (slide_plan에서 챕터/섹션 추론)
    _chapter  = 1
    _section  = 1
    _subsect  = 1
    if slide_plan:
        sec_title = slide_plan.get("content", {}).get("section_title", "")
        m = re.match(r"(\d+)\.?(\d*)", sec_title.strip())
        if m:
            _chapter = int(m.group(1)) if m.group(1) else 1
            _section = int(m.group(2)) if m.group(2) else 1

    def _set(sid, text):
        sp = _find_shape_by_id(root, sid)
        if sp is None: return
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None: return
        orig_rPr = next(
            (_copy.deepcopy(r.find(f"{{{ns_a}}}rPr"))
             for r in sp.findall(f".//{{{ns_a}}}r")
             if r.find(f"{{{ns_a}}}rPr") is not None),
            None,
        )
        if orig_rPr is not None:
            orig_rPr.set("lang", _ppt_lang()); orig_rPr.set("dirty", "0")
        for p in txBody.findall(f"{{{ns_a}}}p"):
            for r in p.findall(f"{{{ns_a}}}r"): p.remove(r)
            end = p.find(f"{{{ns_a}}}endParaRPr")
            idx = list(p).index(end) if end is not None else len(p)
            r_new = ET.Element(f"{{{ns_a}}}r")
            if orig_rPr: r_new.append(_copy.deepcopy(orig_rPr))
            ET.SubElement(r_new, f"{{{ns_a}}}t").text = text
            p.insert(idx, r_new); break

    # item_body / img_slot id 그룹 수집 (열 높이 통일용)
    body_id_groups: list[list[str]] = []

    for slot in slots:
        stype      = slot.get("type", "")
        ids        = slot.get("ids", [])
        width_emu  = slot.get("width_emu", 2000000)
        font_pt    = slot.get("font_pt", 12)
        max_lines  = slot.get("max_lines", 4)

        if stype == "clear":
            for sid in ids:
                _set(sid, "")
        elif stype == "img_slot":
            values    = content.get(slot.get("content_key", "image_descriptions"), [])
            slot_fmt  = slot.get("format", img_fmt)
            for i, sid in enumerate(ids):
                val  = values[i] if i < len(values) else ""
                if val and str(val).strip().startswith("[이미지:"):
                    text = val
                elif val and str(val).strip():
                    text = slot_fmt.replace("{description}", str(val))
                else:
                    # image_descriptions 누락 → 파생 캡션(빈 이미지영역 구조적 방지, 생성 로직)
                    text = _image_caption(content, i)
                # 이미지 캡션은 '…' 없이 깔끔히 끊는다(상세 설명 전체는 plan에 보존 — 이미지생성용)
                truncated = _truncate_to_lines(text, width_emu, font_pt, max_lines, ellipsis=False)
                _set_image_slot_text(root, sid, truncated)
                # img_slot은 높이 고정 — spAutoFit 미적용
            body_id_groups.append(list(ids))
        elif stype == "insight":
            values = content.get(slot.get("content_key", "insights"), [])
            for i, sid in enumerate(ids):
                val = values[i] if i < len(values) else ""
                _set(sid, _truncate_to_lines(val, width_emu, font_pt, max_lines))
        elif stype == "sub_heading":
            val = content.get(slot.get("content_key", "sub_heading")) or _make_sub_heading(content)
            # _apply_common_zones가 이미 prefix 붙인 경우 중복 방지
            if val and not val.startswith("|"):
                prefix = f"| {_chapter}.{_section}.{_subsect}"
                full = f"{prefix} {val}"
            else:
                full = val or ""
            _set(ids[0] if ids else "", _truncate_to_lines(full, width_emu, font_pt, max_lines))
        elif stype == "item_body":
            values = content.get(slot.get("content_key", "descriptions"), [])
            for i, sid in enumerate(ids):
                val = values[i] if i < len(values) else ""
                _set(sid, _truncate_to_lines(val, width_emu, font_pt, max_lines))
                _apply_spAutoFit(root, sid)  # 텍스트 맞게 높이 자동 확장
            body_id_groups.append(list(ids))
        else:  # item_title, …
            values = content.get(slot.get("content_key", "items"), [])
            for i, sid in enumerate(ids):
                val = values[i] if i < len(values) else ""
                _set(sid, _truncate_to_lines(val, width_emu, font_pt, max_lines))

    # item_body와 img_slot이 같은 열에 있을 때 최대 cy 기준으로 높이 통일
    if len(body_id_groups) >= 2:
        _equalize_row_heights(root, body_id_groups)


def _edit_slide13(xml_path: Path, slide_plan: dict) -> None:
    """slide13 (3열 아이콘카드): Zone1+2 공통 + JSON slots 기반 Zone3 채우기."""
    tmpl    = slide_plan.get("template_file", "slide13.xml")
    content = slide_plan.get("content", {})
    body    = content.get("body", {})
    if isinstance(body, dict):
        for k, v in body.items():
            content.setdefault(k, v)
    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, tmpl)
    import copy as _copy; ns_p, ns_a = _NS_P, _NS_A
    def _set_local(sid, text):
        sp = _find_shape_by_id(root, sid)
        if sp is None: return
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None: return
        orig_rPr = next((_copy.deepcopy(r.find(f"{{{ns_a}}}rPr")) for r in sp.findall(f".//{{{ns_a}}}r") if r.find(f"{{{ns_a}}}rPr") is not None), None)
        if orig_rPr is not None: orig_rPr.set("lang", _ppt_lang()); orig_rPr.set("dirty", "0")
        for p in txBody.findall(f"{{{ns_a}}}p"):
            for r in p.findall(f"{{{ns_a}}}r"): p.remove(r)
            end = p.find(f"{{{ns_a}}}endParaRPr"); idx = list(p).index(end) if end is not None else len(p)
            r_new = ET.Element(f"{{{ns_a}}}r")
            if orig_rPr: r_new.append(_copy.deepcopy(orig_rPr))
            ET.SubElement(r_new, f"{{{ns_a}}}t").text = text; p.insert(idx, r_new); break
    z = _zone(tmpl)
    for sid in (z.get("body", {}).get("sub_heading") or []):
        sub_text = _make_sub_heading(content)
        _set_local(sid, _truncate_to_lines(sub_text, 9_000_000, 16, 2) if sub_text else "")
    _s13  = _load_slide_shape_ids().get("slide13", {})
    slots = _s13.get("slots", [])
    if slots:
        _apply_slots(root, content, slots)
    else:
        # 폴백: slots 키 없을 때 구버전 동작
        item_title_ids = _s13.get("item_title_ids", ["18","19","21"])
        item_body_ids  = _s13.get("item_body_ids",  ["62","38","69"])
        clear_ids      = _s13.get("clear_ids",      ["27","29","30"])
        items = (content.get("items") or content.get("bullets") or [])[:3]
        descs = content.get("descriptions", [])
        ttl   = _s13.get("item_title", {"width_emu":2000000,"font_pt":14,"max_lines":2})
        bdy   = _s13.get("item_body",  {"width_emu":3000000,"font_pt":12,"max_lines":6})
        for i, (t_id, b_id) in enumerate(zip(item_title_ids, item_body_ids)):
            _set_local(t_id, _truncate_to_lines(items[i] if i < len(items) else "", ttl["width_emu"], ttl["font_pt"], ttl["max_lines"]))
            _set_local(b_id, _truncate_to_lines(descs[i] if i < len(descs) else "", bdy["width_emu"], bdy["font_pt"], bdy["max_lines"]))
        for sid in clear_ids:
            _set_local(sid, "")
    _clear_residual_placeholders(root); _write_xml(root, xml_path)


def _edit_slide15_v2(xml_path: Path, slide_plan: dict) -> None:
    """slide15 (3열 아이콘카드): slide13과 동일 구조, _apply_slots 기반.
    기존 edit_slide15()는 content/slide_plan 파라미터 혼용 버그로 교체."""
    tmpl    = "slide15.xml"
    content = slide_plan.get("content", {})
    body    = content.get("body", {})
    if isinstance(body, dict):
        for k, v in body.items():
            content.setdefault(k, v)
    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, tmpl)
    import copy as _copy; ns_p, ns_a = _NS_P, _NS_A
    def _set_local(sid, text):
        sp = _find_shape_by_id(root, sid)
        if sp is None: return
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None: return
        orig_rPr = next((_copy.deepcopy(r.find(f"{{{ns_a}}}rPr")) for r in sp.findall(f".//{{{ns_a}}}r") if r.find(f"{{{ns_a}}}rPr") is not None), None)
        if orig_rPr is not None: orig_rPr.set("lang", _ppt_lang()); orig_rPr.set("dirty", "0")
        for p in txBody.findall(f"{{{ns_a}}}p"):
            for r in p.findall(f"{{{ns_a}}}r"): p.remove(r)
            end = p.find(f"{{{ns_a}}}endParaRPr"); idx = list(p).index(end) if end is not None else len(p)
            r_new = ET.Element(f"{{{ns_a}}}r")
            if orig_rPr: r_new.append(_copy.deepcopy(orig_rPr))
            ET.SubElement(r_new, f"{{{ns_a}}}t").text = text; p.insert(idx, r_new); break
    z = _zone(tmpl)
    for sid in (z.get("body", {}).get("sub_heading") or []):
        # ID=41 cx=9289092(거의 전체 폭) — 전체 sub_heading 텍스트 사용
        sub_text = content.get("sub_heading") or _make_sub_heading(content)
        _set_local(sid, sub_text)
        # Vision Fix가 sub_heading 박스에 남긴 여분 단락 제거
        sp = _find_shape_by_id(root, sid)
        if sp is not None:
            txb = sp.find(f"{{{ns_p}}}txBody")
            if txb is not None:
                for xp in txb.findall(f"{{{ns_a}}}p")[1:]:
                    txb.remove(xp)
    _s15  = _load_slide_shape_ids().get("slide15", {})
    slots = _s15.get("slots", [])
    if slots:
        _apply_slots(root, content, slots)
    _clear_residual_placeholders(root); _write_xml(root, xml_path)


def _edit_slide21(xml_path: Path, slide_plan: dict) -> None:
    """slide21 (이미지+2블록): Zone1+2 공통 + Zone3 bullet+이미지안내."""
    content = slide_plan.get("content", {})
    body    = content.get("body", {})
    bullets = (content.get("bullets") or content.get("items")
               or (body.get("bullets") if isinstance(body, dict) else None) or [])
    body_text = content.get("section_desc", "")
    title     = slide_plan.get("title", "")
    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, "slide21.xml")
    import copy as _copy; ns_p, ns_a = _NS_P, _NS_A
    def _set(sid, text):
        sp = _find_shape_by_id(root, sid)
        if sp is None: return
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None: return
        orig_rPr = next((_copy.deepcopy(r.find(f"{{{ns_a}}}rPr")) for r in sp.findall(f".//{{{ns_a}}}r") if r.find(f"{{{ns_a}}}rPr") is not None), None)
        if orig_rPr is not None: orig_rPr.set("lang", _ppt_lang()); orig_rPr.set("dirty","0")
        for p in txBody.findall(f"{{{ns_a}}}p"):
            for r in p.findall(f"{{{ns_a}}}r"): p.remove(r)
            end = p.find(f"{{{ns_a}}}endParaRPr"); idx = list(p).index(end) if end is not None else len(p)
            r_new = ET.Element(f"{{{ns_a}}}r")
            if orig_rPr: r_new.append(_copy.deepcopy(orig_rPr))
            ET.SubElement(r_new, f"{{{ns_a}}}t").text = text; p.insert(idx, r_new); break
    _s21 = _load_slide_shape_ids().get("slide21", {})
    _b1_id = _s21.get("body1_id", "12")
    _b2_id = _s21.get("body2_id", "14")
    _b3_id = _s21.get("body3_id", "15")
    _im_id = _s21.get("img_id",   "18")
    _b1_21 = _s21.get("body1", {"width_emu":3000000,"font_pt":14,"max_lines":2})
    _b2_21 = _s21.get("body2", {"width_emu":2500000,"font_pt":12,"max_lines":3})
    _b3_21 = _s21.get("body3", {"width_emu":2500000,"font_pt":12,"max_lines":3})
    _set(_b1_id, _truncate_to_lines(bullets[0] if bullets else body_text, _b1_21["width_emu"], _b1_21["font_pt"], _b1_21["max_lines"]))
    _set(_b2_id, _truncate_to_lines(bullets[1] if len(bullets) > 1 else body_text, _b2_21["width_emu"], _b2_21["font_pt"], _b2_21["max_lines"]))
    _set(_b3_id, _truncate_to_lines(bullets[2] if len(bullets) > 2 else "", _b3_21["width_emu"], _b3_21["font_pt"], _b3_21["max_lines"]))
    # 이미지 슬롯: 항상 "[이미지: …]" 가이드 주입 (plan 누락 시 section_title/items에서 파생) — 빈 이미지영역 구조적 방지
    _set_image_slot_text(root, _im_id, _image_caption(content, 0))
    _clear_residual_placeholders(root); _write_xml(root, xml_path)


def _resolve_layout_shape_bounds(xml_path: Path, sid: str) -> dict | None:
    """
    레이아웃 상속 shape(cx=0)의 실제 좌표를 슬라이드 레이아웃 XML에서 읽어온다.
    반환: {'x':int,'y':int,'cx':int,'cy':int} 또는 None
    """
    try:
        import xml.etree.ElementTree as _ET2
        ns_p = _NS_P; ns_a = _NS_A
        # _rels에서 레이아웃 파일 경로 찾기
        rels_path = xml_path.parent / '_rels' / (xml_path.name + '.rels')
        if not rels_path.exists(): return None
        rels_root = _ET2.parse(rels_path).getroot()
        rns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        layout_target = next((r.get('Target','') for r in rels_root
                              if 'slideLayout' in r.get('Target','')), None)
        if not layout_target: return None
        layout_path = (xml_path.parent.parent / layout_target.lstrip('../')).resolve()
        if not layout_path.exists(): return None
        layout_root = _ET2.parse(layout_path).getroot()
        for sp in layout_root.iter(f'{{{ns_p}}}sp'):
            cpr = sp.find(f'{{{ns_p}}}nvSpPr/{{{ns_p}}}cNvPr')
            if cpr is None or cpr.get('id') != str(sid): continue
            xfrm = sp.find(f'.//{{{ns_a}}}xfrm')
            if xfrm is None: continue
            off = xfrm.find(f'{{{ns_a}}}off')
            ext = xfrm.find(f'{{{ns_a}}}ext')
            if off is None or ext is None: continue
            return {'x': int(off.get('x',0)), 'y': int(off.get('y',0)),
                    'cx': int(ext.get('cx',0)), 'cy': int(ext.get('cy',0))}
    except Exception:
        pass
    return None


def _materialize_layout_shape(root, sid: str, xml_path: Path) -> bool:
    """
    레이아웃 상속 shape에 레이아웃의 실제 좌표를 명시적으로 설정.
    이후 _auto_resize_textbox 등이 동작 가능해짐.
    반환: 성공 여부
    """
    bounds = _resolve_layout_shape_bounds(xml_path, sid)
    if not bounds or bounds['cx'] == 0: return False
    sp = _find_shape_by_id(root, sid)
    if sp is None: return False
    spPr = sp.find(f'{{{_NS_P}}}spPr')
    if spPr is None: return False
    xfrm = spPr.find(f'{{{_NS_A}}}xfrm')
    if xfrm is None:
        xfrm = ET.SubElement(spPr, f'{{{_NS_A}}}xfrm')
        ET.SubElement(xfrm, f'{{{_NS_A}}}off', x=str(bounds['x']), y=str(bounds['y']))
        ET.SubElement(xfrm, f'{{{_NS_A}}}ext', cx=str(bounds['cx']), cy=str(bounds['cy']))
    else:
        off = xfrm.find(f'{{{_NS_A}}}off')
        ext = xfrm.find(f'{{{_NS_A}}}ext')
        if off is not None:
            off.set('x', str(bounds['x'])); off.set('y', str(bounds['y']))
        if ext is not None:
            ext.set('cx', str(bounds['cx'])); ext.set('cy', str(bounds['cy']))
    return True


_BODY_TITLE_FONT_PT = 20.0   # 사이드바 body_title 실측 폰트 크기 (Pretendard SemiBold 20pt)


def _auto_resize_textbox(root, sid: str, text: str, default_font_pt: float = _BODY_TITLE_FONT_PT,
                         xml_path: Path | None = None) -> None:
    """
    텍스트가 텍스트박스를 초과해 자동 줄바꿈될 경우 cy를 동적으로 확장.
    - font_pt: shape의 rPr.sz에서 읽음. sz 미설정(레이아웃 상속)이면 _BODY_TITLE_FONT_PT 사용
    - box_capacity: cx_pt / (font_pt * 0.6) — 한국어 CJK 실측 보정값
    - line_height_emu: _SIDEBAR_LINE_HEIGHT_EMU (실측 1줄=314,603 EMU) 사용
    - cy는 절대 줄이지 않음 — 확장만 허용
    """
    import math as _math
    sp = _find_shape_by_id(root, sid)
    if sp is None: return
    spPr = sp.find(f"{{{_NS_P}}}spPr")
    xfrm = spPr.find(f"{{{_NS_A}}}xfrm") if spPr else None
    if xfrm is None: return
    ext = xfrm.find(f"{{{_NS_A}}}ext")
    if ext is None: return
    cx = int(ext.get('cx', 0))
    if cx == 0:
        if xml_path is not None and _materialize_layout_shape(root, sid, xml_path):
            xfrm = spPr.find(f"{{{_NS_A}}}xfrm")
            ext = xfrm.find(f"{{{_NS_A}}}ext") if xfrm else None
            if ext is None: return
            cx = int(ext.get('cx', 0))
        if cx == 0: return

    # 폰트 크기: rPr.sz 우선, 없으면 _BODY_TITLE_FONT_PT (레이아웃 상속 = body_title 기본값)
    rPr = next((r.find(f"{{{_NS_A}}}rPr") for r in sp.findall(f".//{{{_NS_A}}}r")
                if r.find(f"{{{_NS_A}}}rPr") is not None), None)
    sz_h = (int(rPr.get('sz')) if rPr is not None and rPr.get('sz') else None)
    font_pt = (sz_h / 100.0) if sz_h else default_font_pt

    # 줄 수 계산: CJK 한글 실측 기준 0.6 보정 (기존 0.8은 과대 추정)
    cx_pt = cx / 12700.0
    box_cap = cx_pt / (font_pt * 0.6)

    def _cw(c: str) -> float:
        if '가' <= c <= '힣' or '一' <= c <= '鿿': return 1.0
        if c in (' ', '\t', '\n'): return 0.35
        return 0.55

    w = sum(_cw(c) for c in (text or ""))
    lines = max(1, _math.ceil(w / box_cap))

    # cy 계산: 실측 1줄=314,603 EMU 기준 사용 (font_pt 기반 계산보다 안정적)
    new_cy = lines * _SIDEBAR_LINE_HEIGHT_EMU
    orig_cy = int(ext.get('cy', 0))
    if new_cy > orig_cy:
        ext.set('cy', str(new_cy))
        # cy 확장 후 인접 shape 겹침 감지 및 자동 해소
        _resolve_overlaps(root, sid)


def _slide_set_helper(root, ns_p, ns_a, sid, text):
    """편집기 공통 shape 텍스트 설정 헬퍼. 텍스트 설정 후 자동 줄바꿈 발생 시 cy 동적 확장."""
    import copy as _copy
    sp = _find_shape_by_id(root, sid)
    if sp is None: return
    txBody = sp.find(f"{{{ns_p}}}txBody")
    if txBody is None: return
    orig_rPr = next((_copy.deepcopy(r.find(f"{{{ns_a}}}rPr")) for r in sp.findall(f".//{{{ns_a}}}r")
                     if r.find(f"{{{ns_a}}}rPr") is not None), None)
    if orig_rPr is not None: orig_rPr.set("lang", _ppt_lang()); orig_rPr.set("dirty","0")
    import copy as _copy2
    for p in txBody.findall(f"{{{ns_a}}}p"):
        for r in p.findall(f"{{{ns_a}}}r"): p.remove(r)
        end = p.find(f"{{{ns_a}}}endParaRPr"); idx = list(p).index(end) if end is not None else len(p)
        r_new = ET.Element(f"{{{ns_a}}}r")
        if orig_rPr: r_new.append(_copy2.deepcopy(orig_rPr))
        ET.SubElement(r_new, f"{{{ns_a}}}t").text = text; p.insert(idx, r_new); break
    # 자동 줄바꿈 발생 시 cy 동적 확장
    if text:
        _auto_resize_textbox(root, sid, text)


def _edit_slide35(xml_path: Path, slide_plan: dict) -> None:
    """slide35 (Before→After 버블): Zone1+2 공통 + Zone3 버블 키워드."""
    content = slide_plan.get("content", {})
    body    = content.get("body", {})
    before  = (content.get("before") or (body.get("before") if isinstance(body,dict) else None) or content.get("as_is",[]))[:3]
    after   = (content.get("after")  or (body.get("after")  if isinstance(body,dict) else None) or content.get("to_be",[]))[:4]
    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, "slide35.xml")
    ns_p, ns_a = _NS_P, _NS_A
    _s35 = _load_slide_shape_ids().get("slide35", {})
    _slide_set_helper(root, ns_p, ns_a, _s35.get("before_label_id", "39"), content.get("before_label","Before"))
    _slide_set_helper(root, ns_p, ns_a, _s35.get("after_label_id",  "38"), content.get("after_label","After"))
    _b35 = _s35.get("bubble", {"width_emu":1500000,"font_pt":16,"max_lines":2})
    _bef_ids35 = _s35.get("before_ids", ["17","16","35"])
    _aft_ids35 = _s35.get("after_ids",  ["18","21","13","19"])
    # Before 원형: ID 17(좌), 16(중), 35(우) — 템플릿 x좌표 기준 Before 영역(x < 5,943,887)
    for i, sid in enumerate(_bef_ids35):
        _slide_set_helper(root, ns_p, ns_a, sid, _truncate_to_lines(before[i] if i<len(before) else "",_b35["width_emu"],_b35["font_pt"],_b35["max_lines"]))
    # After 원형: ID 18(상), 21(중), 13(하좌), 19(하우) — 템플릿 x좌표 기준 After 영역(x ≥ 5,943,887)
    for i, sid in enumerate(_aft_ids35):
        _slide_set_helper(root, ns_p, ns_a, sid, _truncate_to_lines(after[i] if i<len(after) else "",_b35["width_emu"],_b35["font_pt"],_b35["max_lines"]))
    _clear_residual_placeholders(root); _write_xml(root, xml_path)


def _edit_slide36(xml_path: Path, slide_plan: dict) -> None:
    """slide36 (As-is/To-be 벤다이어그램): Zone1+2 공통 + Zone3 벤다이어그램 키워드.
    As-is 원형: 13,15,16 / 라벨: 7(14pt 1행) / 설명 본문: 29(12pt 다행)
    To-be 원형: 21,22,23 / 라벨: 27 / 설명 본문: 30
    """
    content = slide_plan.get("content", {})
    body    = content.get("body", {})
    as_is   = (content.get("as_is")  or (body.get("as_is")  if isinstance(body,dict) else None) or content.get("before",[]))[:4]
    to_be   = (content.get("to_be")  or (body.get("to_be")  if isinstance(body,dict) else None) or content.get("after", []))[:4]
    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, "slide36.xml")
    ns_p, ns_a = _NS_P, _NS_A
    _s36 = _load_slide_shape_ids().get("slide36", {})
    _bbl36 = _s36.get("bubble", {"width_emu":1500000,"font_pt":16,"max_lines":2})
    _dsc36 = _s36.get("desc",   {"width_emu":1500000,"font_pt":16,"max_lines":3})
    _asis_ids36 = _s36.get("asis_ids",     ["13","15","16"])
    _asis_dsc36 = _s36.get("asis_desc_id", "29")   # 설명 본문(12pt 다행) — 라벨(7) 아님
    _asis_lbl36 = _s36.get("asis_label_id", "7")    # 'explain' 라벨(14pt 1행)
    _tobe_ids36 = _s36.get("tobe_ids",     ["21","22","23"])
    _tobe_dsc36 = _s36.get("tobe_desc_id", "30")
    _tobe_lbl36 = _s36.get("tobe_label_id", "27")
    _lbl36 = _s36.get("label", {"width_emu":2514600,"font_pt":14,"max_lines":1})
    # 비교 라벨: compare_labels 제공 시 좌/우 라벨 교체(예: Anthropic/OpenAI), 없으면 템플릿 As-is/To-be 유지
    _cmp_labels = content.get("compare_labels") or content.get("labels") or []
    _cmp_ids = _zone("slide36.xml").get("body", {}).get("compare_label", ["11", "28"])
    for i, sid in enumerate(_cmp_ids):
        if i < len(_cmp_labels) and str(_cmp_labels[i]).strip():
            _slide_set_helper(root, ns_p, ns_a, sid, str(_cmp_labels[i]).strip())
    # 원형 키워드는 작아 '…'가 지저분 → ellipsis 없이 깔끔히(짧은 키워드 권장)
    for i, sid in enumerate(_asis_ids36):
        _slide_set_helper(root, ns_p, ns_a, sid, _truncate_to_lines(as_is[i] if i<len(as_is) else "",_bbl36["width_emu"],_bbl36["font_pt"],_bbl36["max_lines"], ellipsis=False))
    for i, sid in enumerate(_tobe_ids36):
        _slide_set_helper(root, ns_p, ns_a, sid, _truncate_to_lines(to_be[i] if i<len(to_be) else "",_bbl36["width_emu"],_bbl36["font_pt"],_bbl36["max_lines"], ellipsis=False))
    # 우측 설명 본문(29/30, 12pt 다행): explains[0/1] 우선 → 없으면 as_is[3]/to_be[3] fallback.
    # ⚠ 과거 버그: explain을 as_is[3]→shape7(작은 'explain' 라벨)에 써서, 콘텐츠가 원형 키워드 3개만 주면
    #   실제 본문(29/30)은 영구 공백이 됐다(엔진이 29/30을 아예 안 건드림). 본문=29/30, 라벨=7/27로 분리.
    _explains36 = _coerce_list(content.get("explains") or (body.get("explains") if isinstance(body, dict) else None) or [])
    _asis_ex = _explains36[0] if len(_explains36) > 0 else (as_is[3] if len(as_is) > 3 else "")
    _tobe_ex = _explains36[1] if len(_explains36) > 1 else (to_be[3] if len(to_be) > 3 else "")
    _slide_set_helper(root, ns_p, ns_a, _asis_dsc36, _truncate_to_lines(_asis_ex, _dsc36["width_emu"], _dsc36["font_pt"], _dsc36["max_lines"]))
    _slide_set_helper(root, ns_p, ns_a, _tobe_dsc36, _truncate_to_lines(_tobe_ex, _dsc36["width_emu"], _dsc36["font_pt"], _dsc36["max_lines"]))
    # 라벨(7/27, 1행): explain_labels 우선 → compare_labels+' 접근' → 둘 다 없으면 빈 값('explain' 영문 placeholder 제거)
    _ex_labels = content.get("explain_labels") or []
    def _ex_label36(i):
        if i < len(_ex_labels) and str(_ex_labels[i]).strip(): return str(_ex_labels[i]).strip()
        if i < len(_cmp_labels) and str(_cmp_labels[i]).strip(): return f"{str(_cmp_labels[i]).strip()} 접근"
        return ""
    _slide_set_helper(root, ns_p, ns_a, _asis_lbl36, _truncate_to_lines(_ex_label36(0), _lbl36["width_emu"], _lbl36["font_pt"], _lbl36["max_lines"], ellipsis=False))
    _slide_set_helper(root, ns_p, ns_a, _tobe_lbl36, _truncate_to_lines(_ex_label36(1), _lbl36["width_emu"], _lbl36["font_pt"], _lbl36["max_lines"], ellipsis=False))
    _clear_residual_placeholders(root); _write_xml(root, xml_path)


def _edit_slide39(xml_path: Path, slide_plan: dict) -> None:
    """slide39 (4열 상세 흐름도): Zone1+2 공통 + Zone3 keyword/solution/detail/service 4컬럼."""
    content   = slide_plan.get("content", {})
    body      = content.get("body", {})
    keywords  = (content.get("keywords")  or (body.get("keywords")  if isinstance(body,dict) else None) or content.get("items") or [])[:3]
    solutions = (content.get("solutions") or (body.get("solutions") if isinstance(body,dict) else None) or content.get("descriptions") or [])[:3]
    details   = (content.get("details")   or (body.get("details")   if isinstance(body,dict) else None) or [])[:3]
    services  = (content.get("services")  or content.get("details2")
                 or (body.get("services") if isinstance(body,dict) else None) or [])[:3]
    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, "slide39.xml")
    ns_p, ns_a = _NS_P, _NS_A
    _s39 = _load_slide_shape_ids().get("slide39", {})
    _kw39  = _s39.get("keyword",  {"width_emu":1500000,"font_pt":12,"max_lines":2})
    _sol39 = _s39.get("solution", {"width_emu":2000000,"font_pt":12,"max_lines":3})
    _det39 = _s39.get("detail",   {"width_emu":2000000,"font_pt":12,"max_lines":3})
    _svc39 = _s39.get("service",  {"width_emu":2000000,"font_pt":12,"max_lines":3})
    _kw_ids39  = _s39.get("keyword_ids",  ["13","14","15"])
    _sol_ids39 = _s39.get("solution_ids", ["7","10","11"])
    _det_ids39 = _s39.get("detail_ids",   ["16","17","18"])
    _svc_ids39 = _s39.get("service_ids",  ["19","24","25"])
    for i, sid in enumerate(_kw_ids39):
        _slide_set_helper(root, ns_p, ns_a, sid, _truncate_to_lines(keywords[i] if i<len(keywords) else "",_kw39["width_emu"],_kw39["font_pt"],_kw39["max_lines"]))
    for i, sid in enumerate(_sol_ids39):
        _slide_set_helper(root, ns_p, ns_a, sid, _truncate_to_lines(solutions[i] if i<len(solutions) else "",_sol39["width_emu"],_sol39["font_pt"],_sol39["max_lines"]))
    import copy as _copy_d
    for i, sid in enumerate(_det_ids39):
        raw = details[i] if i < len(details) else ""
        # 다중 항목 지원: list 또는 줄바꿈 문자열 → 여러 단락 렌더 (p32 — 항목 2개 이상 가능)
        if isinstance(raw, list):
            sub = [str(s).strip() for s in raw if str(s).strip()]
        else:
            sub = [s.strip() for s in str(raw).split("\n") if s.strip()]
        sub = [_truncate_to_lines(s, _det39["width_emu"], _det39["font_pt"], 2)
               for s in sub[:_det39.get("max_lines", 3)]] or [""]
        sp_det = _find_shape_by_id(root, sid)
        if sp_det is None: continue
        txBody = sp_det.find(f"{{{ns_p}}}txBody")
        if txBody is None: continue
        paras = txBody.findall(f"{{{ns_a}}}p")
        orig_rPr = next((_copy_d.deepcopy(r.find(f"{{{ns_a}}}rPr"))
                         for p in paras for r in p.findall(f"{{{ns_a}}}r")
                         if r.find(f"{{{ns_a}}}rPr") is not None), None)
        if orig_rPr is not None:
            orig_rPr.set("lang", _ppt_lang()); orig_rPr.set("dirty", "0")
        # 항목 수만큼 단락 채우고(서식 보존) 나머지 단락 제거 — 빈 줄·placeholder 잔여 방지
        for k, para in enumerate(paras):
            if k < len(sub):
                for r in para.findall(f"{{{ns_a}}}r"): para.remove(r)
                end = para.find(f"{{{ns_a}}}endParaRPr")
                idx = list(para).index(end) if end is not None else len(para)
                r_new = ET.Element(f"{{{ns_a}}}r")
                if orig_rPr is not None: r_new.append(_copy_d.deepcopy(orig_rPr))
                ET.SubElement(r_new, f"{{{ns_a}}}t").text = sub[k]
                para.insert(idx, r_new)
            else:
                txBody.remove(para)
    for i, sid in enumerate(_svc_ids39):
        _slide_set_helper(root, ns_p, ns_a, sid, _truncate_to_lines(services[i] if i<len(services) else "",_svc39["width_emu"],_svc39["font_pt"],_svc39["max_lines"]))
    _clear_residual_placeholders(root); _write_xml(root, xml_path)


def _edit_slide24(xml_path: Path, slide_plan: dict) -> None:
    """slide24/22 (2블록 텍스트): Zone1+2 공통 + Zone3 본문(bullets/body).
    Zone1: ID=8(대제목), ID=9(중제목 인덱스 prefix 포함)
    Zone2: body_title (slide24=ID=16, slide22=ID=14)
    Zone3: ID=7(bullets 또는 body), ID=10(body2)"""
    import copy as _copy
    ns_p, ns_a = _NS_P, _NS_A
    tmpl    = slide_plan.get("template_file", "slide24.xml")
    content = slide_plan.get("content", {})
    bullets = content.get("bullets") or content.get("items") or []
    body    = content.get("body", "")

    try:
        tree = ET.parse(xml_path)
        root = tree.getroot()
    except ET.ParseError:
        return

    # Zone1 + Zone2: ID=8(대제목), ID=9(중제목), body_title, body_desc 공통 처리
    _apply_common_zones(root, slide_plan, tmpl)

    def _set_text(sp, text):
        txBody = sp.find(f"{{{ns_p}}}txBody") if sp is not None else None
        if txBody is None: return
        orig_rPr = None
        for r in sp.findall(f".//{{{ns_a}}}r"):
            rPr_e = r.find(f"{{{ns_a}}}rPr")
            if rPr_e is not None:
                orig_rPr = _copy.deepcopy(rPr_e)
                orig_rPr.set("lang", _ppt_lang()); orig_rPr.set("dirty","0")
                break
        for p in txBody.findall(f"{{{ns_a}}}p"):
            for r in p.findall(f"{{{ns_a}}}r"): p.remove(r)
            end = p.find(f"{{{ns_a}}}endParaRPr")
            idx = list(p).index(end) if end is not None else len(p)
            r_new = ET.Element(f"{{{ns_a}}}r")
            if orig_rPr: r_new.append(_copy.deepcopy(orig_rPr))
            ET.SubElement(r_new, f"{{{ns_a}}}t").text = text
            p.insert(idx, r_new)
            break

    def _set_bullets(sp, items):
        txBody = sp.find(f"{{{ns_p}}}txBody") if sp is not None else None
        if txBody is None: return
        orig_rPr = None
        for r in sp.findall(f".//{{{ns_a}}}r"):
            rPr_e = r.find(f"{{{ns_a}}}rPr")
            if rPr_e is not None:
                orig_rPr = _copy.deepcopy(rPr_e)
                orig_rPr.set("lang", _ppt_lang()); orig_rPr.set("dirty","0")
                break
        paras = txBody.findall(f"{{{ns_a}}}p")
        for i, para in enumerate(paras):
            for r in para.findall(f"{{{ns_a}}}r"): para.remove(r)
            text = items[i] if i < len(items) else ""
            end = para.find(f"{{{ns_a}}}endParaRPr")
            idx = list(para).index(end) if end is not None else len(para)
            r_new = ET.Element(f"{{{ns_a}}}r")
            if orig_rPr: r_new.append(_copy.deepcopy(orig_rPr))
            ET.SubElement(r_new, f"{{{ns_a}}}t").text = text
            para.insert(idx, r_new)

    # Zone3: 본문 내용 (_apply_common_zones의 body_desc 이후 override)
    # 템플릿별 shape_id 우선 조회 (slide22는 하단 배너 ID=11 사용) → 없으면 slide24 기본값
    _all_sids = _load_slide_shape_ids()
    _s24 = _all_sids.get(tmpl.replace(".xml", ""), {}) or _all_sids.get("slide24", {})
    sp7  = _find_shape_by_id(root, _s24.get("body1_id", "7"))
    sp10 = _find_shape_by_id(root, _s24.get("body2_id", "10"))

    if bullets and sp7:
        _set_bullets(sp7, bullets)
    elif body and sp7:
        _set_text(sp7, body)
    if body and sp10:
        _set_text(sp10, body)

    # 이미지 슬롯(zone_map image_slots)에 가이드 캡션 주입 — slide22 등 (slide24는 image_slots 없음 → 미적용)
    for _img_sid in _zone(tmpl).get("body", {}).get("image_slots", []):
        _set_image_slot_text(root, _img_sid, _image_caption(content, 0))

    _clear_residual_placeholders(root)
    _write_xml(root, xml_path)


def _edit_slide30(xml_path: Path, slide_plan: dict) -> None:
    """slide30 (4단계 스텝): Zone1+2 공통 + Zone3 Step1~4 헤더 + 본문 동적 텍스트박스.

    steps 항목이 '헤더 — 본문' 형식이면 '—' 기준으로 분리:
      헤더(12자 이하) → ID=28~31 Step 박스
      본문 → 각 컬럼 본문 영역에 동적 sp 추가
    step_headers + step_descs로도 받음.
    """
    import copy as _copy
    content = slide_plan.get("content", {})
    body    = content.get("body", {})
    steps   = (content.get("steps") or (body.get("steps") if isinstance(body,dict) else None)
               or content.get("bullets") or content.get("items") or [])

    # 헤더/본문 분리: plan에 step_headers 있으면 그걸 쓰고, 없으면 steps에서 '—'로 분리
    step_headers = content.get("step_headers") or []
    step_descs   = content.get("step_descs") or []
    if not step_headers:
        for s in steps:
            parts = s.split("—", 1) if "—" in s else [s, ""]
            step_headers.append(parts[0].strip())
            step_descs.append(parts[1].strip() if len(parts) > 1 else "")

    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, "slide30.xml")
    ns_p, ns_a = _NS_P, _NS_A

    _s30 = _load_slide_shape_ids().get("slide30", {})
    _hdr_ids30 = _s30.get("header_ids", ["28","29","30","31"])
    col_xs    = _s30.get("col_xs",  [2_657_592, 5_145_661, 7_445_690, 9_887_657])
    body_xs   = _s30.get("body_xs", col_xs)   # 본문/밴드 전용 x (배경도형 기준 중앙 정렬)
    col_cx    = _s30.get("col_cx",  2_240_346)
    body_cxs  = _s30.get("body_cxs", [col_cx] * 4)  # 컬럼별 본문/밴드 너비 (동일 여백 적용)
    body_y    = _s30.get("body_y",  1_575_000)
    body_cy   = _s30.get("body_cy", 5_100_000)
    desc_cy   = _s30.get("desc_cy", body_cy)   # step_descs 영역 높이 (harness에 없으면 전체 사용)
    meta_y    = _s30.get("meta_y",  None)       # band 시작 y
    _b1h      = _s30.get("band1_h",  375_000)
    _b2h      = _s30.get("band2_h",  375_000)
    band_hs   = [_b1h, _b2h,
                 _s30.get("band3_h", _b2h),
                 _s30.get("band4_h", _b2h)]
    band_gap  = _s30.get("band_gap",  50_000)
    band_bg   = _s30.get("band_bg",   "F4F6FE")
    band_bdr  = _s30.get("band_border","3C41E6")
    band_lsz  = _s30.get("band_label_sz", 900)
    band_vsz  = _s30.get("band_value_sz", 1000)
    _sth30    = _s30.get("step_header", {"width_emu":1800000,"font_pt":16,"max_lines":2})
    _sbfsz30  = _s30.get("step_body_font_sz", 1100)
    step_meta = content.get("step_meta") or []

    for i, sid in enumerate(_hdr_ids30):
        _slide_set_helper(root, ns_p, ns_a, sid,
                          _truncate_to_lines(step_headers[i] if i<len(step_headers) else "", _sth30["width_emu"], _sth30["font_pt"], _sth30["max_lines"]))

    # ID=19: 장식용 사각형 — Vision Fix가 임의로 채우는 것을 방지하기 위해 항상 비움
    _slide_set_helper(root, ns_p, ns_a, "19", "")

    # spTree 직접 접근
    spTree = root.find(f"{{{ns_p}}}cSld/{{{ns_p}}}spTree")
    if spTree is None:
        spTree = root.find(f".//{{{ns_p}}}spTree")

    # 기존 step_body / step_meta / band 제거 (재편집 시 중복 방지)
    _dynamic_ids30 = (
        {str(400+k) for k in range(4)} |
        {str(410+k) for k in range(4)} |
        {str(420+k) for k in range(4)} |
        {str(430+k) for k in range(4)} |
        {str(440+k) for k in range(4)} |
        {str(450+k) for k in range(4)} |
        {str(500+k) for k in range(10)} |
        {str(600+k) for k in range(32)}
    )
    if spTree is not None:
        for sp in list(spTree):
            cnvpr = sp.find(f".//{{{ns_p}}}cNvPr")
            if cnvpr is not None and cnvpr.get("id", "") in _dynamic_ids30:
                spTree.remove(sp)

    lang = _ppt_lang()
    _font_attrs = 'panose="02000703000000020004" pitchFamily="2" charset="-127"'

    # ── 컬럼별 desc 박스 (ID 500~503) — 하네스: col_desc_* ──────────────
    # step_descs[i] → 각 컬럼(1단계~4단계) 헤더 아래 개별 텍스트박스
    _cd_id0  = int(_s30.get("col_desc_id_start", 500))
    _cd_y    = _s30.get("col_desc_y",   1280000)
    _cd_cy   = _s30.get("col_desc_cy",  2750000)
    _cd_fsz  = _s30.get("col_desc_font_sz", 1000)
    _cd_lnspc = _s30.get("col_desc_lnspc_pct", 100000)  # 기본 100%, 하네스에서 변경 가능
    _cd_ppr  = f'<a:pPr><a:lnSpc><a:spcPct val="{_cd_lnspc}"/></a:lnSpc></a:pPr>'
    _cd_rpr  = (f'<a:rPr lang="{lang}" altLang="en-US" sz="{_cd_fsz}" dirty="0">'
                f'<a:solidFill><a:schemeClr val="tx1"><a:lumMod val="95000"/><a:lumOff val="5000"/></a:schemeClr></a:solidFill>'
                f'<a:latin typeface="Pretendard" {_font_attrs}/>'
                f'<a:ea typeface="Pretendard" {_font_attrs}/>'
                f'<a:cs typeface="Pretendard" {_font_attrs}/></a:rPr>')
    # ── 헤더(ID=28~31) spAutoFit 활성화 — 텍스트 넘치면 높이 자동 확장 ──
    _hdr_max_cy = 0
    for sid in _hdr_ids30:
        sp_hdr = _find_shape_by_id(root, sid)
        if sp_hdr is not None:
            _apply_spAutoFit(root, sid)
            # 실제 텍스트 길이 기반 예상 cy 계산 (2줄 기준)
            txt_hdr = next((t.text for t in sp_hdr.findall(f".//{{{ns_a}}}t") if t.text), "")
            est = _estimate_lines(txt_hdr, _sth30.get("width_emu", 1800000), _sth30.get("font_pt", 16))
            est_cy = max(338554, est * 314603)
            _hdr_max_cy = max(_hdr_max_cy, est_cy)
    # 모든 헤더 높이를 최대값으로 통일
    if _hdr_max_cy > 0:
        _equalize_row_heights(root, [[sid] for sid in _hdr_ids30])

    # ── 동적 bullet 계산 ──────────────────────────────────────────────────────
    # col_desc 박스는 항상 grid 바로 위까지 꽉 채움 (grid_y는 harness 고정값)
    _g_y_fixed   = _s30.get("grid_y", 4100000)
    _g_gap       = _s30.get("grid_gap", 50000)
    _unified_cy  = _g_y_fixed - _cd_y - _g_gap   # col_desc 박스 고정 높이
    _unified_cy  = max(_unified_cy, _cd_cy)       # 최소값 보장
    _cd_font_pt  = _cd_fsz / 100.0                # 1000 hundredths → 10 pt
    _cd_line_h   = int(_cd_font_pt * 12700 * (_cd_lnspc / 100000))  # EMU/line
    _cd_insets   = 91440 + 45720                  # tIns + bIns
    _cd_max_fit  = max(2, (_unified_cy - _cd_insets) // max(1, _cd_line_h))  # 박스에 최대 들어가는 줄 수
    _bullet_sym  = _s30.get("step_bullet", "▪ ")

    # 각 컬럼 bullet 목록 구성 (min 3개, max _cd_max_fit 줄 이하)
    _col_bullets = []
    for ci in range(4):
        _cx_ci   = body_cxs[ci] if ci < len(body_cxs) else col_cx
        desc_txt = step_descs[ci] if ci < len(step_descs) else ""
        raw      = [l.strip() for l in str(desc_txt).split('\n') if l.strip()] if desc_txt else []
        while len(raw) < 3:
            raw.append("")
        # 렌더 줄 수 기준으로 최대치 이하 bullet만 유지
        kept, used = [], 0
        for b in raw:
            btext    = _bullet_sym + b if b else ""
            rendered = _estimate_lines(btext, _cx_ci, _cd_font_pt) if btext else 1
            if used + rendered > _cd_max_fit:
                break
            kept.append(b)
            used += rendered
        _col_bullets.append(kept)

    if spTree is not None:
        for ci in range(4):
            col_x  = body_xs[ci] if ci < len(body_xs) else col_xs[ci]
            col_cx = body_cxs[ci] if ci < len(body_cxs) else col_cx
            bullets = _col_bullets[ci]
            if any(bullets):
                _paras = "".join(
                    f'<a:p>{_cd_ppr}<a:r>{_cd_rpr}<a:t>{_bullet_sym}{b}</a:t></a:r></a:p>'
                    for b in bullets if b
                )
            else:
                _paras = '<a:p/>'
            sp_xml = f"""<p:sp xmlns:p="{ns_p}" xmlns:a="{ns_a}">
  <p:nvSpPr><p:cNvPr id="{_cd_id0 + ci}" name="col_desc_{ci}"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{col_x}" y="{_cd_y}"/><a:ext cx="{col_cx}" cy="{_unified_cy}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" lIns="91440" tIns="91440" rIns="91440" bIns="45720"/>
    <a:lstStyle/>{_paras}
  </p:txBody>
</p:sp>"""
            spTree.append(ET.fromstring(sp_xml))

    # ── 하단 grid 4열×4행 박스 (zone_lower_grid) — 하네스: grid_* ──────────
    # 컬럼별 독립 x/cx: 배경 도형(흰/회색 교대) 기준 여백 적용값
    _g_col_xs  = _s30.get("grid_col_xs",  body_xs)   # 컬럼별 x (배경 도형 여백 포함)
    _g_col_cxs = _s30.get("grid_col_cxs", body_cxs)  # 컬럼별 cx
    _g_y      = _g_y_fixed  # harness 고정값 (col_desc 박스가 여기까지 꽉 채움)
    _g_tcy    = _s30.get("grid_total_cy", 2750000)
    _g_rgap   = _s30.get("grid_row_gap",  36000)
    _g_cols   = _s30.get("grid_cols",     4)
    _g_rows   = _s30.get("grid_rows",     4)
    _g_bgs    = _s30.get("grid_bg_colors", ["FFFFFF", "F0F0F0", "FFFFFF", "F0F0F0"])
    _g_bdr    = _s30.get("grid_border",   "3C41E6")
    _g_lsz    = _s30.get("grid_label_sz", 800)
    _g_vsz    = _s30.get("grid_font_sz",  1000)
    _g_sid0   = int(_s30.get("grid_start_id", 600))

    _cell_cy = (_g_tcy - _g_rgap * (_g_rows - 1)) // _g_rows

    if step_meta and spTree is not None:
        # 행 레이블: 첫 번째 dict의 키 순서로 결정
        _all_labels = []
        for m in step_meta:
            if isinstance(m, dict):
                for k in m.keys():
                    if k not in _all_labels:
                        _all_labels.append(k)
        row_labels = _all_labels[:_g_rows]

        for col_i in range(_g_cols):
            meta = step_meta[col_i] if col_i < len(step_meta) else {}
            _bg = _g_bgs[col_i % len(_g_bgs)]
            _bx = _g_col_xs[col_i] if col_i < len(_g_col_xs) else body_xs[col_i]
            _cell_cx = _g_col_cxs[col_i] if col_i < len(_g_col_cxs) else body_cxs[col_i]
            for row_i, label in enumerate(row_labels):
                val = str(meta.get(label, "")) if isinstance(meta, dict) else ""
                _cy_off = row_i * (_cell_cy + _g_rgap)
                _by = _g_y + _cy_off
                bid = _g_sid0 + row_i * _g_cols + col_i
                _lrpr = (f'<a:rPr lang="{lang}" altLang="en-US" sz="{_g_lsz}" b="1" dirty="0">'
                         f'<a:solidFill><a:srgbClr val="{_g_bdr}"/></a:solidFill>'
                         f'<a:latin typeface="Pretendard" {_font_attrs}/>'
                         f'<a:ea typeface="Pretendard" {_font_attrs}/>'
                         f'<a:cs typeface="Pretendard" {_font_attrs}/></a:rPr>')
                _vrpr = (f'<a:rPr lang="{lang}" altLang="en-US" sz="{_g_vsz}" dirty="0">'
                         f'<a:solidFill><a:srgbClr val="111111"/></a:solidFill>'
                         f'<a:latin typeface="Pretendard" {_font_attrs}/>'
                         f'<a:ea typeface="Pretendard" {_font_attrs}/>'
                         f'<a:cs typeface="Pretendard" {_font_attrs}/></a:rPr>')
                sp_xml = f"""<p:sp xmlns:p="{ns_p}" xmlns:a="{ns_a}">
  <p:nvSpPr><p:cNvPr id="{bid}" name="grid_{row_i}_{col_i}"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{_bx}" y="{_by}"/><a:ext cx="{_cell_cx}" cy="{_cell_cy}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{_bg}"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="{_g_bdr}"/></a:solidFill></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" lIns="{_s30.get('grid_cell_lins',91440)}" tIns="{_s30.get('grid_cell_tins',91440)}" rIns="{_s30.get('grid_cell_rins',91440)}" bIns="{_s30.get('grid_cell_bins',91440)}" anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r>{_lrpr}<a:t>{label}</a:t></a:r></a:p>
    <a:p><a:pPr algn="ctr"/><a:r>{_vrpr}<a:t>{val}</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""
                spTree.append(ET.fromstring(sp_xml))

    _clear_residual_placeholders(root); _write_xml(root, xml_path)


def _edit_slide32(xml_path: Path, slide_plan: dict) -> None:
    """slide32 (상단텍스트+하단3열): Zone1+2 공통 + Zone3 본문설명글(#26) + 하단3열(#19,20) bullets/이미지."""
    content = slide_plan.get("content", {})
    # body > bullets > section_desc 순서로 ID=26에 넣을 텍스트 결정
    body = (content.get("body")
            or "\n".join(content.get("bullets", []))
            or content.get("section_desc", ""))
    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, "slide32.xml")
    _s32 = _load_slide_shape_ids().get("slide32", {})
    _md32_id = _s32.get("main_desc_id", "26")
    _md32    = _s32.get("main_desc", {"width_emu":4000000,"font_pt":12,"max_lines":8})
    if body:
        _slide_set_helper(root, _NS_P, _NS_A, _md32_id,
                          _truncate_to_lines(body, _md32["width_emu"], _md32["font_pt"], _md32["max_lines"]))
    _clear_residual_placeholders(root); _write_xml(root, xml_path)


def _edit_slide38(xml_path: Path, slide_plan: dict) -> None:
    """slide38 (3행 흐름도 keyword→Solution→Service): Zone1+2 공통 + Zone3 흐름 박스."""
    content   = slide_plan.get("content", {})
    body      = content.get("body", {})
    keywords  = (content.get("keywords")  or (body.get("keywords")  if isinstance(body,dict) else None) or content.get("bullets",[]))[:3]
    solutions = (content.get("solutions") or (body.get("solutions") if isinstance(body,dict) else None) or content.get("items",  []))[:3]
    # 플랜이 'details'로 생성하는 경우가 많아 services 폴백으로 받음 (빈 3번째 컬럼 방지)
    services  = (content.get("services")  or content.get("details")
                 or (body.get("services") if isinstance(body,dict) else None)
                 or (body.get("details")  if isinstance(body,dict) else None) or [])[:3]
    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, "slide38.xml")
    ns_p, ns_a = _NS_P, _NS_A
    _s38 = _load_slide_shape_ids().get("slide38", {})
    _kw38  = _s38.get("keyword",  {"width_emu":1500000,"font_pt":12,"max_lines":2})
    _sol38 = _s38.get("solution", {"width_emu":2000000,"font_pt":12,"max_lines":3})
    _svc38 = _s38.get("service",  {"width_emu":2000000,"font_pt":12,"max_lines":3})
    _kw_ids38  = _s38.get("keyword_ids",  ["13","14","15"])
    _sol_ids38 = _s38.get("solution_ids", ["7","10","11"])
    _svc_ids38 = _s38.get("service_ids",  ["19","24","25"])
    for i, sid in enumerate(_kw_ids38):
        _slide_set_helper(root, ns_p, ns_a, sid,
                          _truncate_to_lines(keywords[i] if i<len(keywords) else "",_kw38["width_emu"],_kw38["font_pt"],_kw38["max_lines"]))
    for i, sid in enumerate(_sol_ids38):
        _slide_set_helper(root, ns_p, ns_a, sid,
                          _truncate_to_lines(solutions[i] if i<len(solutions) else "",_sol38["width_emu"],_sol38["font_pt"],_sol38["max_lines"]))
    for i, sid in enumerate(_svc_ids38):
        _slide_set_helper(root, ns_p, ns_a, sid,
                          _truncate_to_lines(services[i] if i<len(services) else "",_svc38["width_emu"],_svc38["font_pt"],_svc38["max_lines"]))
    _clear_residual_placeholders(root); _write_xml(root, xml_path)


def _edit_slots_slide(xml_path: Path, slide_plan: dict, tmpl_key: str) -> None:
    """slots 하네스 기반 범용 슬라이드 편집기.
    slide_shape_ids.json의 'slots' 배열을 읽어 _apply_slots로 위임.
    slide9-17 계열이 공통으로 사용한다."""
    content = slide_plan.get("content", {})
    body = content.get("body", {})
    if isinstance(body, dict):
        for k, v in body.items():
            content.setdefault(k, v)
    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, tmpl_key)
    slide_key = tmpl_key.replace(".xml", "")
    slots = _load_slide_shape_ids().get(slide_key, {}).get("slots", [])
    if slots:
        _apply_slots(root, content, slots, slide_plan=slide_plan)
    _clear_residual_placeholders(root); _write_xml(root, xml_path)


def _edit_slide9(xml_path: Path, slide_plan: dict) -> None:
    """slide9 (3열 이미지+제목+설명): Zone1+2 공통 + slots 하네스 기반."""
    _edit_slots_slide(xml_path, slide_plan, "slide9.xml")


def _edit_slide10(xml_path: Path, slide_plan: dict) -> None:
    """slide10 (3열+인사이트): Zone1+2 공통 + slots 하네스 기반."""
    _edit_slots_slide(xml_path, slide_plan, "slide10.xml")


def _edit_slide11(xml_path: Path, slide_plan: dict) -> None:
    """slide11 (3열 이미지카드): Zone1+2 공통 + slots 하네스 기반."""
    _edit_slots_slide(xml_path, slide_plan, "slide11.xml")


def _edit_slide12(xml_path: Path, slide_plan: dict) -> None:
    """slide12 (3열+서브헤딩): Zone1+2 공통 + slots 하네스 기반."""
    _edit_slots_slide(xml_path, slide_plan, "slide12.xml")


def _edit_slide14(xml_path: Path, slide_plan: dict) -> None:
    """slide14 (4열 아이콘+인사이트): Zone1+2 공통 + slots 하네스 기반."""
    _edit_slots_slide(xml_path, slide_plan, "slide14.xml")


def _edit_slide16(xml_path: Path, slide_plan: dict) -> None:
    """slide16 (4열 아이콘+서브헤딩): Zone1+2 공통 + slots 하네스 기반."""
    _edit_slots_slide(xml_path, slide_plan, "slide16.xml")


def _edit_slide17(xml_path: Path, slide_plan: dict) -> None:
    """slide17 (4열 이미지+서브헤딩): Zone1+2 공통 + slots 하네스 기반."""
    _edit_slots_slide(xml_path, slide_plan, "slide17.xml")


def _edit_slide25(xml_path: Path, slide_plan: dict) -> None:
    """slide25 (이미지+우측3열): descriptions[0] → overview, descriptions[1:] → item_descs로 분리 후 zonemap 위임."""
    content = slide_plan.get("content", {})
    descs = _coerce_list(content.get("descriptions") or content.get("bullets") or [])
    if descs:
        content = {**content, "overview": descs[0], "descriptions": descs[1:]}
        slide_plan = {**slide_plan, "content": content}
    _edit_zonemap_slide(xml_path, slide_plan)


def _edit_slide26(xml_path: Path, slide_plan: dict) -> None:
    """slide26 (이미지+3항목): descriptions[0] → main_desc, descriptions[1:] → item_descs로 분리 후 zonemap 위임."""
    content = slide_plan.get("content", {})
    descs = _coerce_list(content.get("descriptions") or content.get("bullets") or [])
    if descs:
        content = {**content, "main_desc": descs[0], "descriptions": descs[1:]}
        slide_plan = {**slide_plan, "content": content}
    _edit_zonemap_slide(xml_path, slide_plan)


def _edit_slide27(xml_path: Path, slide_plan: dict) -> None:
    """slide27 (이미지+우측3행, body_title ≤2줄 variant): zone map 기반 직접 편집."""
    _edit_zonemap_slide(xml_path, slide_plan)


def _edit_slide28(xml_path: Path, slide_plan: dict) -> None:
    """slide28 (이미지+우측3행, body_title ≥3줄 variant): zone map 기반 직접 편집."""
    _edit_zonemap_slide(xml_path, slide_plan)


def _edit_slide34(xml_path: Path, slide_plan: dict) -> None:
    """slide34 (2이미지+키워드): Zone1+2 공통 + 이미지슬롯×2 + 키워드×4 + 설명."""
    content  = slide_plan.get("content", {})
    body     = content.get("body", {})
    keywords = _coerce_list(content.get("keywords") or content.get("items") or
               (body.get("keywords") if isinstance(body, dict) else None) or [])
    descs    = _coerce_list(content.get("descriptions") or content.get("bullets") or
               (body.get("descriptions") if isinstance(body, dict) else None) or [])
    img_descs = _coerce_list(content.get("image_descriptions") or
                (body.get("image_descriptions") if isinstance(body, dict) else None) or [])
    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, "slide34.xml")
    import copy as _copy; ns_p, ns_a = _NS_P, _NS_A
    _s34 = _load_slide_shape_ids().get("slide34", {})
    _kw34  = _s34.get("keyword",   {"width_emu":1800000,"font_pt":12,"max_lines":2})
    _md34  = _s34.get("main_desc", {"width_emu":7000000,"font_pt":12,"max_lines":4})
    _img34 = _s34.get("img_slot",  {"width_emu":2800000,"font_pt":11,"max_lines":4})
    pics_ids     = _s34.get("pics_ids",     ["26","27"])
    main_desc_id = _s34.get("main_desc_id", "12")
    kw_ids       = _s34.get("kw_ids",       ["19","14","13","18"])
    for i, sid in enumerate(kw_ids):
        _slide_set_helper(root, ns_p, ns_a, sid,
                          _truncate_to_lines(keywords[i] if i < len(keywords) else "", _kw34["width_emu"], _kw34["font_pt"], _kw34["max_lines"]))
    _slide_set_helper(root, ns_p, ns_a, main_desc_id,
                      _truncate_to_lines(descs[0] if descs else "", _md34["width_emu"], _md34["font_pt"], _md34["max_lines"]))
    _img_fmt34 = _load_common_formatting().get("image_slot_format", "[이미지: {description}]")
    for i, sid in enumerate(pics_ids):
        img_txt = img_descs[i] if i < len(img_descs) else ""
        if img_txt:
            _set_image_slot_text(root, sid, _truncate_to_lines(_img_fmt34.format(description=img_txt), _img34["width_emu"], _img34["font_pt"], _img34["max_lines"]))
        else:
            _slide_set_helper(root, ns_p, ns_a, sid, "")
    _clear_residual_placeholders(root); _write_xml(root, xml_path)


def _edit_slide37(xml_path: Path, slide_plan: dict) -> None:
    """slide37 (3구역 텍스트): Zone1+2 공통 + 3개 텍스트 구역."""
    content = slide_plan.get("content", {})
    body    = content.get("body", {})
    bullets = _coerce_list(content.get("bullets") or content.get("items") or
              (body.get("bullets") if isinstance(body, dict) else None) or [])
    descs   = _coerce_list(content.get("descriptions") or
              (body.get("descriptions") if isinstance(body, dict) else None) or [])
    texts   = bullets if bullets else descs
    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, "slide37.xml")
    import copy as _copy; ns_p, ns_a = _NS_P, _NS_A
    _s37 = _load_slide_shape_ids().get("slide37", {})
    _d0_37  = _s37.get("desc_0",   {"width_emu":7000000,"font_pt":14,"max_lines":3})
    _d12_37 = _s37.get("desc_1_2", {"width_emu":4000000,"font_pt":12,"max_lines":4})
    desc_ids = _s37.get("desc_ids", ["12","48","49"])
    _slide_set_helper(root, ns_p, ns_a, desc_ids[0],
                      _truncate_to_lines(texts[0] if texts else "", _d0_37["width_emu"], _d0_37["font_pt"], _d0_37["max_lines"]))
    _slide_set_helper(root, ns_p, ns_a, desc_ids[1],
                      _truncate_to_lines(texts[1] if len(texts) > 1 else "", _d12_37["width_emu"], _d12_37["font_pt"], _d12_37["max_lines"]))
    _slide_set_helper(root, ns_p, ns_a, desc_ids[2],
                      _truncate_to_lines(texts[2] if len(texts) > 2 else "", _d12_37["width_emu"], _d12_37["font_pt"], _d12_37["max_lines"]))
    _clear_residual_placeholders(root); _write_xml(root, xml_path)


def _edit_slide42(xml_path: Path, slide_plan: dict) -> None:
    """slide42 (대형 본문): 템플릿 구조 그대로.
    - ID=28(body_desc, cx=8.8M, cy≈3줄): section_desc — 상세 설명 3줄 이내
    - ID=43/47/51(item_descs, cx=2.4M, cy≈3줄): descriptions 3개 — 각 원형 아래 content
    zone_map: body_desc="28", body.item_descs=["43","47","51"]
    harness/slide_shape_ids.json slide42.item_descs_ids / item_desc 우선 적용.
    """
    content = slide_plan.get("content", {})
    descs   = _coerce_list(content.get("descriptions") or [])
    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError: return
    _apply_common_zones(root, slide_plan, "slide42.xml")
    ns_p, ns_a = _NS_P, _NS_A
    # harness 우선, 동일 fallback
    _s42 = _load_slide_shape_ids().get("slide42", {})
    _item_ids = _s42.get("item_descs_ids", ["43", "47", "51"])
    _item_cfg  = _s42.get("item_desc", {"width_emu": 2_435_125, "font_pt": 12, "max_lines": 3})
    _cx42   = _item_cfg.get("width_emu", 2_435_125)
    _pt42   = _item_cfg.get("font_pt", 12)
    _lines42 = _item_cfg.get("max_lines", 3)
    # descriptions 3개 → 각 원형 아래 content 박스
    for i, sid in enumerate(_item_ids):
        txt = descs[i] if i < len(descs) else ""
        if txt:
            _slide_set_helper(root, ns_p, ns_a, sid,
                              _truncate_to_lines(txt, _cx42, _pt42, _lines42))
    _clear_residual_placeholders(root); _write_xml(root, xml_path)


# ── _SLIDE_EDITORS 자동 등록 ─────────────────────────────────────
# 함수명 규칙 _edit_slideN → slideN.xml 자동 매핑. 수동 누락 버그 방지.
import re as _re_auto
_SLIDE_EDITORS: dict = {}
for _fn_name, _fn_obj in list(globals().items()):
    _m = _re_auto.match(r'^_edit_(slide\d+)$', _fn_name)
    if _m:
        _SLIDE_EDITORS[f"{_m.group(1)}.xml"] = _fn_obj
# 예외: naming convention 밖의 케이스만 수동 등록
_SLIDE_EDITORS["slide22.xml"] = _edit_slide24          # slide22는 slide24 함수 재사용
_SLIDE_EDITORS["slide15.xml"] = _edit_slide15_v2      # slide13 패턴으로 재구현 (content/slide_plan 파라미터 버그 수정)
del _fn_name, _fn_obj, _m


# ── 존 맵 기반 제너릭 편집기 ──────────────────────────────────────

# 본문구역 하위 존 role → plan content 필드 + 채움 규칙
#   key      : content에서 리스트를 꺼낼 필드(우선순위)
#   maxlines : 셀 줄 수
#   empty    : True면 비워둠(아이콘 자리)
_ZONE_FILL_RULES: dict[str, dict] = {
    "item_titles":   {"keys": ["items", "item_titles", "keywords"], "cx": 2_400_000, "pt": 14, "lines": 2},
    "item_descs":    {"keys": ["descriptions", "descs", "bullets", "details"], "cx": 2_400_000, "pt": 12, "lines": 4},
    "image_slots":   {"keys": ["image_descriptions", "images"], "cx": 2_800_000, "pt": 11, "lines": 4,
                      "prefix": "🖼 "},   # 이미지 자리 → 어떤 이미지인지 설명 텍스트
    "icon_slots":    {"empty": True},
    "insights":      {"keys": ["insights", "insight"], "cx": 2_400_000, "pt": 11, "lines": 3},
    "keywords":      {"keys": ["keywords", "items"], "cx": 1_800_000, "pt": 12, "lines": 2},
    "flow_keyword":  {"keys": ["keywords"], "cx": 2_400_000, "pt": 12, "lines": 2},
    "flow_solution": {"keys": ["solutions"], "cx": 2_700_000, "pt": 12, "lines": 3},
    "flow_service":  {"keys": ["services", "details"], "cx": 3_500_000, "pt": 12, "lines": 2},
    "detail_items":  {"keys": ["details", "sub_items"], "cx": 1_600_000, "pt": 11, "lines": 5},
    "steps":         {"keys": ["steps"], "cx": 2_200_000, "pt": 12, "lines": 2},
    "sub_titles":    {"keys": ["sub_titles", "callouts", "labels"], "cx": 2_500_000, "pt": 12, "lines": 2},
    "explains":      {"keys": ["explains", "descriptions"], "cx": 2_500_000, "pt": 11, "lines": 3},
    "banner":        {"keys": ["banner", "insight", "body"], "cx": 9_500_000, "pt": 14, "lines": 2},
    "sub_heading":   {"keys": ["sub_heading", "subtitle"], "cx": 9_000_000, "pt": 16, "lines": 2},
    "overview":      {"keys": ["overview"], "cx": 7_000_000, "pt": 12, "lines": 3},
    "main_desc":     {"keys": ["main_desc", "overview"], "cx": 7_000_000, "pt": 12, "lines": 3},
}


def _make_sub_heading(content: dict) -> str:
    """plan content에서 sub_heading 자동 생성.
    형식: | {section_ref} {item1·item2·item3} — {sub_heading_tail}
    section_ref: section_title 앞 번호 + '.1' (예: '4.2.' → '4.2.1')
    꼬리(대시 뒤)는 plan의 sub_heading_tail(명사형 종결 요약구)에서만 가져온다.
    section_desc(=~합니다 완결문)에서 파생하지 않는다 — 파생하면 소제목이 서술형 종결이 되어 '명사형' 규칙 위반.
    """
    import re
    sec_title = content.get("section_title", "")
    items = _coerce_list(content.get("items") or content.get("item_titles") or [])

    if not items:
        return ""

    m = re.match(r'^(\d+(?:\.\d+)*)\.?\s', sec_title)
    section_ref = f"{m.group(1)}.1" if m else ""
    items_str = "·".join(str(x) for x in items[:3])

    head = f"| {section_ref} {items_str}".strip() if section_ref else f"| {items_str}"
    # 명사형 종결 요약구(tail)는 plan이 sub_heading_tail로 명시 제공. 미제공 시 head(키워드)만 반환.
    tail = str(content.get("sub_heading_tail") or "").strip()
    if tail:
        return f"{head} — {tail}"
    return head


def _coerce_list(val) -> list[str]:
    if val is None:
        return []
    if isinstance(val, str):
        return [val]
    if isinstance(val, list):
        out = []
        for v in val:
            if isinstance(v, dict):
                out.append(v.get("content") or v.get("text") or v.get("label") or "")
            else:
                out.append(str(v))
        return out
    return [str(val)]


def _image_caption(content: dict, idx: int = 0) -> str:
    """이미지 슬롯 가이드 캡션 생성 (생성 로직 — 빈 이미지영역 구조적 방지).
    plan의 image_descriptions[idx]가 있으면 그것을, 없으면 section_title/items에서 파생한다.
    항상 '[이미지: …]' 포맷으로 반환. AHE_PRINCIPLES §1(verify/correct) — plan 누락에도 가이드 보장."""
    fmt = _load_common_formatting().get("image_slot_format", "[이미지: {description}]")
    descs = content.get("image_descriptions") or content.get("images") or []
    if idx < len(descs) and str(descs[idx]).strip():
        return fmt.format(description=str(descs[idx]).strip())
    # fallback: 같은 열의 item/keyword(열별 distinct) → section_title → 범용. 이미지생성 프롬프트로 쓸 수 있게 서술형으로
    items = content.get("items") or content.get("keywords") or []
    if idx < len(items) and str(items[idx]).strip():
        return fmt.format(description=f"{str(items[idx]).strip()}의 핵심 구조와 구성요소를 노드·화살표로 도식화한 다이어그램")
    sec = content.get("section_title", "")
    bare = re.sub(r'^\d+(\.\d+)*[\s.。]*', '', sec).strip() if sec else ""
    if bare:
        return fmt.format(description=f"{bare}의 전체 흐름과 핵심 요소·관계를 시각화한 다이어그램")
    if items:
        return fmt.format(description=f"{str(items[0]).strip()}의 핵심 구조를 도식화한 다이어그램")
    return fmt.format(description="핵심 개념과 구성요소·관계를 나타낸 설명용 다이어그램")


def _set_image_slot_text(root, sid: str, text: str) -> None:
    """
    image_slots shape에 텍스트 + 올바른 서식을 직접 주입.
    원본 rPr이 없는 shape(템플릿 style 상속만 있는 경우)도 정상 처리.
    서식: sz=1100, 맑은 고딕, #3C41E6(파란색), 가운데 정렬
    """
    import copy as _copy
    ns_p, ns_a = _NS_P, _NS_A
    sp = _find_shape_by_id(root, sid)
    if sp is None: return
    txBody = sp.find(f"{{{ns_p}}}txBody")
    if txBody is None: return
    for p in txBody.findall(f"{{{ns_a}}}p"):
        pPr = p.find(f"{{{ns_a}}}pPr")
        if pPr is None:
            pPr = ET.Element(f"{{{ns_a}}}pPr"); p.insert(0, pPr)
        pPr.set("algn", "ctr")
        for r in p.findall(f"{{{ns_a}}}r"): p.remove(r)
        r_new = ET.Element(f"{{{ns_a}}}r")
        rPr = ET.SubElement(r_new, f"{{{ns_a}}}rPr",
                            lang=_ppt_lang(), altLang="en-US", sz="1100", dirty="0")
        fill = ET.SubElement(rPr, f"{{{ns_a}}}solidFill")
        ET.SubElement(fill, f"{{{ns_a}}}srgbClr", val="3C41E6")
        ET.SubElement(rPr, f"{{{ns_a}}}latin", typeface="맑은 고딕")
        ET.SubElement(rPr, f"{{{ns_a}}}ea", typeface="맑은 고딕")
        ET.SubElement(r_new, f"{{{ns_a}}}t").text = text
        end = p.find(f"{{{ns_a}}}endParaRPr")
        idx = list(p).index(end) if end is not None else len(p)
        p.insert(idx, r_new)
        break


def _edit_zonemap_slide(xml_path: Path, slide_plan: dict) -> None:
    """존 맵(layout_zone_map.json)을 읽어 본문구역 하위 존을 채우는 제너릭 편집기.
    전용 편집기가 없는 모든 본문 슬라이드(이미지 그리드 등)에 사용."""
    template_file = slide_plan.get("template_file", "")
    z = _zone(template_file)
    if not z:
        return _edit_layout_slide(xml_path, slide_plan)
    content = slide_plan.get("content", {})
    body = content.get("body", {})
    bdict = body if isinstance(body, dict) else {}

    try:
        tree = ET.parse(xml_path); root = tree.getroot()
    except ET.ParseError:
        return
    _apply_common_zones(root, slide_plan, template_file)
    ns_p, ns_a = _NS_P, _NS_A

    def fetch(keys: list[str]) -> list[str]:
        for k in keys:
            v = content.get(k)
            if not v and isinstance(bdict, dict):
                v = bdict.get(k)
            if v:
                return _coerce_list(v)
        return []

    zone_rules = _load_zone_fill_rules() or _ZONE_FILL_RULES
    body_zones = z.get("body", {})
    for role, ids in body_zones.items():
        if not ids:
            continue
        rule = zone_rules.get(role)
        if rule is None:
            continue
        if rule.get("empty"):
            for sid in ids:
                _slide_set_helper(root, ns_p, ns_a, sid, "")
            continue
        if role == "charts":
            continue  # 차트는 별도 처리(Excel)
        values = fetch(rule["keys"])
        # sub_heading: plan에 필드 없으면 items + section_desc로 자동 생성
        if role == "sub_heading" and not values:
            auto = _make_sub_heading(content)
            if auto:
                values = [auto]
        # image_caption: section_title에서 숫자 prefix + em dash 이후 제거해 짧은 제목 추출
        if role == "image_caption" and not values:
            sec_t = content.get("section_title", "")
            if sec_t:
                import re as _re2
                bare = _re2.sub(r'^\d+(\.\d+)*[\s.]*', '', sec_t).strip()
                short = _shorten_toc_item(bare)
                values = [short] if short else [bare]
        prefix = rule.get("prefix", "")
        is_image_slot = (role == "image_slots")
        if is_image_slot:
            _img_fmt = _load_common_formatting().get("image_slot_format", "[이미지: {description}]")
        for i, sid in enumerate(ids):
            raw_val = values[i] if i < len(values) else ""
            txt = raw_val
            if txt:
                if is_image_slot:
                    txt = _img_fmt.format(description=str(raw_val))
                else:
                    txt = prefix + str(txt)
                # 실제 shape cx 읽기 (rule 기본값 대신 shape 실측 우선)
                sp_shape = _find_shape_by_id(root, sid)
                actual_cx = rule["cx"]
                if sp_shape is not None:
                    xfrm = sp_shape.find(f".//{{{ns_a}}}xfrm")
                    if xfrm is not None:
                        ext = xfrm.find(f"{{{ns_a}}}ext")
                        if ext is not None:
                            actual_cx = int(ext.get("cx", rule["cx"]))
                # 이미지 슬롯은 '…' 없이 깔끔히 끊는다(전체 상세는 plan 보존), 그 외는 기존대로 말줄임
                txt = _truncate_to_lines(txt, actual_cx, rule["pt"], rule["lines"], ellipsis=not is_image_slot)
                # … 말줄임 발생 시 자동 요약 재시도
                if txt.endswith("…"):
                    short = _shorten_toc_item(str(raw_val))
                    txt2 = _truncate_to_lines(prefix + short if prefix else short, actual_cx, rule["pt"], rule["lines"])
                    if not txt2.endswith("…"):
                        txt = txt2
            if is_image_slot and txt:
                # image_slots: 원본 rPr이 없는 경우가 많아 서식을 직접 주입
                _set_image_slot_text(root, sid, txt)
            else:
                _slide_set_helper(root, ns_p, ns_a, sid, txt)

    _clear_residual_placeholders(root)
    _write_xml(root, xml_path)


# ── 레이아웃 placeholder 기반 본문 편집기 (레거시 폴백) ──────────────

def _edit_layout_slide(xml_path: Path, slide_plan: dict) -> None:
    """
    GS Neotek 템플릿 slide22 전용 편집기.

    slide22 레이아웃:
      ID=8  — 대제목 헤더 (상단 왼쪽)
      ID=9  — 중제목 헤더 (상단 중앙)
      ID=10 — 이미지 영역 (회색 박스, 편집 안 함)
      ID=11 — 핵심 설명 bullet 3개 (하단, sz=1600) ← 메인 콘텐츠
      ID=14 — 좌측 사이드바 레이블 (짧은 키워드)
      ID=16 — 좌측 사이드바 설명 (한 문장)

    편집 전략:
      - ID=8  → slide title (대제목)
      - ID=11 → bullets (최대 3개, 이후는 말줄임)
      - ID=14 → bullets[0] 키워드 요약 (30자 이내)
      - ID=16 → body 한 줄 요약
      - 나머지 → placeholder 클리어
    """
    import copy as _copy
    ns_p = _NS_P
    ns_a = _NS_A
    title   = slide_plan.get("title", "")
    content = slide_plan.get("content", {})
    bullets = (content.get("bullets") or content.get("items") or
               content.get("steps") or [])
    body    = content.get("body", "")

    try:
        tree = ET.parse(xml_path)
        root = tree.getroot()
    except ET.ParseError:
        return

    def _set_placeholder(sp: ET.Element, text: str) -> None:
        """shape의 첫 번째 paragraph를 text로 교체 (rPr 보존)."""
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None:
            return
        for p in txBody.findall(f"{{{ns_a}}}p"):
            first_r = p.find(f"{{{ns_a}}}r")
            orig_rPr = None
            if first_r is not None:
                rPr_e = first_r.find(f"{{{ns_a}}}rPr")
                if rPr_e is not None:
                    orig_rPr = _copy.deepcopy(rPr_e)
                    orig_rPr.set("lang", _ppt_lang())
                    orig_rPr.set("dirty", "0")
            for r in p.findall(f"{{{ns_a}}}r"):
                p.remove(r)
            end_rpr = p.find(f"{{{ns_a}}}endParaRPr")
            idx = list(p).index(end_rpr) if end_rpr is not None else len(p)
            r_new = ET.Element(f"{{{ns_a}}}r")
            if orig_rPr is not None:
                r_new.append(orig_rPr)
            else:
                ET.SubElement(r_new, f"{{{ns_a}}}rPr", lang=_ppt_lang(), dirty="0")
            t = ET.SubElement(r_new, f"{{{ns_a}}}t")
            t.text = text
            p.insert(idx, r_new)
            break  # 첫 paragraph만

    def _set_bullets(sp: ET.Element, items: list[str]) -> None:
        """shape의 paragraph들을 bullet 항목으로 채운다."""
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None:
            return
        # 기존 rPr 보존
        orig_rPr = None
        for r in sp.findall(f".//{{{ns_a}}}r"):
            rPr_e = r.find(f"{{{ns_a}}}rPr")
            if rPr_e is not None:
                orig_rPr = _copy.deepcopy(rPr_e)
                orig_rPr.set("lang", _ppt_lang())
                orig_rPr.set("dirty", "0")
                break

        paras = txBody.findall(f"{{{ns_a}}}p")
        for i, item in enumerate(items):
            if i < len(paras):
                p = paras[i]
                for r in p.findall(f"{{{ns_a}}}r"):
                    p.remove(r)
                end_rpr = p.find(f"{{{ns_a}}}endParaRPr")
                idx = list(p).index(end_rpr) if end_rpr is not None else len(p)
                r_new = ET.Element(f"{{{ns_a}}}r")
                if orig_rPr is not None:
                    r_new.append(_copy.deepcopy(orig_rPr))
                t = ET.SubElement(r_new, f"{{{ns_a}}}t")
                t.text = item
                p.insert(idx, r_new)
            else:
                # 새 paragraph 추가
                p_new = ET.SubElement(txBody, f"{{{ns_a}}}p")
                r_new = ET.SubElement(p_new, f"{{{ns_a}}}r")
                if orig_rPr is not None:
                    r_new.append(_copy.deepcopy(orig_rPr))
                t = ET.SubElement(r_new, f"{{{ns_a}}}t")
                t.text = item
                ET.SubElement(p_new, f"{{{ns_a}}}endParaRPr", lang=_ppt_lang())
        # 초과 paragraph 비우기
        for p in paras[len(items):]:
            for r in p.findall(f"{{{ns_a}}}r"):
                for t in r.findall(f"{{{ns_a}}}t"):
                    t.text = ""

    # slide22 실제 폭 (slideLayout5에서 측정):
    #   ID=14 사이드바 레이블: 173pt → 2,197,100 EMU (한글 ~12자/줄)
    #   ID=16 사이드바 설명:   173pt → 2,197,100 EMU (12pt, 한글 ~17자/줄)
    #   ID=11 하단 bullet:     750pt → 9,525,000 EMU (16pt, 한글 ~50자/줄)
    _CX_SIDEBAR = 2_197_100
    _CX_BULLETS = 9_525_000

    # ── ID=8: 대제목 헤더 (넓은 영역 → 전체 제목) ──────────────
    sp8 = _find_shape_by_id(root, _CONTENT_TITLE_ID)
    if sp8 is not None:
        _set_placeholder(sp8, title)

    # ── ID=11: 하단 bullet (넓은 영역 → 각 bullet 1줄 요약) ─────
    sp11 = _find_shape_by_id(root, _CONTENT_BULLETS_ID)
    if sp11 is not None:
        content_list = bullets[:3] if bullets else (
            [s.strip() for s in body.split("·") if s.strip()][:3] if body else []
        )
        # 각 bullet을 750pt 폭에 맞게 1줄로 정리
        trimmed = [_truncate_to_lines(b, _CX_BULLETS, 16, max_lines=1)
                   for b in content_list]
        if trimmed:
            _set_bullets(sp11, trimmed)

    # ── ID=14: 사이드바 레이블 (좁은 영역 → 핵심 키워드 2줄) ────
    sp14 = _find_shape_by_id(root, "14")
    if sp14 is not None:
        # bullets[0]에서 핵심 키워드만 뽑아서 2줄 이내로 요약
        raw = bullets[0] if bullets else title
        label = _truncate_to_lines(raw, _CX_SIDEBAR, 20, max_lines=2)
        _set_placeholder(sp14, label)

    # ── ID=16: 사이드바 설명 (좁은 영역 → body 3줄 요약) ─────────
    sp16 = _find_shape_by_id(root, "16")
    if sp16 is not None:
        raw16 = body or (bullets[1] if len(bullets) > 1 else "")
        desc = _truncate_to_lines(raw16, _CX_SIDEBAR, 12, max_lines=3)
        _set_placeholder(sp16, desc)

    # ── 나머지: placeholder 텍스트만 클리어 ─────────────────────
    _clear_residual_placeholders(root)

    _write_xml(root, xml_path)


# ── 범용 본문 슬라이드 편집기 ─────────────────────────────────

_PLACEHOLDER_TEXTS = (
    _load_placeholder_patterns() or
    re.compile(
        r"(lorem|작성해주세요|중제목을|대제목을|제목을|설명을|설명 타이틀|"
        r"상세 설명|이미지를|01 제목|insight / definition|image|"
        r"\[insert|\bTODO\b|ipsum dolor|"
        r"이미지/\s*영상|이미지/영상|image placeholder|"
        r"nibh euismod|tincidunt ut|elit,\s*sed diam|nonummy|"
        r"짧은 텍스트에 사용|텍스트에 사용|1\.1\s|1\.2\s|1\.3\s|"
        r"insight /|conclusion /|definition|\bicon\b|image\b|"
        r"핵심 설명을 작성|부분 설명|설명 타이틀|01 핵심|02 핵심|03 핵심|"
        r"ppt 대제목|01 중제목|01 컨텐츠|중제목 작성|대제목 작성|컨텐츠 작성|"
        r"solution 0[123]|sevice 0[123]|keyword\b|step[1-4]|"
        r"1\.4\s|텍스트 작성해주세요|텍스트 길어질|작성하여 사용|"
        r"\[아이콘|\[이미지(?!:)|관련 도식|아이콘 이미지|"
        r"아래 확장|최대 [0-9]+\s*줄|여기에 입력|내용을 입력|"
        r"항목 0[1-9]|항목0[1-9]|세부 항목|detail 0[1-9])",
        re.IGNORECASE,
    )
)


def _clear_residual_placeholders(root: ET.Element) -> bool:
    """
    편집 후 남아있는 placeholder 텍스트를 빈 문자열로 교체한다.
    편집된 콘텐츠(한국어 실제 내용)는 건드리지 않는다.
    수정이 일어나면 True 반환.
    """
    ns_p = _NS_P
    ns_a = _NS_A
    modified = False

    for sp in root.findall(f".//{{{ns_p}}}sp"):
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None:
            continue
        for p in txBody.findall(f".//{{{ns_a}}}p"):
            # paragraph 내 모든 run의 텍스트를 합쳐서 판단
            full_text = "".join(
                t.text or "" for t in p.findall(f".//{{{ns_a}}}t")
            )
            if _PLACEHOLDER_TEXTS.search(full_text):
                for t in p.findall(f".//{{{ns_a}}}t"):
                    t.text = ""
                modified = True
            # 개별 run도 검사 (분산 lorem 대응)
            else:
                for r in p.findall(f"{{{ns_a}}}r"):
                    t = r.find(f"{{{ns_a}}}t")
                    if t is not None and t.text and _PLACEHOLDER_TEXTS.search(t.text):
                        t.text = ""
                        modified = True

    return modified


def edit_content_slide(xml_path: Path, slide_plan: dict) -> None:
    """
    cover/toc 외 모든 role의 슬라이드를 편집한다.
    텍스트 shape을 Y 좌표 순으로 정렬한 뒤 content 필드를 매핑.

    매핑 우선순위:
      1번 shape(최상단): title
      2번 shape: subtitle / body / description
      이후 shape: bullets / items / steps (항목 순서대로)
    """
    import copy as _copy

    ns_a = _NS_A
    ns_p = _NS_P

    title   = slide_plan.get("title", "")
    content = slide_plan.get("content", {})

    # 본문 콘텐츠 우선순위 추출
    body_text = (
        content.get("body")
        or content.get("description")
        or content.get("subtitle")
        or ""
    )
    list_items: list[str] = (
        content.get("bullets")
        or content.get("items")
        or content.get("steps")
        or []
    )

    try:
        tree = ET.parse(xml_path)
        root = tree.getroot()
    except ET.ParseError:
        return

    # 텍스트 shape 수집 → Y 좌표 정렬
    shapes_with_y: list[tuple[int, ET.Element]] = []
    for sp in root.findall(f".//{{{ns_p}}}sp"):
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None:
            continue
        xfrm = sp.find(f".//{{{ns_a}}}xfrm")
        off  = xfrm.find(f"{{{ns_a}}}off") if xfrm is not None else None
        y    = int(off.get("y", 0)) if off is not None else 0
        shapes_with_y.append((y, sp))

    shapes_with_y.sort(key=lambda x: x[0])
    shapes = [sp for _, sp in shapes_with_y]

    def _set_text(sp: ET.Element, text: str) -> None:
        """shape의 첫 번째 paragraph의 run 교체 (rPr 보존)."""
        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None:
            return
        for p in txBody.findall(f"{{{ns_a}}}p"):
            # 기존 rPr 복사
            first_r = p.find(f"{{{ns_a}}}r")
            orig_rPr = None
            if first_r is not None:
                rPr_e = first_r.find(f"{{{ns_a}}}rPr")
                if rPr_e is not None:
                    orig_rPr = _copy.deepcopy(rPr_e)
                    orig_rPr.set("lang", _ppt_lang())
                    orig_rPr.set("dirty", "0")

            for r in p.findall(f"{{{ns_a}}}r"):
                p.remove(r)

            end_rpr = p.find(f"{{{ns_a}}}endParaRPr")
            idx = list(p).index(end_rpr) if end_rpr is not None else len(p)

            r_new = ET.Element(f"{{{ns_a}}}r")
            if orig_rPr is not None:
                r_new.append(orig_rPr)
            else:
                ET.SubElement(r_new, f"{{{ns_a}}}rPr", lang=_ppt_lang(), dirty="0")
            t_new = ET.SubElement(r_new, f"{{{ns_a}}}t")
            t_new.text = text
            p.insert(idx, r_new)
            break  # 첫 번째 paragraph만

    # shape 0: 제목
    if shapes:
        _set_text(shapes[0], title)

    # shape 1: 본문 단일 텍스트
    if len(shapes) > 1 and body_text:
        _set_text(shapes[1], body_text)

    # shape 2+: 리스트 항목
    if list_items:
        start = 2 if body_text and len(shapes) > 2 else 1
        for idx, item in enumerate(list_items):
            si = start + idx
            if si < len(shapes):
                _set_text(shapes[si], item)
            else:
                break  # shape 부족 시 중단

    _write_xml(root, xml_path)


# ── 병렬 Plan 생성 (2-Phase) ─────────────────────────────────

def _build_claude_client():
    """
    Claude 클라이언트 빌드 (generate_plan_with_claude와 동일한 우선순위 로직).
    Vertex > Bedrock > Anthropic 순서, ANTHROPIC_VERTEX_PROJECT_ID가 있으면 Vertex 우선.
    """
    import os, anthropic
    vertex_proj   = os.environ.get("ANTHROPIC_VERTEX_PROJECT_ID")
    vertex_region = os.environ.get("CLOUD_ML_REGION")
    aws_region    = os.environ.get("AWS_REGION")
    api_key       = os.environ.get("ANTHROPIC_API_KEY")
    explicit      = os.environ.get("PPT_SKILL_BACKEND", "auto")
    cc_bedrock    = os.environ.get("CLAUDE_CODE_USE_BEDROCK", "0")
    cc_vertex     = os.environ.get("CLAUDE_CODE_USE_VERTEX", "0")

    if explicit == "vertex":
        use_vertex, use_bedrock = True, False
    elif explicit == "bedrock":
        use_vertex, use_bedrock = False, True
    elif explicit == "anthropic":
        use_vertex, use_bedrock = False, False
    elif cc_bedrock == "1":
        use_vertex, use_bedrock = False, True
    elif cc_vertex == "1":
        use_vertex, use_bedrock = True, False
    elif cc_bedrock == "0" and cc_vertex == "0":
        use_vertex, use_bedrock = False, False
    else:
        use_vertex  = bool(vertex_proj or (vertex_region and not aws_region))
        use_bedrock = bool(aws_region) and not use_vertex

    if use_vertex:
        return anthropic.AnthropicVertex(
            project_id=vertex_proj or "",
            region=vertex_region or "us-east5"
        ), None
    if use_bedrock:
        try:
            import boto3
            from botocore.config import Config as _BotoConfig2
            return None, boto3.client("bedrock-runtime", region_name=aws_region or "us-east-1",
                                      config=_BotoConfig2(read_timeout=300, connect_timeout=10))
        except ImportError:
            pass  # boto3 없으면 Anthropic 폴백
    if api_key:
        return anthropic.Anthropic(api_key=api_key), None
    return None, None


def _call_api(client, bedrock, system: str, user: str, max_tokens: int = 4096) -> str:
    """백엔드 무관 단일 API 호출."""
    import json as _json
    model = __import__("os").environ.get("ANTHROPIC_DEFAULT_SONNET_MODEL", "claude-sonnet-4-6")
    messages = [{"role": "user", "content": user}]
    if bedrock:
        body = _json.dumps({"anthropic_version": "bedrock-2023-05-31",
                            "max_tokens": max_tokens, "system": system, "messages": messages})
        resp = bedrock.invoke_model(
            modelId="us.anthropic.claude-sonnet-4-6", body=body)
        return _json.loads(resp["body"].read())["content"][0]["text"].strip()
    # Vertex AI는 us.anthropic. 접두사 미지원
    if hasattr(client, "_client") or type(client).__name__ == "AnthropicVertex":
        if model.startswith("us.anthropic."):
            model = model[len("us.anthropic."):]
    resp = client.messages.create(model=model, max_tokens=max_tokens,
                                   thinking={"type": "adaptive"},
                                   system=system, messages=messages)
    return resp.content[0].text.strip()


def _parse_json_response(raw: str) -> dict | list | None:
    if raw.startswith("```"):
        raw = re.sub(r"^```[a-z]*\n?", "", raw)
        raw = re.sub(r"\n?```$", "", raw)
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        return None


def generate_outline(topic: str, audience: str, n_slides: int,
                     slide_info: list[dict], memory: dict,
                     constraints: list[str]) -> list[dict] | None:
    """
    Phase 1: 슬라이드 구조(아웃라인)만 생성 — 빠른 1회 호출.
    각 슬라이드의 role, template_file, title, key_points만 반환.
    """
    try:
        client, bedrock = _build_claude_client()
    except Exception:
        return None
    if client is None and bedrock is None:
        return None

    available = {s["file"]: s["texts"][:2] for s in slide_info}
    layout_hints = memory.get("slide_layout_hints", {})
    constraint_str = "\n".join(f"- {c}" for c in constraints) if constraints else "없음"

    system = (
        "당신은 PPT 구조 설계자입니다. 슬라이드 아웃라인만 JSON 배열로 출력합니다.\n"
        "각 항목: {index, role, template_file, title, key_points: [str]}\n"
        "role 목록: cover|toc|section|content|three_col|steps|timeline|closing\n"
        "JSON 배열만 출력, 다른 텍스트 없음."
    )
    user = (
        f"주제: {topic}\n청중: {audience}\n슬라이드 수: {n_slides}\n\n"
        f"사용 가능 템플릿: {list(available.keys())[:10]}\n"
        f"레이아웃 힌트: {layout_hints}\n"
        f"품질 제약 조건:\n{constraint_str}\n\n"
        f"{n_slides}장 아웃라인을 JSON 배열로 작성하세요."
    )

    try:
        raw = _call_api(client, bedrock, system, user, max_tokens=2048)
        result = _parse_json_response(raw)
        if isinstance(result, list) and len(result) > 0:
            return result
    except Exception as e:
        print(f"  ⚠ outline 생성 실패: {e}")
    return None


def _generate_single_slide_content(args: tuple) -> dict:
    """
    Phase 2 단일 슬라이드 콘텐츠 생성 — ThreadPoolExecutor에서 호출.
    args: (slide_outline, topic, audience, constraints, client, bedrock)
    """
    slide_outline, topic, audience, constraints, client, bedrock = args
    role  = slide_outline.get("role", "content")
    title = slide_outline.get("title", "")
    keys  = slide_outline.get("key_points", [])

    schema = PLAN_CONTENT_SCHEMA.get(role, {})
    req_fields = schema.get("required", []) + schema.get("recommended", [])
    constraint_str = "\n".join(f"- {c}" for c in constraints) if constraints else "없음"

    system = (
        "당신은 PPT 콘텐츠 전문가입니다. 슬라이드 1장의 content 필드만 JSON으로 출력합니다.\n"
        "JSON 객체만 출력, 다른 텍스트 없음."
    )
    user = (
        f"주제: {topic}\n청중: {audience}\n"
        f"슬라이드 역할: {role}\n제목: {title}\n핵심 키워드: {keys}\n"
        f"필요 필드: {req_fields}\n"
        f"품질 제약:\n{constraint_str}\n\n"
        "위 슬라이드의 content JSON을 작성하세요. "
        "텍스트는 발표 수준으로 구체적이고 전문적이어야 합니다."
    )

    try:
        raw = _call_api(client, bedrock, system, user, max_tokens=1024)
        content = _parse_json_response(raw)
        if isinstance(content, dict):
            return {**slide_outline, "content": content}
    except Exception as e:
        pass  # 폴백: key_points를 bullets로

    # 폴백
    return {**slide_outline, "content": {"bullets": keys, "body": title}}


def generate_plan_parallel(topic: str, audience: str, n_slides: int,
                            slide_info: list[dict], memory: dict,
                            constraints: list[str]) -> dict | None:
    """
    2-Phase 병렬 Plan 생성:
      Phase 1: 아웃라인 1회 호출 (role/title/key_points)
      Phase 2: 슬라이드별 콘텐츠 병렬 호출 (ThreadPoolExecutor)
    슬라이드 수와 무관하게 일정한 응답 시간 유지.
    """
    try:
        client, bedrock = _build_claude_client()
    except Exception:
        return None
    if client is None and bedrock is None:
        return None

    print("  [병렬 생성] Phase 1: 아웃라인 생성...")
    outline = generate_outline(topic, audience, n_slides, slide_info, memory, constraints)
    if not outline:
        return None
    print(f"  ✓ 아웃라인 {len(outline)}장 완료")

    print(f"  [병렬 생성] Phase 2: {len(outline)}장 콘텐츠 병렬 생성...")
    args_list = [
        (slide, topic, audience, constraints, client, bedrock)
        for slide in outline
    ]
    slides: list[dict] = [None] * len(outline)
    max_workers = min(8, len(outline))

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_idx = {
            executor.submit(_generate_single_slide_content, args): i
            for i, args in enumerate(args_list)
        }
        for future in as_completed(future_to_idx):
            idx = future_to_idx[future]
            try:
                slides[idx] = future.result()
            except Exception:
                slides[idx] = {**outline[idx], "content": {}}

    slides = [s for s in slides if s is not None]
    print(f"  ✓ 병렬 콘텐츠 생성 완료 ({len(slides)}장)")

    return {
        "title":    topic,
        "topic":    topic,
        "audience": audience,
        "n_slides": n_slides,
        "slides":   slides,
        "_generated_by": "parallel",
    }


# ── Vision Fix Agent (독립 검증 에이전트) ─────────────────────

def _gather_shape_info(xml_path: Path) -> list[dict]:
    """
    슬라이드 XML에서 text shape 목록을 추출한다.
    반환: [{id, y, cx, current_texts}]  (Y 오름차순 정렬)
    """
    ns_p = _NS_P
    ns_a = _NS_A
    shapes = []
    try:
        tree = ET.parse(xml_path)
        root = tree.getroot()
    except ET.ParseError:
        return shapes

    for sp in root.findall(f".//{{{ns_p}}}sp"):
        cpr = sp.find(f"{{{ns_p}}}nvSpPr/{{{ns_p}}}cNvPr")
        if cpr is None:
            continue
        xfrm = sp.find(f".//{{{ns_a}}}xfrm")
        off  = xfrm.find(f"{{{ns_a}}}off") if xfrm is not None else None
        ext  = xfrm.find(f"{{{ns_a}}}ext") if xfrm is not None else None
        y  = int(off.get("y", 0)) if off is not None else 0
        cx = int(ext.get("cx", 0)) if ext is not None else 0
        x  = int(off.get("x", 0)) if off is not None else 0

        txBody = sp.find(f"{{{ns_p}}}txBody")
        if txBody is None:
            continue
        texts = [
            "".join(t.text or "" for t in p.findall(f".//{{{ns_a}}}t"))
            for p in txBody.findall(f".//{{{ns_a}}}p")
        ]
        texts = [t for t in texts if t.strip()]

        shapes.append({
            "id":   cpr.get("id", "?"),
            "name": cpr.get("name", ""),
            "y":    y,
            "x":    x,
            "cx":   cx,
            "current_texts": texts[:3],  # 최대 3개 미리보기
        })

    return sorted(shapes, key=lambda s: s["y"])


def _call_vision_api(system: str, content: list) -> str | None:
    """Vision(멀티모달) Claude API 호출."""
    try:
        client, bedrock = _build_claude_client()
    except Exception:
        return None
    if client is None and bedrock is None:
        return None
    model = __import__("os").environ.get(
        "ANTHROPIC_DEFAULT_SONNET_MODEL", "claude-sonnet-4-6")
    # Vertex AI는 us.anthropic. 접두사 미지원
    if client is not None and type(client).__name__ == "AnthropicVertex":
        if model.startswith("us.anthropic."):
            model = model[len("us.anthropic."):]
    messages = [{"role": "user", "content": content}]
    try:
        if bedrock:
            import json as _j, boto3
            body = _j.dumps({
                "anthropic_version": "bedrock-2023-05-31",
                "max_tokens": 4096, "system": system, "messages": messages,
            })
            resp = bedrock.invoke_model(
                modelId="us.anthropic.claude-sonnet-4-6", body=body)
            return _j.loads(resp["body"].read())["content"][0]["text"].strip()
        resp = client.messages.create(
            model=model, max_tokens=4096,
            thinking={"type": "adaptive"},
            system=system, messages=messages)
        return resp.content[0].text.strip()
    except Exception as e:
        print(f"  ⚠ Vision API 호출 실패: {e}")
        return None


def _run_vision_fix_agent(
    work_dir: Path,
    qa_images: list[str],
    plan: dict,
    slides_dir: Path,
) -> list[dict]:
    """
    독립 Vision 검증 에이전트.
    QA 이미지 + plan 의도 + shape 구조를 독립 컨텍스트에서 분석해
    슬라이드별 수술적(shape ID 기반) 수정 지시를 생성한다.

    반환 스키마:
    [
      {
        "slide_file": "slide8.xml",
        "slide_index": 3,
        "fixes": [
          {"shape_id": "4", "action": "set_paragraphs",
           "texts": ["항목1", "항목2", "항목3"]},
          {"shape_id": "2", "action": "clear"}
        ]
      }
    ]
    action 종류:
      set_text       - 단일 텍스트 설정
      set_paragraphs - 다중 paragraph 설정
      clear          - 모든 텍스트 제거
    """
    import base64, os

    if not qa_images:
        return []

    try:
        client, bedrock = _build_claude_client()
    except Exception:
        return []
    if client is None and bedrock is None:
        print("  ⚠ Vision Fix Agent: API 없음 — 건너뜀")
        return []

    # 슬라이드 인덱스 → plan 매핑
    plan_by_idx = {s["index"]: s for s in plan.get("slides", [])}
    # 이미지 순서 → slide index 매핑 (이미지는 1-indexed)
    image_paths = [p for p in qa_images if Path(p).exists()]

    system = """당신은 독립적인 PPT 슬라이드 품질 검증 에이전트입니다.
슬라이드 이미지, 계획된 콘텐츠, shape 구조를 보고 정확한 수정 지시를 생성합니다.

반드시 아래 JSON 배열 형식으로만 응답하세요:
[
  {
    "slide_file": "slideN.xml",
    "slide_index": N,
    "has_issues": true/false,
    "issue_summary": "문제 요약",
    "fixes": [
      {
        "shape_id": "ID",
        "action": "set_text|set_paragraphs|clear|resize_textbox",
        "text": "단일 텍스트 (set_text용)",
        "texts": ["항목1", "항목2"],  // set_paragraphs용
        "new_cy": 600000  // resize_textbox용: 새 높이(EMU). 1줄=314603, 2줄=629206
      }
    ]
  }
]

수정 원칙:
- placeholder(작성해주세요, 중제목, 대제목, lorem ipsum 등)는 반드시 제거 또는 교체
- 실제 콘텐츠(plan에 있는 내용)로만 교체
- 이미 올바른 슬라이드는 has_issues=false, fixes=[]
- shape_id는 반드시 제공된 shape 목록의 실제 ID 사용
- ⛔ 폰트 크기(sz) 절대 변경 금지 — set_text/set_paragraphs에 sz 속성 포함 불가. 텍스트가 넘치면 resize_textbox(cy 확장)만 허용
- ⛔ 텍스트 내용 임의 단축 금지 — plan에 있는 실제 텍스트를 줄이거나 생략하지 말 것
  단, 텍스트가 "…"(말줄임표)로 끝나는 경우는 예외: 박스에 맞게 요약한 짧은 텍스트로 교체할 것 (아래 룰 참조)
- ⛔ [이미지: ...] 형식 텍스트는 이미지 자리 설명 텍스트 — 의도적으로 작성된 콘텐츠이므로 절대 제거/수정/문제 보고 금지. 이 텍스트가 보이는 것은 정상이다
- ⛔ 파란색 배경 헤더 박스(item_title)에 있는 한국어/영어 항목명(예: Transformer, RAG, Fine-tuning, Amazon Bedrock 등)은 실제 콘텐츠 — 절대 clear/set_text("")로 제거 금지. 번호·점·prefix 잔류 텍스트(예: "1.1 ", "01 ")만 제거 대상
- ⛔ clear 및 set_text("") 지시 사용 금지 — 내용을 지워야 할 것 같은 경우 fixes에 포함하지 말고, issue_summary에 "삭제 검토 필요: shape_id=X (이유)" 형태로만 기록할 것. 실제 삭제는 사용자 확인 후 처리
- ⛔ 비어있는 shape에 새 텍스트 임의 삽입 금지 — shape에 텍스트가 없는 경우, 그 shape이 이미지 슬롯(placeholder "이미지"/"image")이 아니라면 채우지 말 것. 의도적으로 비워둔 장식용 도형일 수 있음
- ⛔ 제공된 shape 목록(shapes 필드)에 없는 shape ID로 수정 지시 생성 금지
- shape ID=14(사이드바 레이블, 173pt 폭): 텍스트 overflow가 보여도 수정 금지 — 이미 너비에 맞게 요약됨
- shape ID=16(사이드바 설명, 173pt 폭): 동일하게 수정 금지
- shape ID=17(사이드바 소제목, 173pt 폭): 텍스트 overflow가 보여도 수정 금지 — 이미 자동 요약됨
- shape ID=18(사이드바 설명, 173pt 폭): 텍스트 overflow가 보여도 수정 금지 — 이미 자동 요약됨
- 사이드바는 의도적으로 짧은 요약 텍스트를 담는 좁은 영역임. ID=17,18은 사이드바 영역으로 overflow가 시각적으로 보여도 수정 대상이 아님
- shape ID=9: "01 중제목 작성", "01 컨텐츠 작성" 같은 영문/숫자 placeholder만 비울 것. "01 Anthropic 창립 배경", "02 Claude 모델 계보" 같이 숫자 prefix + 한국어 실제 콘텐츠가 있으면 절대 수정 금지 — 이 숫자 prefix는 의도적으로 시스템이 붙인 슬라이드 번호
- ⛔ shape ID=8 절대 수정 금지 — 시스템이 목차에서 자동 추출한 챕터 대제목. Vision Fix 대상이 아님
- ⛔ shape ID=9 수정 시 주의 — "01 중제목 작성" 같은 순수 영문/숫자 placeholder만 비울 것. "01 Anthropic 창립 배경"처럼 숫자 + 실제 한국어 내용이 있으면 절대 수정 금지
- shape ID=8은 대제목(목차 챕터명), ID=9는 중제목(세부 분류명)으로 의도적으로 설정된 콘텐츠임
- "Solution 01", "keyword", "Sevice 01", "Step1"~"Step4" 등 템플릿 안내 텍스트는 placeholder — 비울 것
- ⛔ [이미지: ...] 형식 텍스트는 이미지 슬롯 설명 — 절대 삭제/수정/비우기 금지. "image", "이미지" placeholder와 다름 — 이미 올바른 콘텐츠임
- "| 1.1.1 ..." 형식 텍스트(파이프 기호로 시작하는 sub_heading)는 의도적으로 작성된 소제목 — 절대 수정 금지
- 이미 실제 콘텐츠(한국어 설명문)가 있는 shape는 수정 금지

가시성 검증 (시각적으로 판단):
- 텍스트가 배경색과 대비가 낮아 거의 안 보이는 경우: issue_summary에 "가시성 낮음: shape_id=X (다크 배경에 다크 텍스트)" 형태로 보고. fixes는 비워둘 것 (폰트 색 변경은 별도 처리)
- 다크(네이비/다크블루) 배경 영역의 텍스트가 흰색/밝은색이 아니면 가시성 문제로 보고
- 라이트(흰색/회색) 배경 영역의 텍스트가 흰색/밝은색이면 가시성 문제로 보고

텍스트박스 크기 및 여백 검증 (시각적으로 판단):
- 텍스트가 2줄 이상으로 넘어가는데 텍스트박스 높이가 1줄 크기로 보이면:
  action="resize_textbox", new_cy=629206(2줄) 또는 943809(3줄)로 높이 조정 지시
  단, sidebar(ID=14,16,17,18,25) shape과 sub_heading(ID=41, cy=337169 고정) shape은 resize 금지
- resize_textbox는 최대 1~2줄 초과(new_cy 최대 943809 = 3줄)만 허용. 그 이상(4줄 이상) 초과 시에는 resize 금지 — 텍스트가 너무 길어 구조적으로 박스에 맞지 않는 경우이므로 issue_summary에 기록만 할 것
- 텍스트박스가 배경 도형(회색/컬러 사각형)을 벗어난 경우:
  action="resize_textbox"로 배경 도형 안으로 제한 지시
- 텍스트박스는 배경 도형 안쪽에서 최소 50,000 EMU(약 4mm) 여백이 있어야 함
  배경 도형 경계를 넘거나 딱 붙은 경우 issue로 보고

말줄임표(…) 오버플로 감지 및 수정:
- 텍스트가 "…"(말줄임표, U+2026 또는 "...")로 끝나는 shape을 발견하면 오버플로 상태임
- 해당 shape의 텍스트를 박스 크기에 맞는 짧은 요약 텍스트로 교체하도록 set_text 지시 생성
- 요약 시 핵심 키워드·개념을 유지하되 조사·접속어·부연설명을 제거해 1~2줄로 축약
- "…"가 있는 shape은 반드시 issue로 보고하고 수정 지시를 생성할 것

이미지 슬롯 비어있음 감지 및 채우기:
- 이미지 슬롯이란: 슬라이드 내 회색/파란색 배경의 사각형 영역에 텍스트가 없거나 "image", "이미지" 같은 placeholder만 있는 박스
- 이미지 슬롯이 비어 있거나 placeholder 텍스트만 있는 경우: "[이미지: {슬라이드 주제와 해당 항목에 맞는 구체적인 이미지 설명}]" 형식으로 set_text 지시 생성
- 예시: "[이미지: LLM 학습 파이프라인 다이어그램 — 데이터 수집, 전처리, 파인튜닝, 평가 단계를 보여주는 플로우차트]"
- ⛔ [이미지: ...] 형식 텍스트가 이미 있는 이미지 슬롯은 수정하지 말 것 (이미 처리된 슬롯)
- 슬라이드별 이미지 슬롯 shape ID는 제공된 shape 구조에서 파악할 것"""

    # 슬라이드별 정보 구성
    content_parts = []
    slide_contexts = []

    for i, img_path in enumerate(image_paths[:10]):
        slide_idx = i + 1
        slide_plan = plan_by_idx.get(slide_idx, {})
        slide_file = slide_plan.get("template_file", "")
        xml_path   = slides_dir / slide_file if slide_file else None
        shape_info = _gather_shape_info(xml_path) if xml_path and xml_path.exists() else []

        slide_contexts.append({
            "slide_index": slide_idx,
            "slide_file":  slide_file,
            "role":        slide_plan.get("role", "content"),
            "title":       slide_plan.get("title", ""),
            "content":     slide_plan.get("content", {}),
            "shapes":      shape_info,
        })

        # img_slot IDs 추출 (이미지 슬롯 감지 힌트로 제공)
        _sids_data = _load_slide_shape_ids()
        _slide_key = slide_file.replace(".xml", "") if slide_file else ""
        _slide_sids = _sids_data.get(_slide_key, {})
        _img_slot_ids = []
        if isinstance(_slide_sids, dict):
            for _sl in _slide_sids.get("slots", []):
                if _sl.get("type") == "img_slot":
                    _img_slot_ids.extend(_sl.get("ids", []))

        b64 = base64.standard_b64encode(Path(img_path).read_bytes()).decode()
        img_slot_hint = f"\n이미지 슬롯 shape IDs: {_img_slot_ids}" if _img_slot_ids else ""
        content_parts.append({
            "type": "text",
            "text": f"=== 슬라이드 {slide_idx} ({slide_plan.get('role','')}) ===\n"
                    f"의도 제목: {slide_plan.get('title','')}\n"
                    f"의도 콘텐츠: {json.dumps(slide_plan.get('content',{}), ensure_ascii=False)[:300]}\n"
                    f"shape 구조: {json.dumps(shape_info, ensure_ascii=False)[:400]}"
                    f"{img_slot_hint}",
        })
        content_parts.append({
            "type": "image",
            "source": {"type": "base64", "media_type": "image/jpeg", "data": b64},
        })

    content_parts.append({
        "type": "text",
        "text": "각 슬라이드를 분석해 수정 지시 JSON을 반환하세요. "
                "placeholder가 없고 의도한 콘텐츠가 올바르게 표시된 슬라이드는 has_issues=false."
    })

    # 슬라이드별 개별 분석 (토큰 한도 초과 방지)
    # 전용 편집기가 처리하는 슬라이드는 Vision Fix 대상 제외
    # — toc: edit_toc_slide가 7-paragraph 구조를 완전 관리, Vision Fix가 수정하면 구조 파괴
    # — cover: edit_cover_slide가 담당
    # — closing: 편집 금지 (레이아웃에 "감사합니다." 내장)
    SKIP_ROLES = {"cover", "closing", "toc"}
    all_results = []
    for ctx in slide_contexts:
        if ctx.get("role") in SKIP_ROLES:
            continue
        img_path = image_paths[ctx["slide_index"] - 1] \
            if ctx["slide_index"] - 1 < len(image_paths) else None
        if not img_path or not Path(img_path).exists():
            continue

        # 이 슬라이드의 img_slot IDs 추출
        _sids_ctx = _load_slide_shape_ids()
        _slide_key_ctx = ctx["slide_file"].replace(".xml", "") if ctx["slide_file"] else ""
        _slide_sids_ctx = _sids_ctx.get(_slide_key_ctx, {})
        _img_slot_ids_ctx = []
        if isinstance(_slide_sids_ctx, dict):
            for _sl in _slide_sids_ctx.get("slots", []):
                if _sl.get("type") == "img_slot":
                    _img_slot_ids_ctx.extend(_sl.get("ids", []))
        _img_slot_hint_ctx = f"\n이미지 슬롯 shape IDs: {_img_slot_ids_ctx}" if _img_slot_ids_ctx else ""

        b64 = base64.standard_b64encode(Path(img_path).read_bytes()).decode()
        per_slide_content = [
            {
                "type": "text",
                "text": (
                    f"슬라이드 {ctx['slide_index']} ({ctx['role']})\n"
                    f"의도 제목: {ctx['title']}\n"
                    f"의도 콘텐츠: {json.dumps(ctx['content'], ensure_ascii=False)[:300]}\n"
                    f"shape 구조: {json.dumps(ctx['shapes'], ensure_ascii=False)[:400]}"
                    f"{_img_slot_hint_ctx}"
                ),
            },
            {
                "type": "image",
                "source": {"type": "base64", "media_type": "image/jpeg", "data": b64},
            },
            {
                "type": "text",
                "text": (
                    "이 슬라이드에 placeholder(작성해주세요, lorem ipsum 등) 또는 레이아웃 문제가 있으면 "
                    "JSON 객체로 수정 지시를 반환하세요. 문제 없으면 {\"has_issues\": false}."
                ),
            },
        ]

        try:
            raw = _call_vision_api(system, per_slide_content)
            if not raw or not raw.strip():
                continue
            # JSON 블록 추출 (```json ... ```, { ... }, 또는 전체 텍스트)
            text = raw.strip()
            # 방법 1: ```json ... ``` 코드 블록
            m = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", text, re.DOTALL)
            if m:
                text = m.group(1)
            else:
                # 방법 2: 첫 번째 { ... } 객체 추출
                m2 = re.search(r"\{.*\}", text, re.DOTALL)
                if m2:
                    text = m2.group(0)
            parsed = json.loads(text)
            if isinstance(parsed, dict):
                # ctx의 정확한 파일명으로 강제 설정 (에이전트가 잘못 추측할 수 있음)
                parsed["slide_file"]  = ctx["slide_file"]
                parsed["slide_index"] = ctx["slide_index"]
                all_results.append(parsed)
        except (json.JSONDecodeError, Exception) as e:
            print(f"  ⚠ slide {ctx['slide_index']} Vision 분석 실패: {e}")

    issues_found = [r for r in all_results if r.get("has_issues")]
    if issues_found:
        print(f"  [Vision Fix] {len(issues_found)}/{len(all_results)}장 수정 필요")
        for r in issues_found:
            print(f"    slide {r.get('slide_index')} ({r.get('slide_file')}): "
                  f"{r.get('issue_summary', '')[:60]}")
    else:
        print(f"  [Vision Fix] {len(all_results)}장 분석 — 수정 불필요")
    return all_results


def _apply_fix_instructions(work_dir: Path, fix_instructions: list[dict]) -> bool:
    """
    Vision Fix Agent의 수정 지시를 슬라이드 XML에 적용한다.
    수정이 일어나면 True 반환.
    set_text로 수정된 텍스트는 instruction에 '_plan_updates' 키로 기록된다.
    """
    import copy as _copy
    ns_p = _NS_P
    ns_a = _NS_A
    slides_dir = work_dir / "unpacked" / "ppt" / "slides"
    any_modified = False

    for instruction in fix_instructions:
        if not instruction.get("has_issues"):
            continue
        fixes = instruction.get("fixes", [])
        if not fixes:
            continue

        slide_file = instruction.get("slide_file", "")
        xml_path   = slides_dir / slide_file
        if not xml_path.exists():
            continue

        try:
            tree = ET.parse(xml_path)
            root = tree.getroot()
        except ET.ParseError:
            continue

        modified = False
        plan_updates: list[dict] = []  # set_text 수정 내용 추적 (plan.json 반영용)
        shape_map = {}
        for sp in root.findall(f".//{{{ns_p}}}sp"):
            cpr = sp.find(f"{{{ns_p}}}nvSpPr/{{{ns_p}}}cNvPr")
            if cpr is not None:
                shape_map[cpr.get("id", "")] = sp

        # 이미지 슬롯 shape ID 목록 (맑은고딕 파란색 서식 적용 대상)
        _sids_data_fix = _load_slide_shape_ids()
        _slide_key_fix = slide_file.replace(".xml", "")
        _slide_sids_fix = _sids_data_fix.get(_slide_key_fix, {})
        _img_slot_ids_fix: set = set()
        if isinstance(_slide_sids_fix, dict):
            for _sl in _slide_sids_fix.get("slots", []):
                if _sl.get("type") == "img_slot":
                    _img_slot_ids_fix.update(_sl.get("ids", []))

        for fix in fixes:
            shape_id = str(fix.get("shape_id", ""))
            action   = fix.get("action", "")
            sp = shape_map.get(shape_id)
            if sp is None:
                continue

            txBody = sp.find(f"{{{ns_p}}}txBody")
            if txBody is None:
                continue

            # 기존 rPr·pPr 보존용 (정렬·색상 등 서식 유지)
            first_rPr = None
            first_pPr = None
            for p_elem in txBody.findall(f"{{{ns_a}}}p"):
                if first_pPr is None:
                    pPr_e = p_elem.find(f"{{{ns_a}}}pPr")
                    if pPr_e is not None:
                        first_pPr = _copy.deepcopy(pPr_e)
                for r in p_elem.findall(f"{{{ns_a}}}r"):
                    rPr_e = r.find(f"{{{ns_a}}}rPr")
                    if rPr_e is not None and first_rPr is None:
                        first_rPr = _copy.deepcopy(rPr_e)
                        first_rPr.set("lang", _ppt_lang())
                        first_rPr.set("dirty", "0")
                if first_rPr and first_pPr:
                    break

            def _make_para(text_str: str) -> ET.Element:
                p = ET.Element(f"{{{ns_a}}}p")
                # 원본 pPr(algn 등) 복사 — 없으면 빈 pPr
                if first_pPr is not None:
                    p.append(_copy.deepcopy(first_pPr))
                else:
                    ET.SubElement(p, f"{{{ns_a}}}pPr")
                r_new = ET.SubElement(p, f"{{{ns_a}}}r")
                if first_rPr is not None:
                    r_new.append(_copy.deepcopy(first_rPr))
                else:
                    ET.SubElement(r_new, f"{{{ns_a}}}rPr",
                                  lang=_ppt_lang(), dirty="0")
                t = ET.SubElement(r_new, f"{{{ns_a}}}t")
                t.text = text_str
                ET.SubElement(p, f"{{{ns_a}}}endParaRPr",
                              lang=_ppt_lang(), dirty="0")
                return p

            if action == "clear":
                # 내용 삭제는 사용자 확인 없이 실행 금지 — 경고만 출력
                print(f"  ⚠ [Vision Fix] 삭제 지시 차단: shape {shape_id} in {slide_file} "
                      f"(이유: {fix.get('reason', instruction.get('issue_summary','')[:40])})")
                continue

            elif action == "set_text":
                text = fix.get("text", "")
                # 빈 텍스트로 내용 삭제 시도 차단
                if not text.strip():
                    print(f"  ⚠ [Vision Fix] 빈 텍스트 삭제 지시 차단: shape {shape_id} in {slide_file}")
                    continue
                # 이미지 슬롯이면 맑은고딕 파란색 서식 적용
                if shape_id in _img_slot_ids_fix:
                    _set_image_slot_text(root, shape_id, text)
                else:
                    for p in txBody.findall(f"{{{ns_a}}}p"):
                        txBody.remove(p)
                    txBody.append(_make_para(text))
                plan_updates.append({"shape_id": shape_id, "text": text})
                modified = True

            elif action == "set_paragraphs":
                texts = fix.get("texts", [])
                if not texts:
                    continue
                for p in txBody.findall(f"{{{ns_a}}}p"):
                    txBody.remove(p)
                for t in texts:
                    txBody.append(_make_para(t))
                modified = True

            elif action == "resize_textbox":
                # 텍스트박스 cy 동적 조정 + 연관 shape 재배치
                new_cy = fix.get("new_cy")
                if new_cy and sp is not None:
                    # 레이아웃 상속 shape이면 실제 좌표 먼저 명시화
                    _materialize_layout_shape(root, shape_id, xml_path)
                    spPr = sp.find(f"{{{ns_p}}}spPr")
                    if spPr is not None:
                        xfrm = spPr.find(f"{{{ns_a}}}xfrm")
                        if xfrm is not None:
                            ext = xfrm.find(f"{{{ns_a}}}ext")
                            if ext is not None:
                                ext.set("cy", str(new_cy))
                                # body_title이면 body_desc도 재배치
                                text_for_resize = fix.get("text", "")
                                if text_for_resize:
                                    _resize_sidebar_and_reposition_desc(
                                        root, shape_id, None, text_for_resize
                                    )
                                modified = True

        if modified:
            _write_xml(root, xml_path)
            any_modified = True
            if plan_updates:
                instruction["_plan_updates"] = plan_updates
            print(f"  ✓ Vision Fix 적용: {slide_file} ({len(fixes)}개 지시)")

    return any_modified


def _sync_plan_with_fixes(plan: dict, fix_instructions: list[dict], plan_path: Path) -> None:
    """
    Vision Fix에서 set_text로 수정된 텍스트를 plan.json의 content에 반영한다.
    shape_id → content 필드 매핑은 slide_shape_ids.json 기반으로 추론한다.
    """
    sids_data = _load_slide_shape_ids()
    plan_by_file = {s["template_file"]: s for s in plan.get("slides", [])}
    changed = False

    for instruction in fix_instructions:
        updates = instruction.get("_plan_updates")
        if not updates:
            continue
        slide_file = instruction.get("slide_file", "")
        slide_plan = plan_by_file.get(slide_file)
        if not slide_plan:
            continue
        content = slide_plan.setdefault("content", {})
        slide_key = slide_file.replace(".xml", "")
        slide_sids = sids_data.get(slide_key, {})
        if not isinstance(slide_sids, dict):
            continue

        # shape_id → (slot_type, content_key, index) 매핑 구축
        sid_map: dict[str, tuple[str, str, int]] = {}
        for slot in slide_sids.get("slots", []):
            stype = slot.get("type", "")
            ckey  = slot.get("content_key", "")
            for idx, sid in enumerate(slot.get("ids", [])):
                sid_map[str(sid)] = (stype, ckey, idx)

        for upd in updates:
            sid  = str(upd.get("shape_id", ""))
            text = upd.get("text", "")
            if sid not in sid_map:
                continue
            stype, ckey, idx = sid_map[sid]
            if stype in ("item_title", "item_body", "img_slot", "insight"):
                # 리스트 필드 — 인덱스 위치 업데이트
                lst = content.get(ckey, [])
                if isinstance(lst, list):
                    while len(lst) <= idx:
                        lst.append("")
                    lst[idx] = text
                    content[ckey] = lst
                    changed = True
            elif stype in ("sub_heading", "section_title", "section_desc"):
                content[ckey] = text
                changed = True

    if changed:
        plan_path.write_text(json.dumps(plan, ensure_ascii=False, indent=2))
        print("  ✓ plan.json 업데이트 (Vision Fix 반영)")


# ── 메인 루프 ────────────────────────────────────────────────

def _success_rate(runs: list) -> float | None:
    """qa_ok이 확정(bool)인 run만으로 성공률 계산. None(판정 보류)·레거시(필드 부재)는 제외."""
    rated = [r for r in runs if isinstance(r.get("qa_ok"), bool)]
    if not rated:
        return None
    return round(sum(1 for r in rated if r["qa_ok"]) / len(rated), 3)


def _record_run_experience(topic: str, plan: dict, vision_issues: int,
                           qa_done: bool = True) -> str | None:
    """AHE 경험 관찰성(❷) + auto-update: 매 실행 후 long_term_memory에 run 기록을 append하고
    메타(total_runs/success_rate/last_updated)를 갱신한다. evolution/last_run_digest.json에도 요약 기록.
    실패해도 생성 결과엔 영향 없도록 전부 try/except로 감싼다 (AHE_PRINCIPLES §5 경험 관찰성).

    qa_done=False (skill 경로, inline_vision_qa=False)면 엔진이 QA를 안 한 것이므로
    qa_ok=None(판정 보류)로 기록한다. 실제 판정은 오케스트레이터의 독립 QA 에이전트가
    update_last_run_qa()로 채운다 (확증편향 차단 + success_rate 정확성, F1).

    반환값: 기록한 run의 run_id (off-by-one 없이 update_last_run_qa(run_id=...)로
    정확히 그 run의 qa_ok를 닫기 위함). 기록 실패 시 None (#10)."""
    try:
        from datetime import datetime as _dt
        import uuid as _uuid
        slides = plan.get("slides", [])
        run_id = _uuid.uuid4().hex
        rec = {
            "run_id": run_id,
            "ts": _dt.now().isoformat(timespec="seconds"),
            "topic": topic,
            "n_slides": len(slides),
            "templates": [s.get("template_file", "") for s in slides],
            "vision_issues": int(vision_issues),
            "qa_done": bool(qa_done),
            # 엔진이 QA했을 때만 확정값. skill 경로는 None → 독립 QA가 나중에 채움
            "qa_ok": (int(vision_issues) == 0) if qa_done else None,
        }
        mem_path = SKILL_DIR / "harness" / "long_term_memory.json"
        mem = json.loads(mem_path.read_text(encoding="utf-8"))
        runs = mem.setdefault("runs", [])
        runs.append(rec)
        mem["runs"] = runs[-50:]  # 최근 50건만 유지 (progressive disclosure — 토큰 절약)
        mem["total_runs"] = mem.get("total_runs", 0) + 1
        sr = _success_rate(mem["runs"])
        if sr is not None:
            mem["success_rate"] = sr
        mem["last_updated"] = rec["ts"][:10]
        mem_path.write_text(json.dumps(mem, ensure_ascii=False, indent=2), encoding="utf-8")

        evo_dir = SKILL_DIR / "evolution"
        evo_dir.mkdir(exist_ok=True)
        (evo_dir / "last_run_digest.json").write_text(
            json.dumps(rec, ensure_ascii=False, indent=2), encoding="utf-8")
        _qa = rec["qa_ok"] if rec["qa_ok"] is not None else "보류(독립 QA 대기)"
        print(f"  ✓ AHE 경험 기록: total_runs={mem['total_runs']}, "
              f"success_rate={mem.get('success_rate')}, qa_ok={_qa}, "
              f"run_id={run_id[:8]} → harness/long_term_memory.json")
        return run_id
    except Exception as _e:
        print(f"  ⚠ AHE 경험 기록 실패(무시): {_e}")
        return None


def update_last_run_qa(qa_ok: bool, run_id: str | None = None) -> None:
    """오케스트레이터(독립 QA 에이전트)의 실제 판정을 해당 run 기록에 반영한다 (F1).
    skill 경로에서 _record_run_experience가 qa_ok=None으로 남긴 것을 확정값으로 채우고
    success_rate를 재계산한다. SKILL.md 11.5단계에서 독립 QA 종료 후 호출.

    어떤 run을 닫을지 결정하는 순서 (#10 off-by-one 회피):
      1. run_id가 주어지면 → 그 run_id와 일치하는 run을 정확히 닫는다
         (생성과 QA 판정 사이에 다른 run이 기록돼도 안전).
      2. run_id가 없으면 → qa_ok가 아직 None(판정 보류)인 가장 최근 run을 닫는다
         (블라인드 runs[-1]은 보류 run이 아닐 수 있어 오판 위험 → 보류 run 우선).
      3. 보류 run도 없으면 → runs[-1] (레거시 호환)."""
    try:
        mem_path = SKILL_DIR / "harness" / "long_term_memory.json"
        mem = json.loads(mem_path.read_text(encoding="utf-8"))
        runs = mem.get("runs", [])
        if not runs:
            print("  ⚠ QA 판정 반영: run 기록 없음"); return

        target = None
        if run_id is not None:
            target = next((r for r in reversed(runs)
                           if r.get("run_id") == run_id), None)
            if target is None:
                print(f"  ⚠ QA 판정 반영: run_id={run_id[:8]} 매칭 실패 — "
                      f"보류 run으로 폴백")
        if target is None:
            # 판정 보류(qa_ok=None) 상태인 가장 최근 run을 우선 닫는다
            target = next((r for r in reversed(runs)
                           if r.get("qa_ok") is None), None)
        if target is None:
            target = runs[-1]  # 레거시 호환

        target["qa_ok"] = bool(qa_ok)
        target["qa_done"] = True
        sr = _success_rate(runs)
        if sr is not None:
            mem["success_rate"] = sr
        mem_path.write_text(json.dumps(mem, ensure_ascii=False, indent=2), encoding="utf-8")
        _rid = target.get("run_id", "?")
        _rid = _rid[:8] if isinstance(_rid, str) else _rid
        print(f"  ✓ 독립 QA 판정 반영: qa_ok={qa_ok}, run_id={_rid} "
              f"→ success_rate={mem.get('success_rate')}")
    except Exception as _e:
        print(f"  ⚠ QA 판정 반영 실패(무시): {_e}")


def run_ppt_generation(
    topic: str,
    template_path: Path,
    work_dir: Path,
    audience: str = "전문가",
    n_slides: int = 10,
    plan_override: dict | None = None,
    cleanup_work_dir: bool = True,
    layout_from_pptx: Path | None = None,
    inline_vision_qa: bool = True,
) -> Path:
    """
    analyze_template → generate_plan → edit_slide 루프 → pack → verify
    layout_from_pptx가 지정되면 해당 PPTX의 슬라이드 순서를 그대로 사용한다.
    최종 output.pptx 경로를 반환한다.

    inline_vision_qa (AHE_PRINCIPLES §2 생성≠판정):
      True (기본, headless 폴백) — 엔진이 인라인 Vision Fix Agent로 자기 결과를 검증·수정.
        같은 프로세스의 자기-검증이라 확증편향 위험이 있어, 오케스트레이터(Claude Code 세션)가
        없는 headless 실행에서만 권장.
      False (skill 경로 권장) — 인라인 vision QA를 건너뛴다. QA는 SKILL.md 지시대로
        오케스트레이터가 '생성 컨텍스트를 모르는 독립 격리 에이전트'로 수행해 확증편향을 차단한다.
      어느 경우든 결정적 검증(verifier_rules·placeholder·폰트)은 항상 수행된다.
    """
    print(f"\n[PPT Generator] topic={topic!r}, slides={n_slides}, audience={audience!r}")

    # ── 언팩 ─────────────────────────────────────
    unpacked = work_dir / "unpacked"
    subprocess.run(
        [sys.executable, str(SKILL_DIR / "scripts" / "office" / "unpack.py"),
         str(template_path), str(unpacked)],
        check=True, capture_output=True,
    )
    print("  ✓ 언팩 완료")

    # ── 분석 ─────────────────────────────────────
    slide_info = analyze_template(work_dir)
    print(f"  ✓ 슬라이드 {len(slide_info)}개 분석 완료")

    # ── 하네스 로드 ───────────────────────────────
    memory: dict = {}
    memory_path = SKILL_DIR / "harness" / "long_term_memory.json"
    if memory_path.exists():
        memory = json.loads(memory_path.read_text())

    # Vision 분석에서 축적된 planning constraints 로드
    constraints: list[str] = memory.get("planning_constraints", [])
    known_fixes: list[dict] = memory.get("known_failure_fixes", [])
    if constraints:
        print(f"  ✓ planning constraints {len(constraints)}개 로드")
    if known_fixes:
        print(f"  ✓ known_failure_fixes {len(known_fixes)}개 로드")

    # ── layout_from_pptx 처리: banned slide 감지 후 사용자 확인 ────────
    _layout_order: list[str] | None = None
    if layout_from_pptx is not None:
        raw_order = _extract_layout_order(Path(layout_from_pptx))
        cat = _load_slide_catalog()
        banned = set(cat.get("banned_slides", {}).keys())
        # 대체 후보: verified + unverified 중 banned가 아닌 것
        _all = (list(cat.get("verified_slides", {}).keys()) +
                list(cat.get("unverified_slides", {}).keys()) or
                cat.get("allowed_content_slides", []))
        allowed = [s for s in _all if s not in banned]
        replaced: list[tuple[str, str]] = []
        resolved: list[str] = []
        for tmpl in raw_order:
            if tmpl in banned:
                # 같은 계열 추천 (슬라이드 번호 근접 순)
                num = int(re.search(r"\d+", tmpl).group()) if re.search(r"\d+", tmpl) else 99
                candidate = min(allowed, key=lambda x: abs(int(re.search(r"\d+", x).group()) - num) if re.search(r"\d+", x) else 99)
                replaced.append((tmpl, candidate))
                resolved.append(candidate)
            else:
                resolved.append(tmpl)
        if replaced:
            print("\n  ⚠️  참조 PPTX에 banned 슬라이드가 포함되어 있습니다:")
            for orig, repl in replaced:
                ban_reason = cat.get("banned_slides", {}).get(orig, "사용 금지")
                print(f"    - {orig} ({ban_reason}) → 추천 대체: {repl}")
            print("  계속 진행합니다 (추천 슬라이드로 자동 대체). 다른 슬라이드를 원하면 layout_from_pptx 없이 직접 plan을 지정하세요.")
        print(f"  ✓ 레이아웃 고정 모드 ({len(resolved)}장): {resolved}")
        _layout_order = resolved
        n_slides = len(resolved)

    # ── 계획 생성 ─────────────────────────────────
    if plan_override is not None:
        plan = plan_override
        print(f"  ✓ 기존 plan 주입 ({len(plan.get('slides', []))}개 슬라이드) — API 호출 생략")
    else:
        # 15장 초과: 병렬 2-phase 생성 / 이하: 단일 호출
        plan = None
        if n_slides > 14 and _layout_order is None:
            print(f"  [병렬 모드] {n_slides}장 → 2-phase 병렬 생성")
            plan = generate_plan_parallel(topic, audience, n_slides,
                                           slide_info, memory, constraints)

        if plan is None:
            plan = generate_plan_with_claude(topic, audience, n_slides,
                                              slide_info, memory,
                                              layout_order=_layout_order)
            if plan:
                print("  ✓ Claude API 단일 호출로 계획 생성")
            else:
                plan = generate_plan(topic, audience, n_slides, slide_info)
                print("  ✓ 규칙 기반 계획 생성 (폴백)")

    # ── Closing 슬라이드 보장 ─────────────────────
    # 총 N장 = 표지(1) + 목차(1) + 본문(N-3) + 마무리(1)
    # Claude가 closing을 빠뜨리면 자동 추가
    has_closing = any(s.get("role") == "closing" for s in plan.get("slides", []))
    if not has_closing:
        closing_idx = len(plan["slides"]) + 1
        plan["slides"].append({
            "index": closing_idx,
            "template_file": "slide46.xml",
            "role": "closing",
            "title": "감사합니다",
            "content": {}
        })
        # 초과 content 슬라이드 제거 (총 n_slides 유지)
        content_slides = [s for s in plan["slides"] if s.get("role") not in ("cover","toc","closing")]
        max_content = n_slides - 3  # cover+toc+closing
        if len(content_slides) > max_content:
            to_remove = len(content_slides) - max_content
            removed = 0
            new_slides = []
            for s in reversed(plan["slides"]):
                if s.get("role") not in ("cover","toc","closing") and removed < to_remove:
                    removed += 1
                else:
                    new_slides.insert(0, s)
            plan["slides"] = new_slides
        # index 재부여
        for i, s in enumerate(plan["slides"], 1):
            s["index"] = i
        print(f"  ✓ 마무리 슬라이드 자동 추가 (총 {len(plan['slides'])}장, 본문 {n_slides-3}장)")

    # ── 챕터 대제목 후처리: section_title 번호 기준으로 title 자동 보정 ──
    # LLM이 title에 슬라이드 제목을 쓰는 경우를 방지
    toc_slide = next((s for s in plan["slides"] if s.get("role") == "toc"), None)
    if toc_slide:
        toc_items = toc_slide.get("content", {}).get("items", [])
        _dynamic_chapter_map = {}
        for idx, item in enumerate(toc_items, 1):
            _dynamic_chapter_map[str(idx)] = item if isinstance(item, str) else str(item)
    else:
        # toc 슬라이드가 없으면 동적 챕터맵 없음 → _infer_chapter_title이 plan.title로 폴백
        # (과거 _CHAPTER_TITLE_MAP 전역은 정의된 적 없어 NameError였음 — ML 전용 chapter_map.json도 의도적으로 미사용)
        _dynamic_chapter_map = {}

    # _infer_chapter_title이 chapter_map.json(ML 전용)을 읽지 않도록
    # 현재 문서의 챕터 맵을 캐시에 주입 → 렌더링 시 올바른 헤더 표시
    global _CHAPTER_MAP_CACHE
    if _dynamic_chapter_map:
        _CHAPTER_MAP_CACHE = dict(_dynamic_chapter_map)

    for s in plan["slides"]:
        if s.get("role") in ("cover", "toc", "closing"): continue
        sec_title = s.get("content", {}).get("section_title", "")
        if not sec_title: continue
        parts = sec_title.split(".")
        if parts and parts[0].strip().isdigit():
            ch = _dynamic_chapter_map.get(parts[0].strip())
            if ch and s.get("title") != ch:
                s["title"] = ch

    # ── section_title 중복 방지: title과 동일하면 번호 prefix 추가 ──
    for s in plan["slides"]:
        if s.get("role") in ("cover", "toc", "closing"): continue
        content = s.get("content", {})
        sec_title = content.get("section_title", "")
        slide_title = s.get("title", "")
        if not sec_title or not slide_title: continue
        # 번호 제거 후 비교
        bare_sec = re.sub(r'^\d+(\.\d+)*[\s.。]*', '', sec_title).strip()
        bare_title = re.sub(r'^\d+(\.\d+)*[\s.。]*', '', slide_title).strip()
        if bare_sec == bare_title and bare_sec:
            # section_title에 인덱스가 없으면 챕터 번호 prefix 부여
            if not re.match(r'^\d', sec_title.strip()):
                # 챕터 번호는 동적 맵에서 역추출 (하드코딩 "1." 금지 — N챕터 오라벨 방지, F3)
                _ch_num = next((k for k, v in _dynamic_chapter_map.items()
                                if v == slide_title), None) or "1"
                content["section_title"] = f"{_ch_num}. {sec_title}"
            else:
                # 같은 텍스트인데 번호만 있는 경우 → 보조 키워드 추가 불가, 경고만
                print(f"  ⚠️  slide {s.get('index')}: section_title이 title과 동일 — LLM 재생성 권장")

    # ── 중제목 순번(subtitle_seq) 주입 — 대제목 기준 01/02/03 카운터 ──
    _subtitle_counter: dict[str, int] = {}
    for s in plan["slides"]:
        if s.get("role") in ("cover", "toc", "closing"): continue
        ch_title = s.get("title", "")
        if not ch_title: continue
        _subtitle_counter[ch_title] = _subtitle_counter.get(ch_title, 0) + 1
        s["subtitle_seq"] = _subtitle_counter[ch_title]

    # ── Layout variant 자동 선택 (JSON 규칙 기반) ──
    # slide27 → slide28 등 body_title 길이에 따라 자동 전환
    variant_changes = _resolve_layout_variants(plan)
    for vc in variant_changes:
        print(f"  ✓ variant 전환: {vc}")

    # ── Plan 제약 강제 (banned slides 자동 교체) ──
    plan, constraint_changes = enforce_plan_constraints(plan, slide_info)

    # ── 표지 날짜를 오늘 날짜로 자동 설정 ───────────
    from datetime import date as _date
    today_str = _date.today().strftime("%Y.%m.%d")
    for s in plan["slides"]:
        if s.get("role") == "cover":
            s.setdefault("content", {})["date"] = today_str

    # ── 중복 template_file 복사 ───────────────────
    # 같은 template_file을 여러 슬라이드가 쓰면 편집이 덮어써짐
    # → 복사본을 만들어 각 슬라이드가 독립 파일을 갖게 함
    seen_templates: dict[str, int] = {}
    slides_dir_tmp = work_dir / "unpacked" / "ppt" / "slides"
    for s in plan["slides"]:
        tmpl = s.get("template_file", "")
        if not tmpl or s.get("role") in ("cover", "toc", "closing"):
            continue
        if tmpl in seen_templates:
            # 복사본 생성
            seen_templates[tmpl] += 1
            suffix = seen_templates[tmpl]
            new_name = tmpl.replace(".xml", f"_c{suffix}.xml")
            src = slides_dir_tmp / tmpl
            dst = slides_dir_tmp / new_name
            if src.exists() and not dst.exists():
                shutil.copy2(src, dst)
                # rels 파일도 복사
                src_rels = slides_dir_tmp / "_rels" / f"{tmpl}.rels"
                dst_rels = slides_dir_tmp / "_rels" / f"{new_name}.rels"
                if src_rels.exists():
                    shutil.copy2(src_rels, dst_rels)
                # 올바른 위치(closing 슬라이드 앞)에 삽입 → 순서 보장
                closing_file = next(
                    (sv.get("template_file") for sv in plan.get("slides", [])
                     if sv.get("role") == "closing"), None)
                _add_slide_to_presentation(work_dir, new_name,
                                           before_filename=closing_file)
                _register_slide_content_type(work_dir, new_name)
            s["template_file"] = new_name
            print(f"  ✓ 중복 방지: slide {s['index']} {tmpl} → {new_name}")
        else:
            seen_templates[tmpl] = 1

    # ── TOC 항목 및 페이지번호 설정 ─────────────────
    # 목차 구조: 1개 TOC 항목 = 1개 챕터, 챕터당 N개 슬라이드 가능
    # page_nums = 각 챕터의 첫 번째 슬라이드 번호
    # TOC 항목 수 ≤ content 슬라이드 수 보장 (항목당 최소 1개 슬라이드)
    content_slides = [s for s in plan["slides"]
                      if s.get("role") not in ("cover", "toc", "closing")]
    n_content = len(content_slides)

    for s in plan["slides"]:
        if s.get("role") == "toc":
            c = s.setdefault("content", {})
            items = (c.get("items")
                     or [i["text"] if isinstance(i, dict) else i
                         for i in c.get("toc_items", [])]
                     or c.get("bullets", []))

            # TOC 항목 수가 content 슬라이드 수를 초과하면 자름
            # (항목당 최소 1개 슬라이드 보장)
            if len(items) > n_content:
                items = items[:n_content]

            # 챕터별 첫 번째 슬라이드 번호 계산
            # content slide를 TOC 항목 수로 균등 분배하고 첫 페이지만 표시
            n_items = len(items)
            if n_items > 0:
                # 각 챕터의 시작 슬라이드 인덱스 (균등 분배)
                slides_per_chapter = n_content / n_items
                page_nums = [
                    f"{content_slides[min(int(i * slides_per_chapter), n_content - 1)]['index']:02d}"
                    for i in range(n_items)
                ]
            else:
                page_nums = []

            c["items"]     = items
            c["page_nums"] = page_nums
            print(f"  ✓ TOC: {n_items}개 챕터, 페이지번호={page_nums} (content {n_content}장)")
    for ch in constraint_changes:
        print(f"  ⚠ 제약 교체: {ch}")

    # ── Plan 검증 ─────────────────────────────────
    ok, warnings = validate_plan(plan)
    for w in warnings[:10]:
        print(f"  ⚠ plan 검증: {w}")
    if not ok:
        subtitle_missing = [w for w in warnings if "subtitle" in w]
        if subtitle_missing:
            print(f"  ✗ subtitle(중제목) 누락 {len(subtitle_missing)}건 — Zone1 ID=9 비어있음. "
                  f"plan 재생성 필요 (LLM이 subtitle 필드를 포함해야 함)")
        else:
            print("  ⚠ 필수 content 필드 누락 — 편집 시 폴백 처리됨")

    plan_path = work_dir / "plan.json"
    plan_path.write_text(json.dumps(plan, ensure_ascii=False, indent=2))
    print(f"  ✓ plan.json 저장 ({len(plan['slides'])}개 슬라이드)")

    # ── 편집 루프 ────────────────────────────────
    slides_dir = work_dir / "unpacked" / "ppt" / "slides"
    for slide_plan in plan["slides"]:
        xml_path = slides_dir / slide_plan.get("template_file", "")

        # TOC 슬라이드: ID=42는 "발표 제목을 작성해주세요" placeholder → 전체 발표 제목 주입
        if slide_plan.get("role") == "toc" and "prs_title" not in slide_plan:
            slide_plan["prs_title"] = plan.get("title", slide_plan.get("title", ""))

        # known_fixes 자동 적용 (편집 전)
        if xml_path.exists():
            apply_known_fixes(xml_path, slide_plan, known_fixes)

        success = edit_slide(work_dir, slide_plan)
        if not success:
            print(f"  → slide {slide_plan['index']} 재시도...")
            edit_slide(work_dir, slide_plan)

    # ── 사용 슬라이드만 남기기 (plan에 없는 슬라이드 제거) ────
    _trim_to_plan_slides(work_dir, plan)

    # ── 차트 데이터 주입 (slide40/41/43 Excel 연동) — 패킹 전 unpacked 차트에 적용 ──
    _update_chart_data(work_dir, plan)

    # ── 패킹 ─────────────────────────────────────
    output = work_dir / "output.pptx"
    # 복사본 슬라이드가 있으면 스키마 검증 우회 (원본에 없는 파일이라 오탐)
    has_copies = any("_c" in s.get("template_file", "") for s in plan["slides"])
    if not pack_output(work_dir, output, skip_validation=has_copies):
        raise RuntimeError("패킹 실패")

    # ── 섹션 정리 ─────────────────────────────────
    restructure_sections(output)

    # ── Verifier 실행 ─────────────────────────────
    print("\n  [Verifier] 규칙 검증 중...")
    violations = execute_verifier_rules(output, work_dir)
    critical = [v for v in violations if v["severity"] == "CRITICAL"]

    # CRITICAL 위반이 있으면 해당 슬라이드 재편집 (최대 1회)
    if critical:
        print(f"  → CRITICAL 위반 {len(critical)}건 감지 — 영향 슬라이드 재편집")
        retry_files = {v.get("file") for v in critical if v.get("file")}
        for slide_plan in plan["slides"]:
            if slide_plan.get("template_file") in retry_files or not retry_files:
                xml_path = slides_dir / slide_plan.get("template_file", "")
                if xml_path.exists():
                    apply_known_fixes(xml_path, slide_plan, known_fixes)
                edit_slide(work_dir, slide_plan)
        # 재패킹 후 재검증
        pack_output(work_dir, output)
        restructure_sections(output)
        violations2 = execute_verifier_rules(output, work_dir)
        remaining = [v for v in violations2 if v["severity"] == "CRITICAL"]
        if remaining:
            print(f"  ⚠ 재시도 후 CRITICAL {len(remaining)}건 잔존 (best-effort 출력)")

    # ── 시각 QA (인라인) ──────────────────────────────────
    # AHE_PRINCIPLES §2: skill 경로(inline_vision_qa=False)는 인라인 자기-검증을 건너뛰고
    # 오케스트레이터가 독립 격리 QA 에이전트로 검증(확증편향 차단). headless는 인라인 폴백 사용.
    images = []
    if inline_vision_qa:
        images = visual_qa(work_dir, output)
        if images:
            print(f"  → QA 이미지 {len(images)}장 생성")
    else:
        print("  → 인라인 vision QA 생략 (독립 QA 에이전트가 검증 — AHE_PRINCIPLES §2 생성≠판정)")

    # ── Vision Fix Agent 루프 (최대 3회) ────────
    MAX_VISION_ITER = 3
    if images:
        for vision_iter in range(1, MAX_VISION_ITER + 1):
            print(f"\n  [Vision Fix Agent] 검증 {vision_iter}/{MAX_VISION_ITER}회...")
            fix_instructions = _run_vision_fix_agent(
                work_dir   = work_dir,
                qa_images  = images,
                plan       = plan,
                slides_dir = slides_dir,
            )

            has_issues = any(f.get("has_issues") for f in fix_instructions)
            if not has_issues:
                print("  ✓ Vision 검증 통과 — 모든 슬라이드 정상")
                break

            applied = _apply_fix_instructions(work_dir, fix_instructions)
            if not applied:
                print("  ⚠ 적용 가능한 수정 없음 — 구조적 한계로 best-effort 종료")
                break

            # Vision Fix 결과를 plan.json에 반영 — 다음 라운드 덮어씌움 방지
            _sync_plan_with_fixes(plan, fix_instructions, plan_path)

            print(f"  → 수정 적용 후 재패킹 (이터레이션 {vision_iter})...")
            _trim_to_plan_slides(work_dir, plan)
            pack_output(work_dir, output)
            restructure_sections(output)

            # 재검증 후 QA 이미지 갱신
            violations = execute_verifier_rules(output, work_dir)
            critical   = [v for v in violations if v["severity"] == "CRITICAL"]
            if critical:
                print(f"  ⚠ CRITICAL {len(critical)}건 잔존")
            else:
                print(f"  ✓ Verifier 통과 (이터레이션 {vision_iter})")

            images = visual_qa(work_dir, output)  # QA 이미지 갱신

            if vision_iter == MAX_VISION_ITER:
                print(f"  ⚠ 최대 반복({MAX_VISION_ITER}회) 도달 — best-effort 출력")

    # ── Vision QA 최종 결과 요약 (미통과 페이지 명시) ──────────
    if "fix_instructions" in dir() and fix_instructions:
        failed = [(f.get("slide_index"), f.get("slide_file"), f.get("issue_summary", ""))
                  for f in fix_instructions if f.get("has_issues")]
        if failed:
            print(f"\n  ❌ QA 미통과 슬라이드 ({len(failed)}건):")
            for idx, fname, summary in failed:
                print(f"    - {idx}p ({fname}): {summary}")
        else:
            print("\n  ✅ 모든 슬라이드 Vision QA 통과")

    # ── 콘텐츠 검증 ──────────────────────────────
    issues = verify_content(output)
    if issues:
        print(f"  ⚠ 플레이스홀더 잔여 ({len(issues)}건):")
        for issue in issues[:3]:
            print(f"    - {issue}")
    else:
        print("  ✓ 콘텐츠 검증 통과")

    # ── 폰트 컴플라이언스 검사 ──────────────────────
    font_warns = check_font_compliance(work_dir / "unpacked" / "ppt" / "slides")
    if font_warns:
        print(f"  ⚠ 폰트 비준수 {len(font_warns)}건 (Pretendard 미지정):")
        for w in font_warns[:3]:
            print(f"    - {w}")
    else:
        print("  ✓ 폰트 검증 통과 (Pretendard)")

    # Vision 이슈 수를 meta로 반환 (조건부 auto-evolve용)
    _vision_critical_total = sum(
        1 for slide in (fix_instructions if "fix_instructions" in dir() else [])
        if slide.get("has_issues")
    )
    # ── Excel 차트 데이터 파일 생성 ─────────────────
    dest_dir_for_excel = work_dir.parent  # runs 상위 디렉토리가 아닌 work_dir 사용
    generate_excel_for_charts(work_dir, plan, work_dir)

    # ── AHE 경험 자동 기록 (❷ 경험 관찰성 + auto-update) ──
    # skill 경로(inline_vision_qa=False)는 qa_ok 보류로 기록 → 독립 QA가 update_last_run_qa로 확정 (F1)
    # 반환된 run_id는 evolution/last_run_digest.json에도 남아 update_last_run_qa(run_id=...) 매칭에 쓰인다 (#10)
    _last_run_id = _record_run_experience(topic, plan, _vision_critical_total, qa_done=inline_vision_qa)

    # ── tmp work_dir 정리 ────────────────────────────
    if cleanup_work_dir:
        import shutil as _shutil3
        safe_name = topic.replace(" ", "_").replace("/", "_")[:60]
        final_output = work_dir.parent.parent / f"{safe_name}.pptx"
        _shutil3.move(str(output), str(final_output))
        _shutil3.rmtree(work_dir, ignore_errors=True)
        print(f"  ✓ 작업 디렉토리 정리 완료 → {final_output}")
        return final_output, _vision_critical_total

    return output, _vision_critical_total
